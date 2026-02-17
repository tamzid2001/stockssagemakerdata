from __future__ import annotations

import base64
import hashlib
import json
import math
import os
import re
import time
from html import unescape
from datetime import date, datetime, timedelta, timezone
from io import BytesIO, StringIO
from statistics import NormalDist
from typing import Any
from urllib.parse import parse_qs, quote_plus, unquote, urlparse

import requests
try:
    from zoneinfo import ZoneInfo
except Exception:  # pragma: no cover - Python < 3.9 fallback
    ZoneInfo = None  # type: ignore

import firebase_admin
from firebase_admin import credentials, firestore, messaging as admin_messaging, storage as admin_storage
from firebase_functions import https_fn, scheduler_fn
from firebase_functions.options import MemoryOption, set_global_options
from firebase_functions.params import SecretParam

try:
    from firebase_admin import remote_config as admin_remote_config  # type: ignore
except Exception:  # pragma: no cover - optional dependency until firebase-admin>=7.x
    admin_remote_config = None

# Bind high-sensitivity API keys via Secret Manager instead of committing them.
OPENAI_API_KEY_SECRET = SecretParam("OPENAI_API_KEY")

set_global_options(max_instances=10, secrets=[OPENAI_API_KEY_SECRET])

SERVICE_ACCOUNT_PATH = os.environ.get(
    "SERVICE_ACCOUNT_PATH",
    os.path.join(os.path.dirname(__file__), "serviceAccountKey.json"),
)
STORAGE_BUCKET = (
    os.environ.get("STORAGE_BUCKET")
    or os.environ.get("FIREBASE_STORAGE_BUCKET")
    # Default bucket names use the GCS bucket, not the public "firebasestorage.app" domain.
    or "quantura-e2e3d.appspot.com"
)
PUBLIC_ORIGIN = (os.environ.get("PUBLIC_ORIGIN") or "https://quantura-e2e3d.web.app").rstrip("/")
ADMIN_EMAIL = "tamzid257@gmail.com"
ALLOWED_STATUSES = {"pending", "in_progress", "fulfilled", "cancelled"}
CONTACT_REQUIRED_FIELDS = {"name", "email", "message"}
FORECAST_SERVICES = {"prophet", "ibm_timemixer"}
BACKTEST_SOURCE_FORMATS = {"python", "tradingview", "metatrader5", "tradelocker"}
FEATURE_VOTE_KEYS = {"uploads", "autopilot"}
FEATURE_VOTE_CHOICES = {"yes", "no"}
REPORT_AGENT_BATCH_SIZE = max(1, min(int(os.environ.get("REPORT_AGENT_BATCH_SIZE", "8") or 8), 40))

ALPACA_API_BASE = os.environ.get("ALPACA_API_BASE", "https://paper-api.alpaca.markets")
ALPACA_DATA_BASE = os.environ.get("ALPACA_DATA_BASE", "https://data.alpaca.markets")
ALPACA_API_KEY = os.environ.get("ALPACA_API_KEY") or os.environ.get("ALPACAAPIKEY")
ALPACA_SECRET_KEY = os.environ.get("ALPACA_SECRET_KEY") or os.environ.get("ALPACASECRETKEY")

IBM_TIMEMIXER_MODEL_ID = os.environ.get("IBM_TIMEMIXER_MODEL_ID", "ibm-granite/granite-timeseries-ttm-r2")
IBM_TIMEMIXER_ENDPOINT = os.environ.get("IBM_TIMEMIXER_ENDPOINT", "").strip()
IBM_TIMEMIXER_API_KEY = os.environ.get("IBM_TIMEMIXER_API_KEY", "").strip()
HUGGINGFACEHUB_API_TOKEN = os.environ.get("HUGGINGFACEHUB_API_TOKEN", "").strip()

SLACK_WEBHOOK_URL = os.environ.get("SLACK_WEBHOOK_URL", "").strip()
FCM_WEB_VAPID_KEY = os.environ.get("FCM_WEB_VAPID_KEY", "").strip()
STRIPE_PUBLIC_KEY = os.environ.get("STRIPE_PUBLIC_KEY", "").strip()
STRIPE_SECRET_KEY = os.environ.get("STRIPE_SECRET_KEY", "").strip()
STRIPE_WEBHOOK_SECRET = os.environ.get("STRIPE_WEBHOOK_SECRET", "").strip()
MASSIVE_API_KEY = os.environ.get("MASSIVE_API_KEY", "").strip()
MASSIVE_BASE_URL = (os.environ.get("MASSIVE_BASE_URL") or "https://api.massive.com").rstrip("/")
UNSPLASH_ACCESS_KEY = str(
    os.environ.get("UNSPLASH_ACCESS_KEY")
    or os.environ.get("UNSPLASH_APPLICATION_ID")
    or ""
).strip()
STRIPE_CONNECT_PLATFORM_FEE_PERCENT = float(os.environ.get("STRIPE_CONNECT_PLATFORM_FEE_PERCENT", "12") or 12)
CREATOR_DEFAULT_SUBSCRIBE_USD = float(os.environ.get("CREATOR_DEFAULT_SUBSCRIBE_USD", "9") or 9)
CREATOR_DEFAULT_THANKS_USD = float(os.environ.get("CREATOR_DEFAULT_THANKS_USD", "5") or 5)
RISK_FREE_RATE = float(os.environ.get("RISK_FREE_RATE", "0.045") or 0.045)
TRENDING_URL = "https://query1.finance.yahoo.com/v1/finance/trending/US"
YAHOO_SEARCH_URL = "https://query2.finance.yahoo.com/v1/finance/search"
DEFAULT_FORECAST_PRICE = 349
OPENAI_API_KEY = OPENAI_API_KEY_SECRET.value.strip() or os.environ.get("OPENAI_API_KEY", "").strip()
AMAZON_NOVA_API_KEY = os.environ.get("AMAZON_NOVA_API_KEY", "").strip()
AMAZON_NOVA_API_ENDPOINT = str(os.environ.get("AMAZON_NOVA_API_ENDPOINT") or "").strip()
AMAZON_NOVA_DEFAULT_MODEL = str(os.environ.get("AMAZON_NOVA_DEFAULT_MODEL") or "amazon.nova-lite-v1:0").strip()
SOCIAL_CONTENT_MODEL = (os.environ.get("SOCIAL_CONTENT_MODEL") or "gpt-5-mini").strip()
SOCIAL_AUTOMATION_ENABLED = str(os.environ.get("SOCIAL_AUTOMATION_ENABLED") or "true").strip().lower() in {
    "1",
    "true",
    "yes",
    "on",
}
SOCIAL_AUTOMATION_TIMEZONE = str(os.environ.get("SOCIAL_AUTOMATION_TIMEZONE") or "America/New_York").strip()
SOCIAL_POSTING_TIMEZONE = str(os.environ.get("SOCIAL_POSTING_TIMEZONE") or SOCIAL_AUTOMATION_TIMEZONE).strip()
SOCIAL_DISPATCH_BATCH_SIZE = max(1, min(int(os.environ.get("SOCIAL_DISPATCH_BATCH_SIZE", "30") or 30), 100))
SOCIAL_DEFAULT_CTA_URL = str(os.environ.get("SOCIAL_DEFAULT_CTA_URL") or PUBLIC_ORIGIN).strip()
SOCIAL_AUTOPILOT_ENABLED = str(os.environ.get("SOCIAL_AUTOPILOT_ENABLED") or "true").strip().lower() in {
    "1",
    "true",
    "yes",
    "on",
}
SOCIAL_AUTOPILOT_USER_ID = str(os.environ.get("SOCIAL_AUTOPILOT_USER_ID") or "quantura_system").strip()
SOCIAL_AUTOPILOT_USER_EMAIL = str(os.environ.get("SOCIAL_AUTOPILOT_USER_EMAIL") or "system@quantura.ai").strip()
SOCIAL_AUTOPILOT_POSTS_PER_CHANNEL = max(
    1,
    min(int(os.environ.get("SOCIAL_AUTOPILOT_POSTS_PER_CHANNEL", "3") or 3), 8),
)
SOCIAL_AUTOPILOT_TOPIC = str(
    os.environ.get("SOCIAL_AUTOPILOT_TOPIC")
    or "Daily Quantura market pulse: top catalysts, risk posture, and setup watchlist"
).strip()
SOCIAL_AUTOPILOT_OBJECTIVE = str(
    os.environ.get("SOCIAL_AUTOPILOT_OBJECTIVE")
    or "Drive qualified users to Quantura forecasting workflows"
).strip()
SOCIAL_AUTOPILOT_AUDIENCE = str(
    os.environ.get("SOCIAL_AUTOPILOT_AUDIENCE")
    or "active investors, analysts, and portfolio operators"
).strip()
SOCIAL_AUTOPILOT_TONE = str(os.environ.get("SOCIAL_AUTOPILOT_TONE") or "institutional, concise, actionable").strip()
SOCIAL_AUTOPILOT_CHANNELS = [
    channel.strip().lower()
    for channel in str(os.environ.get("SOCIAL_AUTOPILOT_CHANNELS") or "x,linkedin,facebook,instagram,tiktok").split(",")
    if channel.strip()
]

TWITTER_BEARER_TOKEN = str(
    os.environ.get("TWITTER_BEARER_TOKEN")
    or os.environ.get("X_BEARER_TOKEN")
    or ""
).strip()
TWITTER_API_KEY = str(
    os.environ.get("TWITTER_API_KEY")
    or os.environ.get("X_API_KEY")
    or os.environ.get("TWITTER_CONSUMER_KEY")
    or ""
).strip()
TWITTER_API_SECRET = str(
    os.environ.get("TWITTER_API_SECRET")
    or os.environ.get("X_API_SECRET")
    or os.environ.get("TWITTER_CONSUMER_SECRET")
    or ""
).strip()
TWITTER_ACCESS_TOKEN = str(
    os.environ.get("TWITTER_ACCESS_TOKEN")
    or os.environ.get("X_ACCESS_TOKEN")
    or ""
).strip()
TWITTER_ACCESS_TOKEN_SECRET = str(
    os.environ.get("TWITTER_ACCESS_TOKEN_SECRET")
    or os.environ.get("X_ACCESS_TOKEN_SECRET")
    or ""
).strip()
LINKEDIN_ACCESS_TOKEN = str(os.environ.get("LINKEDIN_ACCESS_TOKEN") or "").strip()
LINKEDIN_AUTHOR_URN = str(os.environ.get("LINKEDIN_AUTHOR_URN") or "").strip()
FACEBOOK_PAGE_ID = str(os.environ.get("FACEBOOK_PAGE_ID") or "").strip()
FACEBOOK_PAGE_ACCESS_TOKEN = str(os.environ.get("FACEBOOK_PAGE_ACCESS_TOKEN") or "").strip()
INSTAGRAM_BUSINESS_ACCOUNT_ID = str(os.environ.get("INSTAGRAM_BUSINESS_ACCOUNT_ID") or "").strip()
INSTAGRAM_ACCESS_TOKEN = str(os.environ.get("INSTAGRAM_ACCESS_TOKEN") or FACEBOOK_PAGE_ACCESS_TOKEN).strip()
INSTAGRAM_DEFAULT_IMAGE_URL = str(
    os.environ.get("INSTAGRAM_DEFAULT_IMAGE_URL")
    or "https://images.unsplash.com/photo-1535320903710-d993d3d77d29?auto=format&fit=crop&w=1280&q=80"
).strip()
TIKTOK_ACCESS_TOKEN = str(os.environ.get("TIKTOK_ACCESS_TOKEN") or "").strip()
TIKTOK_OPEN_ID = str(os.environ.get("TIKTOK_OPEN_ID") or "").strip()
SOCIAL_POPULAR_CHANNELS = [
    "x",
    "linkedin",
    "facebook",
    "instagram",
    "threads",
    "reddit",
    "tiktok",
    "youtube",
    "pinterest",
]
SOCIAL_CHANNEL_WEBHOOK_ENV = {
    "x": "SOCIAL_WEBHOOK_X",
    "linkedin": "SOCIAL_WEBHOOK_LINKEDIN",
    "facebook": "SOCIAL_WEBHOOK_FACEBOOK",
    "instagram": "SOCIAL_WEBHOOK_INSTAGRAM",
    "threads": "SOCIAL_WEBHOOK_THREADS",
    "reddit": "SOCIAL_WEBHOOK_REDDIT",
    "tiktok": "SOCIAL_WEBHOOK_TIKTOK",
    "youtube": "SOCIAL_WEBHOOK_YOUTUBE",
    "pinterest": "SOCIAL_WEBHOOK_PINTEREST",
}
SCREENER_NOTES_TICKER_HINTS: dict[str, list[str]] = {
    "nancy pelosi": ["NVDA", "AAPL", "MSFT", "GOOGL", "AMZN", "PANW", "TSLA"],
    "pelosi": ["NVDA", "AAPL", "MSFT", "GOOGL", "AMZN"],
    "jeff bezos": ["AMZN", "MSFT", "GOOGL", "NVDA", "SHOP"],
    "bezos": ["AMZN", "MSFT", "GOOGL", "NVDA"],
    "warren buffett": ["AAPL", "BAC", "KO", "OXY", "CVX", "AXP"],
    "bill gates": ["MSFT", "BRK.B", "CAT", "WM"],
    "elon musk": ["TSLA", "NVDA", "GOOGL", "MSFT"],
    "cathie wood": ["TSLA", "COIN", "ROKU", "SQ", "PATH", "CRSP"],
    "ark invest": ["TSLA", "COIN", "ROKU", "SQ", "CRSP", "PATH"],
    "blackrock": ["AAPL", "MSFT", "NVDA", "AMZN", "META", "GOOGL"],
    "vanguard": ["AAPL", "MSFT", "NVDA", "AMZN", "META", "BRK.B"],
    "fidelity": ["AAPL", "MSFT", "NVDA", "AMZN", "GOOGL", "META"],
    "bridgewater": ["SPY", "GLD", "TLT", "IEMG", "EEM", "XLF"],
    "renaissance": ["MSFT", "NVDA", "AAPL", "AMZN", "META"],
    "citadel": ["AAPL", "MSFT", "NVDA", "AMZN", "META", "JPM"],
    "cnbc": ["NVDA", "AAPL", "TSLA", "MSFT", "AMZN", "META"],
    "bloomberg": ["AAPL", "MSFT", "NVDA", "AMZN", "TSLA", "JPM"],
    "reuters": ["XOM", "CVX", "AAPL", "MSFT", "AMZN", "JPM"],
    "wall street journal": ["AAPL", "MSFT", "JPM", "XOM", "UNH", "NVDA"],
    "marketwatch": ["TSLA", "NVDA", "AAPL", "AMD", "PLTR", "META"],
    "motley fool": ["AMZN", "MSFT", "NVDA", "CRWD", "MELI", "SHOP"],
    "seeking alpha": ["AAPL", "MSFT", "NVDA", "O", "SCHD", "VICI"],
    "hedge fund": ["AAPL", "MSFT", "NVDA", "AMZN", "META", "GOOGL"],
}
META_PIXEL_ID = str(os.environ.get("META_PIXEL_ID") or "1643823927053003").strip()
META_CAPI_ACCESS_TOKEN = str(os.environ.get("META_CAPI_ACCESS_TOKEN") or "").strip()
META_GRAPH_API_VERSION = str(os.environ.get("META_GRAPH_API_VERSION") or "v21.0").strip() or "v21.0"
META_ALLOWED_ORIGINS = {
    PUBLIC_ORIGIN,
    "https://quantura-e2e3d.web.app",
    "https://quantura-e2e3d.firebaseapp.com",
    "http://localhost:5000",
    "http://127.0.0.1:5000",
    "http://localhost:5173",
    "http://127.0.0.1:5173",
}

if not firebase_admin._apps:
    options: dict[str, Any] = {}
    if STORAGE_BUCKET:
        options["storageBucket"] = STORAGE_BUCKET
    if os.path.exists(SERVICE_ACCOUNT_PATH):
        try:
            cred = credentials.Certificate(SERVICE_ACCOUNT_PATH)
            firebase_admin.initialize_app(cred, options or None)
        except Exception as exc:
            # Allow checked-in placeholder files (no private key) without breaking runtime.
            print(f"Warning: ignoring invalid service account file at {SERVICE_ACCOUNT_PATH}: {exc}")
            firebase_admin.initialize_app(options=options or None)
    else:
        # IMPORTANT: initialize_app() expects (credential, options). When running in Cloud
        # Functions, we rely on Application Default Credentials and only pass options.
        firebase_admin.initialize_app(options=options or None)

db = firestore.client()

_REMOTE_CONFIG_CACHE: dict[str, Any] = {"template": None, "loadedAt": 0.0}
DEFAULT_LLM_ALLOWED_MODELS = [
    "gpt-5-nano",
    "gpt-5-mini",
    "gpt-5",
    "gpt-5.1",
    "gpt-5.2",
]

DEFAULT_REMOTE_CONFIG: dict[str, str] = {
    "welcome_message": "Welcome to Quantura",
    "watchlist_enabled": "true",
    "forecast_prophet_enabled": "true",
    "forecast_timemixer_enabled": "true",
    "enable_social_leaderboard": "true",
    "forecast_model_primary": "Quantura Horizon",
    "promo_banner_text": "",
    "maintenance_mode": "false",
    "volatility_threshold": "0.05",
    "llm_allowed_models": ",".join(DEFAULT_LLM_ALLOWED_MODELS),
    "ai_usage_tiers": json.dumps(
        {
            "free": {
                "allowed_models": ["gpt-5-nano", "gpt-5-mini"],
                "weekly_limit": 3,
                "daily_limit": 3,
                "volatility_alerts": False,
            },
            "pro": {
                "allowed_models": ["gpt-5-mini", "gpt-5", "gpt-5.1"],
                "weekly_limit": 25,
                "daily_limit": 25,
                "volatility_alerts": True,
            },
            "desk": {
                "allowed_models": ["gpt-5-nano", "gpt-5-mini", "gpt-5", "gpt-5.1", "gpt-5.2"],
                "weekly_limit": 75,
                "daily_limit": 75,
                "volatility_alerts": True,
            },
        }
    ),
    "push_notifications_enabled": "true",
    "webpush_vapid_key": "",
    "stripe_checkout_enabled": "true",
    "stripe_public_key": "",
    "holiday_promo": "false",
    "backtesting_enabled": "true",
    "backtesting_free_daily_limit": "1",
    "backtesting_pro_daily_limit": "25",
}


def _get_remote_config_template(max_age_seconds: int = 300) -> Any | None:
    """Loads a Remote Config ServerTemplate and keeps it cached between invocations."""
    if admin_remote_config is None:
        return None
    now = time.time()
    cached = _REMOTE_CONFIG_CACHE.get("template")
    loaded_at = float(_REMOTE_CONFIG_CACHE.get("loadedAt") or 0.0)

    if cached is None:
        try:
            cached = admin_remote_config.init_server_template(default_config=DEFAULT_REMOTE_CONFIG)
            _REMOTE_CONFIG_CACHE["template"] = cached
        except Exception:
            return None

    if loaded_at and (now - loaded_at) < max_age_seconds:
        return cached

    try:
        cached.load()
        _REMOTE_CONFIG_CACHE["loadedAt"] = now
    except Exception:
        # Keep the old template cache (if any).
        pass
    return cached


def _remote_config_context(
    req: https_fn.CallableRequest | None,
    token: dict[str, Any] | None,
    meta: dict[str, Any] | None = None,
) -> dict[str, Any]:
    """Builds a stable Remote Config evaluation context.

    Note: Python ServerTemplate percent conditions expect `randomization_id`.
    """
    meta = meta or {}
    token = token or {}
    context: dict[str, Any] = {}

    randomization_id = ""
    if req and req.auth:
        randomization_id = req.auth.uid
    if not randomization_id and isinstance(meta.get("sessionId"), str) and meta.get("sessionId"):
        # For unauthenticated calls, fall back to a stable per-device session id.
        randomization_id = meta["sessionId"]
    if randomization_id:
        context["randomization_id"] = randomization_id

    email = token.get("email")
    if isinstance(email, str) and email:
        context["email"] = email.lower()
    if isinstance(meta.get("pagePath"), str) and meta.get("pagePath"):
        context["page_path"] = meta["pagePath"]
    if isinstance(meta.get("timezone"), str) and meta.get("timezone"):
        context["timezone"] = meta["timezone"]
    return context


def _remote_config_param(key: str, default: str = "", context: dict[str, Any] | None = None) -> str:
    template = _get_remote_config_template()
    if template is None:
        return default
    try:
        config = template.evaluate(context or {})
        value = config.get_string(key)
        if value is None:
            return default
        value_str = str(value)
        return value_str if value_str != "" else default
    except Exception:
        return default


def _remote_config_bool(key: str, default: bool = False, context: dict[str, Any] | None = None) -> bool:
    template = _get_remote_config_template()
    if template is None:
        return default
    try:
        config = template.evaluate(context or {})
        return bool(config.get_boolean(key))
    except Exception:
        raw = _remote_config_param(key, "", context=context)
        if not raw:
            return default
        normalized = raw.strip().lower()
        if normalized in {"1", "true", "yes", "on"}:
            return True
        if normalized in {"0", "false", "no", "off"}:
            return False
        return default


def _remote_config_int(key: str, default: int = 0, context: dict[str, Any] | None = None) -> int:
    raw = _remote_config_param(key, str(default), context=context)
    if raw == "":
        return default
    try:
        return int(float(str(raw).strip()))
    except Exception:
        return default


def _remote_config_float(key: str, default: float = 0.0, context: dict[str, Any] | None = None) -> float:
    raw = _remote_config_param(key, str(default), context=context)
    if raw == "":
        return default
    try:
        return float(str(raw).strip())
    except Exception:
        return default


def _remote_config_json_param(key: str, default: dict[str, Any], context: dict[str, Any] | None = None) -> dict[str, Any]:
    raw = _remote_config_param(key, "", context=context)
    if not raw:
        return default
    try:
        value = json.loads(raw)
        return value if isinstance(value, dict) else default
    except Exception:
        return default


def _raise_structured_error(
    code: https_fn.FunctionsErrorCode,
    error_key: str,
    detail: str,
    extras: dict[str, Any] | None = None,
) -> None:
    payload: dict[str, Any] = {"error": error_key, "detail": detail}
    if extras:
        payload.update(extras)
    raise https_fn.HttpsError(code, detail, payload)


def _normalize_ai_model_id(raw: Any) -> str:
    value = str(raw or "").strip()
    if not value:
        return ""
    lowered = value.lower()
    aliases = {
        "gpt-5-2": "gpt-5.2",
        "gpt5.2": "gpt-5.2",
        "gpt-5-1": "gpt-5.1",
        "gpt5.1": "gpt-5.1",
        "gpt5": "gpt-5",
        "gpt5-mini": "gpt-5-mini",
        "gpt5-nano": "gpt-5-nano",
        "gpt-5-thinking": "gpt-5.2",
        "nova-micro": "amazon.nova-micro-v1:0",
        "nova-lite": "amazon.nova-lite-v1:0",
        "nova-pro": "amazon.nova-pro-v1:0",
        "amazon-nova-micro": "amazon.nova-micro-v1:0",
        "amazon-nova-lite": "amazon.nova-lite-v1:0",
        "amazon-nova-pro": "amazon.nova-pro-v1:0",
    }
    if lowered in aliases:
        return aliases[lowered]
    if lowered.startswith("gpt-") and len(lowered) > 4 and lowered[4] == "4":
        return "gpt-5-mini"
    if lowered.startswith("o1"):
        return "gpt-5.1"
    if lowered.startswith("amazon.nova"):
        return lowered
    return value


def _is_supported_llm_model(model_id: str) -> bool:
    if not model_id:
        return False
    normalized = _normalize_ai_model_id(model_id)
    if not normalized:
        return False
    lowered = normalized.lower()
    return lowered.startswith("gpt-5") or lowered.startswith("amazon.nova")


def _model_provider_from_id(model_id: str) -> str:
    lowered = str(model_id or "").strip().lower()
    if lowered.startswith("amazon.nova"):
        return "amazon_nova"
    return "openai"


def _parse_llm_model_allowlist(raw: Any) -> list[str]:
    values: list[Any] = []
    if isinstance(raw, list):
        values = raw
    elif isinstance(raw, str):
        text = raw.strip()
        if not text:
            values = []
        else:
            parsed: Any = None
            try:
                parsed = json.loads(text)
            except Exception:
                parsed = None
            if isinstance(parsed, list):
                values = parsed
            else:
                values = [part.strip() for part in text.split(",")]
    normalized: list[str] = []
    for value in values:
        model_id = _normalize_ai_model_id(value)
        if not _is_supported_llm_model(model_id):
            continue
        if model_id not in normalized:
            normalized.append(model_id)
    return normalized


def _get_llm_allowed_models(context: dict[str, Any] | None = None) -> list[str]:
    defaults = [_normalize_ai_model_id(item) for item in DEFAULT_LLM_ALLOWED_MODELS if _normalize_ai_model_id(item)]
    defaults = [item for item in defaults if _is_supported_llm_model(item)]
    rc_raw = _remote_config_param("llm_allowed_models", ",".join(defaults), context=context)
    parsed = _parse_llm_model_allowlist(rc_raw)
    if not parsed:
        env_raw = str(os.environ.get("LLM_ALLOWED_MODELS") or "").strip()
        parsed = _parse_llm_model_allowlist(env_raw)
    if not parsed:
        parsed = defaults
    return parsed


def _get_ai_usage_tiers(context: dict[str, Any] | None = None) -> dict[str, Any]:
    fallback = {
        "free": {
            "allowed_models": ["gpt-5-nano", "gpt-5-mini"],
            "weekly_limit": 3,
            "daily_limit": 3,
            "volatility_alerts": False,
        },
        "pro": {
            "allowed_models": ["gpt-5-mini", "gpt-5", "gpt-5.1"],
            "weekly_limit": 25,
            "daily_limit": 25,
            "volatility_alerts": True,
        },
        "desk": {
            "allowed_models": ["gpt-5-nano", "gpt-5-mini", "gpt-5", "gpt-5.1", "gpt-5.2"],
            "weekly_limit": 75,
            "daily_limit": 75,
            "volatility_alerts": True,
        },
    }
    parsed = _remote_config_json_param("ai_usage_tiers", fallback, context=context)
    if not isinstance(parsed, dict):
        return fallback
    global_allowed = _get_llm_allowed_models(context=context)
    global_allowed_set = set(global_allowed)
    default_limits = {"free": 3, "pro": 25, "desk": 75}
    for tier_name, tier_data in parsed.items():
        if not isinstance(tier_data, dict):
            continue
        raw_models = tier_data.get("allowed_models")
        normalized: list[str] = []
        if isinstance(raw_models, list):
            for raw in raw_models:
                model_id = _normalize_ai_model_id(raw)
                if not model_id:
                    continue
                if not _is_supported_llm_model(model_id):
                    continue
                if global_allowed_set and model_id not in global_allowed_set:
                    continue
                if model_id not in normalized:
                    normalized.append(model_id)
        if not normalized:
            defaults = fallback.get(str(tier_name).lower(), fallback["free"]).get("allowed_models") or []
            for model_id in defaults:
                if not model_id:
                    continue
                normalized_id = _normalize_ai_model_id(model_id)
                if not normalized_id or (global_allowed_set and normalized_id not in global_allowed_set):
                    continue
                if normalized_id not in normalized:
                    normalized.append(normalized_id)
        if not normalized:
            for model_id in global_allowed:
                if model_id not in normalized:
                    normalized.append(model_id)

        tier_key = str(tier_name).strip().lower()
        default_limit = default_limits.get(tier_key, 3)
        raw_weekly = tier_data.get("weekly_limit")
        if raw_weekly in (None, ""):
            raw_weekly = tier_data.get("daily_limit")
        try:
            weekly_limit = int(float(raw_weekly))
        except Exception:
            weekly_limit = default_limit
        if weekly_limit <= 0:
            weekly_limit = default_limit
        if tier_key == "free":
            weekly_limit = max(1, min(weekly_limit, 5))
        elif tier_key == "pro":
            weekly_limit = max(6, min(weekly_limit, 40))
        elif tier_key == "desk":
            weekly_limit = max(20, min(weekly_limit, 200))

        tier_data["weekly_limit"] = weekly_limit
        tier_data["daily_limit"] = weekly_limit  # Backward-compatible alias for older clients.
        tier_data["volatility_alerts"] = bool(tier_data.get("volatility_alerts"))
        tier_data["allowed_models"] = normalized
    return parsed


def _resolve_paid_plan_tier(uid: str) -> str:
    """Infers a paid tier from the user's orders. Returns free/pro/desk."""
    best_rank = 0
    try:
        docs = db.collection("orders").where("userId", "==", uid).limit(50).stream()
    except Exception:
        return "free"

    for doc in docs:
        data = doc.to_dict() or {}
        status = str(data.get("paymentStatus") or "").strip().lower()
        if status not in {"paid", "succeeded"}:
            continue
        product = str(data.get("product") or "").strip().lower()
        if "desk" in product or "enterprise" in product:
            best_rank = max(best_rank, 2)
        elif "pro" in product or "elite" in product or "forecast" in product:
            best_rank = max(best_rank, 1)
        else:
            best_rank = max(best_rank, 1)

    if best_rank >= 2:
        return "desk"
    if best_rank >= 1:
        return "pro"
    return "free"


def _resolve_ai_tier(
    uid: str,
    token: dict[str, Any],
    context: dict[str, Any] | None = None,
) -> tuple[str, dict[str, Any]]:
    global_allowed = _get_llm_allowed_models(context=context)
    global_allowed_set = set(global_allowed)
    tiers = _get_ai_usage_tiers(context=context)
    if token.get("email") == ADMIN_EMAIL:
        tier_key = "desk"
    else:
        tier_key = _resolve_paid_plan_tier(uid)
        if tier_key not in {"free", "pro", "desk"}:
            tier_key = "free"
    tier = tiers.get(tier_key) if isinstance(tiers.get(tier_key), dict) else {}
    if not tier:
        fallback_limit = 3 if tier_key == "free" else 25 if tier_key == "pro" else 75
        fallback_models = {
            "free": ["gpt-5-nano", "gpt-5-mini"],
            "pro": ["gpt-5-mini", "gpt-5", "gpt-5.1"],
            "desk": ["gpt-5-nano", "gpt-5-mini", "gpt-5", "gpt-5.1", "gpt-5.2"],
        }
        tier = {
            "allowed_models": fallback_models.get(tier_key, fallback_models["free"]),
            "weekly_limit": fallback_limit,
            "daily_limit": fallback_limit,
            "volatility_alerts": tier_key != "free",
        }
    tier_allowed_raw = tier.get("allowed_models") if isinstance(tier, dict) else []
    tier_allowed: list[str] = []
    if isinstance(tier_allowed_raw, list):
        for value in tier_allowed_raw:
            model_id = _normalize_ai_model_id(value)
            if not _is_supported_llm_model(model_id):
                continue
            if global_allowed_set and model_id not in global_allowed_set:
                continue
            if model_id not in tier_allowed:
                tier_allowed.append(model_id)
    if not tier_allowed:
        tier_allowed = list(global_allowed)
    if not tier_allowed:
        tier_allowed = list(DEFAULT_LLM_ALLOWED_MODELS)
    tier["allowed_models"] = tier_allowed
    return tier_key, tier


def _yahoo_headers() -> dict[str, str]:
    # Yahoo endpoints frequently rate-limit requests without a browser-like UA.
    return {
        "User-Agent": "Mozilla/5.0 (Quantura; +https://quantura-e2e3d.web.app)",
        "Accept": "application/json,text/plain,*/*",
        "Accept-Language": "en-US,en;q=0.9",
        "Connection": "keep-alive",
    }


def _force_remove_second_line(csv_text: str) -> str:
    lines = csv_text.splitlines(True)
    if len(lines) > 1:
        del lines[1]
    return "".join(lines)


def _social_channel_webhooks() -> dict[str, str]:
    out: dict[str, str] = {}
    for channel, env_key in SOCIAL_CHANNEL_WEBHOOK_ENV.items():
        out[channel] = str(os.environ.get(env_key) or "").strip()
    return out


def _social_channel_has_webhook(channel: str) -> bool:
    key = str(channel or "").strip().lower()
    if not key:
        return False
    return bool(_social_channel_webhooks().get(key))


def _social_channel_has_direct_credentials(channel: str) -> bool:
    key = str(channel or "").strip().lower()
    if key == "x":
        return bool(TWITTER_API_KEY and TWITTER_API_SECRET and TWITTER_ACCESS_TOKEN and TWITTER_ACCESS_TOKEN_SECRET)
    if key == "linkedin":
        return bool(LINKEDIN_ACCESS_TOKEN and LINKEDIN_AUTHOR_URN)
    if key == "facebook":
        return bool(FACEBOOK_PAGE_ID and FACEBOOK_PAGE_ACCESS_TOKEN)
    if key == "instagram":
        return bool(INSTAGRAM_BUSINESS_ACCOUNT_ID and INSTAGRAM_ACCESS_TOKEN and INSTAGRAM_DEFAULT_IMAGE_URL)
    if key == "tiktok":
        # Direct text-only posting is not supported in this implementation.
        return False
    return False


def _social_channel_posting_route(channel: str) -> str:
    key = str(channel or "").strip().lower()
    has_direct = _social_channel_has_direct_credentials(key)
    has_webhook = _social_channel_has_webhook(key)
    if has_direct and has_webhook:
        return "direct+webhook"
    if has_direct:
        return "direct"
    if has_webhook:
        return "webhook"
    return "none"


def _configured_social_posting_channels(raw_channels: Any = None) -> list[str]:
    channels = _normalize_social_channels(raw_channels)
    return [channel for channel in channels if _social_channel_posting_route(channel) != "none"]


def _normalize_social_channels(raw: Any) -> list[str]:
    if raw is None:
        return list(SOCIAL_POPULAR_CHANNELS)
    if isinstance(raw, str):
        parts = [part.strip().lower() for part in raw.split(",")]
    elif isinstance(raw, list):
        parts = [str(part).strip().lower() for part in raw]
    else:
        parts = []
    channels = [channel for channel in parts if channel in SOCIAL_POPULAR_CHANNELS]
    if not channels:
        return list(SOCIAL_POPULAR_CHANNELS)
    return list(dict.fromkeys(channels))


def _posting_tzinfo():
    if ZoneInfo is not None:
        try:
            return ZoneInfo(SOCIAL_POSTING_TIMEZONE)
        except Exception:
            pass
    return timezone.utc


def _strategic_slot_minutes(channel: str) -> list[int]:
    # Local-time slots chosen for US market participation windows.
    slots = {
        "x": [9 * 60 + 5, 12 * 60 + 20, 16 * 60 + 40],
        "linkedin": [8 * 60 + 45, 12 * 60 + 10, 17 * 60 + 20],
        "facebook": [10 * 60 + 15, 14 * 60 + 35, 19 * 60 + 15],
        "instagram": [11 * 60 + 30, 18 * 60 + 0, 21 * 60 + 10],
        "threads": [12 * 60 + 0, 18 * 60 + 30],
        "reddit": [13 * 60 + 15, 20 * 60 + 0],
        "tiktok": [13 * 60 + 0, 19 * 60 + 30, 22 * 60 + 5],
        "youtube": [17 * 60 + 0],
        "pinterest": [13 * 60 + 45, 20 * 60 + 30],
    }
    key = str(channel or "").strip().lower()
    return slots.get(key) or [14 * 60 + 10]


def _strategic_suggested_post_time(channel: str, index: int, now_utc: datetime | None = None) -> str:
    now_utc = now_utc or datetime.now(timezone.utc)
    local_tz = _posting_tzinfo()
    now_local = now_utc.astimezone(local_tz)
    slots = _strategic_slot_minutes(channel)
    if not slots:
        slots = [14 * 60 + 10]
    slot_idx = max(0, int(index or 0)) % len(slots)
    day_offset = max(0, int(index or 0)) // len(slots)
    base_date = now_local.date() + timedelta(days=day_offset)
    minute_of_day = slots[slot_idx]
    hour = minute_of_day // 60
    minute = minute_of_day % 60
    local_candidate = datetime(
        base_date.year,
        base_date.month,
        base_date.day,
        hour,
        minute,
        tzinfo=local_tz,
    )
    if day_offset == 0 and local_candidate <= now_local + timedelta(minutes=8):
        local_candidate += timedelta(days=1)
    return local_candidate.astimezone(timezone.utc).replace(second=0, microsecond=0).isoformat().replace("+00:00", "Z")


def _render_social_text(
    *,
    body: str,
    headline: str,
    hashtags: list[str],
    cta: str,
    cta_url: str,
) -> tuple[str, list[str]]:
    tags = [str(tag).strip() for tag in hashtags if str(tag).strip()]
    rendered_text = body.strip()
    if headline.strip():
        rendered_text = f"{headline.strip()}\n\n{rendered_text}".strip()
    if tags:
        rendered_text = f"{rendered_text}\n\n{' '.join(tags)}".strip()
    if cta.strip() and cta_url.strip():
        rendered_text = f"{rendered_text}\n\n{cta.strip()}: {cta_url.strip()}".strip()
    return rendered_text, tags


def _post_x_direct(text: str) -> dict[str, Any]:
    has_oauth_user_context = bool(
        TWITTER_API_KEY and TWITTER_API_SECRET and TWITTER_ACCESS_TOKEN and TWITTER_ACCESS_TOKEN_SECRET
    )
    has_bearer = bool(TWITTER_BEARER_TOKEN)
    if not has_oauth_user_context and not has_bearer:
        return {
            "ok": False,
            "pendingCredentials": True,
            "error": (
                "Missing X credentials. Provide TWITTER_API_KEY, TWITTER_API_SECRET, "
                "TWITTER_ACCESS_TOKEN, TWITTER_ACCESS_TOKEN_SECRET (preferred) or TWITTER_BEARER_TOKEN."
            ),
        }
    try:
        auth = None
        headers = {"Content-Type": "application/json"}
        if has_oauth_user_context:
            try:
                from requests_oauthlib import OAuth1  # type: ignore
            except Exception:
                return {
                    "ok": False,
                    "pendingCredentials": False,
                    "error": "requests-oauthlib is required for X user-context posting.",
                }
            auth = OAuth1(
                TWITTER_API_KEY,
                TWITTER_API_SECRET,
                TWITTER_ACCESS_TOKEN,
                TWITTER_ACCESS_TOKEN_SECRET,
            )
        elif has_bearer:
            headers["Authorization"] = f"Bearer {TWITTER_BEARER_TOKEN}"

        response = requests.post(
            "https://api.twitter.com/2/tweets",
            headers=headers,
            auth=auth,
            json={"text": text},
            timeout=20,
        )
        if response.status_code >= 400:
            err_text = str(response.text or "")[:300]
            err_lower = err_text.lower()
            pending = response.status_code in {401, 403} and (
                "unsupported authentication" in err_lower
                or "application-only is forbidden" in err_lower
                or "oauth 2.0 application-only" in err_lower
            )
            return {
                "ok": False,
                "pendingCredentials": pending or has_bearer,
                "error": (
                    "X posting requires OAuth user-context write tokens (TWITTER_ACCESS_TOKEN / "
                    "TWITTER_ACCESS_TOKEN_SECRET) or a posting webhook."
                    if pending
                    else f"X API {response.status_code}: {err_text}"
                ),
            }
        body = response.json() if response.content else {}
        post_id = str((body.get("data") or {}).get("id") or "").strip()
        return {
            "ok": True,
            "externalId": post_id,
            "statusCode": response.status_code,
            "provider": "x_api_oauth1" if has_oauth_user_context else "x_api_bearer",
        }
    except Exception as exc:
        return {"ok": False, "pendingCredentials": False, "error": f"X API exception: {str(exc)[:300]}"}


def _post_linkedin_direct(text: str) -> dict[str, Any]:
    if not LINKEDIN_ACCESS_TOKEN or not LINKEDIN_AUTHOR_URN:
        return {
            "ok": False,
            "pendingCredentials": True,
            "error": "Missing LINKEDIN_ACCESS_TOKEN or LINKEDIN_AUTHOR_URN for LinkedIn posting.",
        }
    payload = {
        "author": LINKEDIN_AUTHOR_URN,
        "lifecycleState": "PUBLISHED",
        "specificContent": {
            "com.linkedin.ugc.ShareContent": {
                "shareCommentary": {"text": text},
                "shareMediaCategory": "NONE",
            }
        },
        "visibility": {"com.linkedin.ugc.MemberNetworkVisibility": "PUBLIC"},
    }
    try:
        response = requests.post(
            "https://api.linkedin.com/v2/ugcPosts",
            headers={
                "Authorization": f"Bearer {LINKEDIN_ACCESS_TOKEN}",
                "Content-Type": "application/json",
                "X-Restli-Protocol-Version": "2.0.0",
            },
            json=payload,
            timeout=20,
        )
        if response.status_code >= 400:
            return {
                "ok": False,
                "pendingCredentials": False,
                "error": f"LinkedIn API {response.status_code}: {response.text[:300]}",
            }
        external_id = str(response.headers.get("x-restli-id") or "").strip()
        return {"ok": True, "externalId": external_id, "statusCode": response.status_code, "provider": "linkedin_api"}
    except Exception as exc:
        return {"ok": False, "pendingCredentials": False, "error": f"LinkedIn API exception: {str(exc)[:300]}"}


def _post_facebook_direct(text: str) -> dict[str, Any]:
    if not FACEBOOK_PAGE_ID or not FACEBOOK_PAGE_ACCESS_TOKEN:
        return {
            "ok": False,
            "pendingCredentials": True,
            "error": "Missing FACEBOOK_PAGE_ID or FACEBOOK_PAGE_ACCESS_TOKEN for Facebook posting.",
        }
    try:
        response = requests.post(
            f"https://graph.facebook.com/{META_GRAPH_API_VERSION}/{FACEBOOK_PAGE_ID}/feed",
            data={"message": text, "access_token": FACEBOOK_PAGE_ACCESS_TOKEN},
            timeout=20,
        )
        if response.status_code >= 400:
            return {
                "ok": False,
                "pendingCredentials": False,
                "error": f"Facebook API {response.status_code}: {response.text[:300]}",
            }
        body = response.json() if response.content else {}
        external_id = str(body.get("id") or "").strip()
        return {"ok": True, "externalId": external_id, "statusCode": response.status_code, "provider": "facebook_api"}
    except Exception as exc:
        return {"ok": False, "pendingCredentials": False, "error": f"Facebook API exception: {str(exc)[:300]}"}


def _post_instagram_direct(text: str) -> dict[str, Any]:
    if not INSTAGRAM_BUSINESS_ACCOUNT_ID or not INSTAGRAM_ACCESS_TOKEN:
        return {
            "ok": False,
            "pendingCredentials": True,
            "error": "Missing INSTAGRAM_BUSINESS_ACCOUNT_ID or INSTAGRAM_ACCESS_TOKEN for Instagram posting.",
        }
    if not INSTAGRAM_DEFAULT_IMAGE_URL:
        return {"ok": False, "pendingCredentials": True, "error": "Missing INSTAGRAM_DEFAULT_IMAGE_URL for Instagram posts."}
    try:
        create_resp = requests.post(
            f"https://graph.facebook.com/{META_GRAPH_API_VERSION}/{INSTAGRAM_BUSINESS_ACCOUNT_ID}/media",
            data={
                "image_url": INSTAGRAM_DEFAULT_IMAGE_URL,
                "caption": text,
                "access_token": INSTAGRAM_ACCESS_TOKEN,
            },
            timeout=25,
        )
        if create_resp.status_code >= 400:
            return {
                "ok": False,
                "pendingCredentials": False,
                "error": f"Instagram media create {create_resp.status_code}: {create_resp.text[:300]}",
            }
        creation_id = str((create_resp.json() or {}).get("id") or "").strip()
        if not creation_id:
            return {"ok": False, "pendingCredentials": False, "error": "Instagram media create did not return an id."}

        publish_resp = requests.post(
            f"https://graph.facebook.com/{META_GRAPH_API_VERSION}/{INSTAGRAM_BUSINESS_ACCOUNT_ID}/media_publish",
            data={
                "creation_id": creation_id,
                "access_token": INSTAGRAM_ACCESS_TOKEN,
            },
            timeout=25,
        )
        if publish_resp.status_code >= 400:
            return {
                "ok": False,
                "pendingCredentials": False,
                "error": f"Instagram publish {publish_resp.status_code}: {publish_resp.text[:300]}",
            }
        body = publish_resp.json() if publish_resp.content else {}
        external_id = str(body.get("id") or "").strip()
        return {"ok": True, "externalId": external_id, "statusCode": publish_resp.status_code, "provider": "instagram_api"}
    except Exception as exc:
        return {"ok": False, "pendingCredentials": False, "error": f"Instagram API exception: {str(exc)[:300]}"}


def _post_tiktok_direct(text: str) -> dict[str, Any]:
    # TikTok's content publishing flow requires pre-uploaded media and user-granted scopes.
    if not TIKTOK_ACCESS_TOKEN or not TIKTOK_OPEN_ID:
        return {
            "ok": False,
            "pendingCredentials": True,
            "error": "Missing TIKTOK_ACCESS_TOKEN or TIKTOK_OPEN_ID for TikTok posting.",
        }
    return {
        "ok": False,
        "pendingCredentials": True,
        "error": "TikTok direct text posting is not supported by this endpoint. Use SOCIAL_WEBHOOK_TIKTOK or media-publish integration.",
    }


def _as_utc_datetime(raw: Any) -> datetime | None:
    if raw is None:
        return None
    if isinstance(raw, datetime):
        return raw.astimezone(timezone.utc) if raw.tzinfo else raw.replace(tzinfo=timezone.utc)
    if hasattr(raw, "to_datetime"):
        try:
            dt = raw.to_datetime()  # Firestore Timestamp
            if isinstance(dt, datetime):
                return dt.astimezone(timezone.utc) if dt.tzinfo else dt.replace(tzinfo=timezone.utc)
        except Exception:
            return None
    try:
        text = str(raw).strip()
        if not text:
            return None
        normalized = text.replace("Z", "+00:00")
        dt = datetime.fromisoformat(normalized)
        return dt.astimezone(timezone.utc) if dt.tzinfo else dt.replace(tzinfo=timezone.utc)
    except Exception:
        return None


def _extract_json_object(raw_text: str) -> dict[str, Any]:
    text = str(raw_text or "").strip()
    if not text:
        return {}
    try:
        parsed = json.loads(text)
        if isinstance(parsed, dict):
            return parsed
    except Exception:
        pass

    fenced = re.search(r"```(?:json)?\\s*(\\{.*\\})\\s*```", text, flags=re.DOTALL)
    if fenced:
        try:
            parsed = json.loads(fenced.group(1))
            if isinstance(parsed, dict):
                return parsed
        except Exception:
            pass

    start = text.find("{")
    end = text.rfind("}")
    if start >= 0 and end > start:
        try:
            parsed = json.loads(text[start : end + 1])
            if isinstance(parsed, dict):
                return parsed
        except Exception:
            return {}
    return {}


def _extract_responses_output_text(payload: dict[str, Any]) -> str:
    output = payload.get("output")
    if not isinstance(output, list):
        return ""
    chunks: list[str] = []
    for item in output:
        if not isinstance(item, dict):
            continue
        if item.get("type") != "message":
            continue
        content = item.get("content")
        if not isinstance(content, list):
            continue
        for part in content:
            if not isinstance(part, dict):
                continue
            part_type = str(part.get("type") or "").strip().lower()
            if part_type in {"output_text", "text"}:
                text = str(part.get("text") or "").strip()
                if text:
                    chunks.append(text)
    return "\n".join(chunks).strip()


def _extract_nova_output_text(payload: dict[str, Any]) -> str:
    if not isinstance(payload, dict):
        return ""
    direct = str(payload.get("outputText") or payload.get("completion") or payload.get("text") or "").strip()
    if direct:
        return direct

    output = payload.get("output")
    if isinstance(output, dict):
        message = output.get("message")
        if isinstance(message, dict):
            content = message.get("content")
            if isinstance(content, list):
                pieces: list[str] = []
                for part in content:
                    if isinstance(part, dict):
                        text = str(part.get("text") or "").strip()
                        if text:
                            pieces.append(text)
                if pieces:
                    return "\n".join(pieces).strip()

    choices = payload.get("choices")
    if isinstance(choices, list) and choices:
        first = choices[0] if isinstance(choices[0], dict) else {}
        message = first.get("message") if isinstance(first, dict) else {}
        if isinstance(message, dict):
            content = message.get("content")
            if isinstance(content, str):
                return content.strip()
            if isinstance(content, list):
                pieces: list[str] = []
                for part in content:
                    if isinstance(part, dict):
                        text = str(part.get("text") or "").strip()
                        if text:
                            pieces.append(text)
                if pieces:
                    return "\n".join(pieces).strip()
    return ""


def _invoke_amazon_nova_text(
    *,
    model_id: str,
    system_prompt: str,
    user_prompt: str,
    max_tokens: int = 800,
    temperature: float = 0.2,
) -> str:
    endpoint = AMAZON_NOVA_API_ENDPOINT
    if not endpoint or not AMAZON_NOVA_API_KEY:
        return ""

    resolved_model = _normalize_ai_model_id(model_id or AMAZON_NOVA_DEFAULT_MODEL)
    if not resolved_model.lower().startswith("amazon.nova"):
        resolved_model = _normalize_ai_model_id(AMAZON_NOVA_DEFAULT_MODEL)

    payload = {
        "model": resolved_model,
        "messages": [
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": user_prompt},
        ],
        "max_tokens": max(128, min(int(max_tokens), 2400)),
        "temperature": max(0.0, min(float(temperature), 1.0)),
    }
    headers = {
        "Content-Type": "application/json",
        # Keep both auth headers for compatibility with API-gateway and OpenAI-compatible Nova proxies.
        "Authorization": f"Bearer {AMAZON_NOVA_API_KEY}",
        "x-api-key": AMAZON_NOVA_API_KEY,
    }
    try:
        response = requests.post(endpoint, headers=headers, json=payload, timeout=30)
        response.raise_for_status()
        body = response.json() if response.text else {}
        return _extract_nova_output_text(body)
    except Exception:
        return ""


def _normalize_symbol_token(raw: Any) -> str:
    text = str(raw or "").strip().upper()
    if text.startswith("$"):
        text = text[1:]
    text = re.sub(r"[^A-Z0-9\.\-]", "", text)
    if not text:
        return ""
    if len(text) > 10:
        return ""
    if not re.fullmatch(r"[A-Z][A-Z0-9\.\-]{0,9}", text):
        return ""
    return text


def _extract_symbols_from_text(raw_text: str) -> list[str]:
    text = str(raw_text or "")
    if not text:
        return []
    candidates = re.findall(r"\$?[A-Za-z][A-Za-z0-9\.\-]{0,9}", text)
    symbols: list[str] = []
    blocked = {
        "THE",
        "THIS",
        "THAT",
        "WITH",
        "FROM",
        "WHAT",
        "ABOUT",
        "PORTFOLIO",
        "FAVORITE",
        "STOCK",
        "STOCKS",
    }
    for token in candidates:
        symbol = _normalize_symbol_token(token)
        if not symbol or symbol in blocked:
            continue
        if symbol not in symbols:
            symbols.append(symbol)
    return symbols


def _yahoo_search_symbols(query: str, *, max_results: int = 8) -> list[str]:
    q = str(query or "").strip()
    if not q:
        return []
    try:
        response = requests.get(
            YAHOO_SEARCH_URL,
            headers=_yahoo_headers(),
            params={"q": q, "quotesCount": max(1, min(max_results, 20)), "newsCount": 0},
            timeout=12,
        )
        response.raise_for_status()
        payload = response.json()
    except Exception:
        return []

    out: list[str] = []
    for quote in payload.get("quotes", []) or []:
        if not isinstance(quote, dict):
            continue
        symbol = _normalize_symbol_token(quote.get("symbol"))
        if not symbol:
            continue
        quote_type = str(quote.get("quoteType") or "").strip().upper()
        if quote_type and quote_type not in {"EQUITY", "ETF", "MUTUALFUND"}:
            continue
        if symbol not in out:
            out.append(symbol)
        if len(out) >= max_results:
            break
    return out


def _normalize_country_code(raw: Any) -> str:
    text = str(raw or "").strip().upper()
    if not text:
        return "US"
    alias = {
        "USA": "US",
        "UNITED STATES": "US",
        "UNITED STATES OF AMERICA": "US",
        "U.S.": "US",
        "U.S.A.": "US",
        "CANADA": "CA",
        "UNITED KINGDOM": "GB",
        "UK": "GB",
        "GREAT BRITAIN": "GB",
        "GERMANY": "DE",
        "FRANCE": "FR",
        "JAPAN": "JP",
        "CHINA": "CN",
        "HONG KONG": "HK",
        "INDIA": "IN",
        "AUSTRALIA": "AU",
        "BRAZIL": "BR",
        "MEXICO": "MX",
        "SOUTH KOREA": "KR",
        "TAIWAN": "TW",
        "SINGAPORE": "SG",
    }
    if text in alias:
        return alias[text]
    if len(text) == 2 and text.isalpha():
        return text
    return "US"


def _is_us_equity_symbol(symbol: str) -> bool:
    try:
        import yfinance as yf  # type: ignore

        info = yf.Ticker(str(symbol or "").strip().upper()).info or {}
        if not isinstance(info, dict):
            return False
        country = str(info.get("country") or "").strip().lower()
        if country in {"united states", "usa", "us"}:
            return True
        exchange = str(info.get("exchange") or info.get("fullExchangeName") or "").strip().lower()
        if any(tag in exchange for tag in ["nasdaq", "nyse", "amex", "cboe", "arca"]):
            return True
    except Exception:
        return False
    return False


def _coerce_event_iso_dates(value: Any) -> list[str]:
    values: list[Any]
    if value is None:
        return []
    if isinstance(value, (list, tuple)):
        values = list(value)
    else:
        values = [value]

    out: list[str] = []
    for item in values:
        if item is None:
            continue
        try:
            if isinstance(item, datetime):
                out.append(item.date().isoformat())
                continue
            if isinstance(item, date):
                out.append(item.isoformat())
                continue
            text = str(item).strip()
            if not text:
                continue
            if "T" in text:
                parsed = datetime.fromisoformat(text.replace("Z", "+00:00"))
                out.append(parsed.date().isoformat())
                continue
            parsed_date = datetime.fromisoformat(text[:10]).date()
            out.append(parsed_date.isoformat())
        except Exception:
            continue
    return list(dict.fromkeys(out))


def _fetch_yahoo_corporate_events_for_ticker(
    ticker: str,
    *,
    start_date: date,
    end_date: date,
) -> list[dict[str, Any]]:
    import yfinance as yf  # type: ignore

    symbol = _normalize_symbol_token(ticker)
    if not symbol:
        return []

    out: list[dict[str, Any]] = []
    tk = yf.Ticker(symbol)
    info: dict[str, Any] = {}
    try:
        raw_info = tk.info or {}
        info = raw_info if isinstance(raw_info, dict) else {}
    except Exception:
        info = {}

    company_name = str(info.get("longName") or info.get("shortName") or symbol).strip()
    country = str(info.get("country") or "").strip()
    exchange = str(info.get("fullExchangeName") or info.get("exchange") or "").strip()

    seen: set[str] = set()
    try:
        calendar_raw = tk.calendar or {}
        if isinstance(calendar_raw, dict):
            for key, raw_value in calendar_raw.items():
                label = str(key or "").strip() or "Corporate Event"
                for dt_iso in _coerce_event_iso_dates(raw_value):
                    try:
                        dt_obj = datetime.fromisoformat(dt_iso).date()
                    except Exception:
                        continue
                    if dt_obj < start_date or dt_obj > end_date:
                        continue
                    event_id = f"{symbol}_{dt_iso}_{label}".lower().replace(" ", "_")
                    if event_id in seen:
                        continue
                    seen.add(event_id)
                    out.append(
                        {
                            "id": event_id,
                            "ticker": symbol,
                            "companyName": company_name,
                            "date": dt_iso,
                            "type": str(label).lower().replace(" ", "_"),
                            "status": "scheduled",
                            "name": label,
                            "country": country,
                            "exchange": exchange,
                            "source": "yahoo_finance",
                        }
                    )
    except Exception:
        pass

    try:
        earnings = tk.get_earnings_dates(limit=8)
        if earnings is not None and getattr(earnings, "empty", True) is False:
            for idx, row in earnings.reset_index().iterrows():
                raw_dt = row.get("Earnings Date") or row.get("Date") or row.get("date")
                for dt_iso in _coerce_event_iso_dates(raw_dt):
                    try:
                        dt_obj = datetime.fromisoformat(dt_iso).date()
                    except Exception:
                        continue
                    if dt_obj < start_date or dt_obj > end_date:
                        continue
                    event_id = f"{symbol}_{dt_iso}_earnings_announcement_date_{idx}".lower()
                    if event_id in seen:
                        continue
                    seen.add(event_id)
                    out.append(
                        {
                            "id": event_id,
                            "ticker": symbol,
                            "companyName": company_name,
                            "date": dt_iso,
                            "type": "earnings_announcement_date",
                            "status": "scheduled",
                            "name": "Earnings Announcement Date",
                            "country": country,
                            "exchange": exchange,
                            "source": "yahoo_finance",
                        }
                    )
    except Exception:
        pass

    out.sort(key=lambda row: (str(row.get("date") or ""), str(row.get("ticker") or "")))
    return out


def _fetch_massive_corporate_events(
    *,
    tickers: list[str],
    start_date: date,
    end_date: date,
    limit: int = 200,
    event_types: list[str] | None = None,
    statuses: list[str] | None = None,
) -> tuple[list[dict[str, Any]], str]:
    if not MASSIVE_API_KEY:
        return [], "Massive TMX API key is not configured."

    normalized = [_normalize_symbol_token(t) for t in (tickers or [])]
    normalized = [t for t in normalized if t]
    if not normalized:
        return [], "No valid tickers were provided."

    params: dict[str, Any] = {
        "date.gte": start_date.isoformat(),
        "date.lte": end_date.isoformat(),
        "sort": "date.asc",
        "limit": max(1, min(int(limit or 200), 1000)),
        "apiKey": MASSIVE_API_KEY,
    }
    if len(normalized) == 1:
        params["ticker"] = normalized[0]
    else:
        params["ticker.any_of"] = ",".join(normalized[:100])

    if event_types:
        clean_types = [str(item).strip() for item in event_types if str(item).strip()]
        if clean_types:
            params["type.any_of"] = ",".join(clean_types[:30])
    if statuses:
        clean_statuses = [str(item).strip() for item in statuses if str(item).strip()]
        if clean_statuses:
            params["status.any_of"] = ",".join(clean_statuses[:20])

    url = f"{MASSIVE_BASE_URL}/tmx/v1/corporate-events"
    try:
        response = requests.get(url, params=params, timeout=18)
        response.raise_for_status()
        payload = response.json() if response.text else {}
        rows = payload.get("results") or []
    except Exception as exc:
        return [], f"Massive TMX request failed: {str(exc)[:180]}"

    out: list[dict[str, Any]] = []
    for idx, row in enumerate(rows):
        if not isinstance(row, dict):
            continue
        ticker = _normalize_symbol_token(row.get("ticker"))
        if not ticker:
            continue
        dt = str(row.get("date") or "").strip()
        if not dt:
            continue
        event_id = str(row.get("tmx_record_id") or "").strip() or f"{ticker}_{dt}_{idx}"
        out.append(
            {
                "id": event_id,
                "ticker": ticker,
                "companyName": str(row.get("company_name") or row.get("name") or ticker).strip(),
                "date": dt[:10],
                "type": str(row.get("type") or "").strip(),
                "status": str(row.get("status") or "").strip() or "scheduled",
                "name": str(row.get("name") or row.get("type") or "Corporate Event").strip(),
                "country": "United States",
                "exchange": str(row.get("exchange") or "").strip(),
                "source": "massive_tmx",
            }
        )
    out.sort(key=lambda item: (str(item.get("date") or ""), str(item.get("ticker") or "")))
    return out[: max(1, min(int(limit or 200), 400))], ""


def _fetch_yahoo_news_query(query: str, *, limit: int = 12) -> list[dict[str, Any]]:
    q = str(query or "").strip()
    if not q:
        return []

    def _thumbnail(item: dict[str, Any]) -> str:
        thumb = item.get("thumbnail")
        if isinstance(thumb, dict):
            resolutions = thumb.get("resolutions")
            if isinstance(resolutions, list):
                best_url = ""
                best_width = -1
                for res in resolutions:
                    if not isinstance(res, dict):
                        continue
                    url = str(res.get("url") or "").strip()
                    width = int(res.get("width") or 0)
                    if url and width >= best_width:
                        best_width = width
                        best_url = url
                return best_url
        return ""

    seen: set[str] = set()
    out: list[dict[str, Any]] = []
    try:
        response = requests.get(
            YAHOO_SEARCH_URL,
            headers=_yahoo_headers(),
            params={"q": q, "newsCount": max(1, min(int(limit or 12), 50)), "quotesCount": 0, "listsCount": 0},
            timeout=12,
        )
        if response.status_code >= 400:
            return []
        payload = response.json() if response.text else {}
        for item in payload.get("news") or []:
            if not isinstance(item, dict):
                continue
            title = str(item.get("title") or "").strip()
            link = str(item.get("link") or item.get("url") or "").strip()
            if not title:
                continue
            key = (link or title).lower()
            if key in seen:
                continue
            seen.add(key)
            out.append(
                {
                    "title": title,
                    "publisher": str(item.get("publisher") or item.get("source") or "").strip(),
                    "link": link,
                    "publishedAt": item.get("providerPublishTime") or item.get("publishedAt"),
                    "summary": str(item.get("summary") or item.get("description") or "").strip(),
                    "thumbnailUrl": _thumbnail(item),
                    "source": "yahoo_finance",
                }
            )
            if len(out) >= limit:
                break
    except Exception:
        return []

    return out


def _fetch_unsplash_random_photos(
    query: str,
    *,
    count: int = 1,
) -> tuple[list[dict[str, Any]], str]:
    q = str(query or "").strip()
    if not q:
        return [], "Query is required."
    if not UNSPLASH_ACCESS_KEY:
        return [], "Unsplash access key is not configured."

    try:
        response = requests.get(
            "https://api.unsplash.com/photos/random",
            headers={
                "Accept-Version": "v1",
                "Authorization": f"Client-ID {UNSPLASH_ACCESS_KEY}",
            },
            params={
                "query": q,
                "orientation": "landscape",
                "content_filter": "high",
                "count": max(1, min(int(count or 1), 8)),
            },
            timeout=12,
        )
        if response.status_code >= 400:
            return [], f"Unsplash API returned HTTP {response.status_code}."
        payload = response.json() if response.text else {}
    except Exception as exc:
        return [], f"Unsplash request failed: {str(exc)[:160]}"

    items = payload if isinstance(payload, list) else [payload]
    photos: list[dict[str, Any]] = []
    for item in items:
        if not isinstance(item, dict):
            continue
        url = (
            str(((item.get("urls") or {}) if isinstance(item.get("urls"), dict) else {}).get("regular") or "")
            .strip()
        )
        if not url:
            url = (
                str(((item.get("urls") or {}) if isinstance(item.get("urls"), dict) else {}).get("full") or "")
                .strip()
            )
        if not url:
            url = (
                str(((item.get("urls") or {}) if isinstance(item.get("urls"), dict) else {}).get("small") or "")
                .strip()
            )
        if not url:
            continue

        link = ""
        links_obj = item.get("links") if isinstance(item.get("links"), dict) else {}
        if isinstance(links_obj, dict):
            link = str(links_obj.get("html") or "").strip()
        if link:
            link = f"{link}{'&' if '?' in link else '?'}utm_source=quantura&utm_medium=referral"

        user_obj = item.get("user") if isinstance(item.get("user"), dict) else {}
        photographer = str(user_obj.get("name") or "").strip() if isinstance(user_obj, dict) else ""
        photographer_link = ""
        if isinstance(user_obj, dict):
            ulinks = user_obj.get("links") if isinstance(user_obj.get("links"), dict) else {}
            if isinstance(ulinks, dict):
                photographer_link = str(ulinks.get("html") or "").strip()
        if photographer_link:
            photographer_link = f"{photographer_link}{'&' if '?' in photographer_link else '?'}utm_source=quantura&utm_medium=referral"

        photos.append(
            {
                "url": url,
                "alt": str(item.get("alt_description") or item.get("description") or "Market imagery from Unsplash").strip(),
                "link": link,
                "photographer": photographer,
                "photographerLink": photographer_link,
            }
        )

    if photos:
        return photos, ""
    return [], "Unsplash returned no usable photos."


_SOCIAL_QUERY_STOPWORDS = {
    "the",
    "and",
    "or",
    "for",
    "with",
    "from",
    "this",
    "that",
    "today",
    "top",
    "news",
    "headlines",
    "stock",
    "stocks",
    "market",
}


def _social_query_terms(query: str) -> list[str]:
    tokens = re.findall(r"[a-z0-9$\.]{2,}", str(query or "").lower())
    out: list[str] = []
    for token in tokens:
        if token in _SOCIAL_QUERY_STOPWORDS:
            continue
        if token not in out:
            out.append(token)
    return out[:8]


def _text_matches_social_query(text: str, query: str) -> bool:
    body = str(text or "").lower()
    terms = _social_query_terms(query)
    if not terms:
        return True
    hits = sum(1 for term in terms if term in body)
    # Soft match: one term for short queries, two terms for longer multi-word queries.
    needed = 1 if len(terms) <= 3 else 2
    return hits >= needed


def _fetch_social_posts_via_web_search(
    query: str,
    *,
    platform: str,
    limit: int = 8,
) -> tuple[list[dict[str, Any]], str]:
    q = str(query or "").strip()
    if not q:
        return [], "Web fallback query is empty."

    domain_map = {
        "x": "x.com",
        "reddit": "reddit.com",
        "facebook": "facebook.com",
        "instagram": "instagram.com",
    }
    domain = domain_map.get(str(platform or "").strip().lower())
    if not domain:
        return [], "Unsupported web fallback platform."

    def _normalize_search_href(raw_href: str) -> str:
        href = unescape(str(raw_href or "")).strip()
        if not href:
            return ""
        if href.startswith("//"):
            href = f"https:{href}"
        href = href.replace("&amp;", "&")
        try:
            parsed = urlparse(href)
            if "duckduckgo.com" in parsed.netloc.lower():
                for key in ("uddg", "rut", "u"):
                    value = parse_qs(parsed.query).get(key, [""])[0]
                    if value:
                        href = unquote(value)
                        parsed = urlparse(href)
                        break
        except Exception:
            pass
        return href

    def _href_matches_domain(href: str) -> bool:
        try:
            parsed = urlparse(str(href or "").strip())
            host = parsed.netloc.lower().strip()
            if host.startswith("www."):
                host = host[4:]
            return bool(host and (host == domain or host.endswith(f".{domain}")))
        except Exception:
            return False

    def _social_permalink_score(target_platform: str, href: str) -> int:
        href_l = str(href or "").lower()
        if not _href_matches_domain(href):
            return -10
        score = 1
        if target_platform == "x":
            if re.search(r"x\.com/[^/?#]+/status/\d+", href_l):
                score += 10
            elif "/i/web/status/" in href_l:
                score += 8
            elif re.search(r"x\.com/[^/?#]+/?$", href_l):
                score += 3
            if any(
                bad in href_l
                for bad in (
                    "about.x.com",
                    "help.x.com",
                    "docs.x.com",
                    "business.x.com",
                    "/search?",
                    "/hashtag/",
                )
            ):
                score -= 7
        elif target_platform == "reddit":
            if "/comments/" in href_l:
                score += 9
            elif "/r/" in href_l:
                score += 5
            elif "/user/" in href_l or "/u/" in href_l:
                score += 2
        elif target_platform == "facebook":
            if "/posts/" in href_l or "/permalink.php" in href_l:
                score += 8
            elif "/reel/" in href_l or "/watch/" in href_l:
                score += 6
            elif "/groups/" in href_l:
                score += 3
        elif target_platform == "instagram":
            if "/p/" in href_l or "/reel/" in href_l or "/tv/" in href_l:
                score += 9
            elif re.search(r"instagram\.com/[^/?#]+/?$", href_l):
                score += 3
        if href_l.endswith((".ico", ".svg", ".png", ".jpg", ".jpeg", ".webp", ".css", ".js")):
            score -= 12
        if href_l.rstrip("/") in {f"https://{domain}", f"http://{domain}"}:
            score -= 6
        return score

    def _extract_social_rows(text_blob: str) -> list[dict[str, Any]]:
        if not text_blob:
            return []
        max_rows = max(1, min(int(limit or 8), 20))
        candidates: list[tuple[int, dict[str, Any]]] = []
        seen_links: set[str] = set()

        def _push_row(raw_href: str, raw_title: str, context_text: str) -> None:
            href = _normalize_search_href(raw_href)
            if not href or href in seen_links:
                return
            href_l = href.lower()
            if not _href_matches_domain(href):
                return
            score = _social_permalink_score(platform, href)
            if score <= 0:
                return
            title = unescape(re.sub(r"<[^>]+>", "", str(raw_title or ""))).strip()
            context = unescape(re.sub(r"<[^>]+>", "", str(context_text or ""))).strip()
            text_probe = f"{title} {context} {href}".strip()
            matches = _text_matches_social_query(text_probe, q)
            if not matches and score < 8:
                return
            if not title:
                parsed = urlparse(href)
                title = (parsed.path or parsed.netloc or platform).strip("/")
            row = {
                "id": hashlib.sha1(href.encode("utf-8")).hexdigest()[:16],
                "title": title[:220],
                "text": (context or title)[:420],
                "permalink": href,
                "source": platform,
                "author": "",
                "authorUsername": "",
                "createdAt": None,
            }
            seen_links.add(href)
            candidates.append((score + (2 if matches else 0), row))

        html_anchor_pattern = re.compile(
            r'<a[^>]+href="(?P<href>[^"]+)"[^>]*>(?P<title>.*?)</a>',
            re.IGNORECASE | re.DOTALL,
        )
        for match in html_anchor_pattern.finditer(text_blob):
            start = max(0, match.start() - 180)
            end = min(len(text_blob), match.end() + 260)
            _push_row(
                match.group("href") or "",
                match.group("title") or "",
                text_blob[start:end],
            )
            if len(candidates) >= max_rows * 3:
                break

        markdown_link_pattern = re.compile(
            r"\[(?P<title>[^\]\n]{2,220})\]\((?P<href>https?://[^)\s]+)\)",
            re.IGNORECASE,
        )
        for match in markdown_link_pattern.finditer(text_blob):
            start = max(0, match.start() - 160)
            end = min(len(text_blob), match.end() + 240)
            _push_row(
                match.group("href") or "",
                match.group("title") or "",
                text_blob[start:end],
            )
            if len(candidates) >= max_rows * 4:
                break

        plain_url_pattern = re.compile(r"https?://[^\s<>'\"\]\)]+", re.IGNORECASE)
        for match in plain_url_pattern.finditer(text_blob):
            start = max(0, match.start() - 140)
            end = min(len(text_blob), match.end() + 220)
            href = match.group(0) or ""
            _push_row(href, href, text_blob[start:end])
            if len(candidates) >= max_rows * 5:
                break

        candidates.sort(key=lambda item: item[0], reverse=True)
        return [row for _, row in candidates[:max_rows]]

    search_query = quote_plus(f"site:{domain} {q}")
    attempts = [
        (
            "duckduckgo",
            f"https://duckduckgo.com/html/?q={search_query}",
            {"User-Agent": "Mozilla/5.0 (Quantura)"},
        ),
        (
            "brave",
            f"https://search.brave.com/search?q={search_query}&source=web",
            {"User-Agent": "Mozilla/5.0 (Quantura)", "Accept-Encoding": "identity"},
        ),
        (
            "jina-duckduckgo",
            f"https://r.jina.ai/http://duckduckgo.com/?q={search_query}",
            {"User-Agent": "Mozilla/5.0 (Quantura)"},
        ),
    ]
    errors: list[str] = []
    for source_name, source_url, source_headers in attempts:
        try:
            response = requests.get(source_url, headers=source_headers, timeout=16)
            if response.status_code >= 400:
                errors.append(f"{source_name} HTTP {response.status_code}")
                continue
            text_blob = response.text or ""
            rows = _extract_social_rows(text_blob)
            if rows:
                return rows, ""
            errors.append(f"{source_name} returned no matching posts")
        except Exception as exc:
            errors.append(f"{source_name} error: {str(exc)[:120]}")

    max_rows = max(1, min(int(limit or 8), 20))
    seed_links: dict[str, list[tuple[str, str]]] = {
        "x": [
            ("Live X search", f"https://x.com/search?q={quote_plus(q)}&src=typed_query&f=live"),
            ("Kobeissi Letter", "https://x.com/KobeissiLetter"),
            ("Seeking Alpha", "https://x.com/SeekingAlpha"),
            ("Benzinga", "https://x.com/Benzinga"),
        ],
        "reddit": [
            ("Reddit live search", f"https://www.reddit.com/search/?q={quote_plus(q)}"),
            ("r/investing", "https://www.reddit.com/r/investing/"),
            ("r/stocks", "https://www.reddit.com/r/stocks/"),
            ("r/wallstreetbets", "https://www.reddit.com/r/wallstreetbets/"),
        ],
        "facebook": [
            ("Facebook search", f"https://www.facebook.com/search/top?q={quote_plus(q)}"),
            ("Investing.com", "https://www.facebook.com/investingcom/"),
            ("CNBC", "https://www.facebook.com/CNBC/"),
            ("Bloomberg", "https://www.facebook.com/Bloomberg/"),
        ],
        "instagram": [
            ("Instagram finance hashtag", "https://www.instagram.com/explore/tags/stocks/"),
            ("@investingcom", "https://www.instagram.com/investingcom/"),
            ("@benzinga", "https://www.instagram.com/benzinga/"),
            ("@wsj", "https://www.instagram.com/wsj/"),
        ],
    }
    if platform in seed_links:
        seeded_rows: list[dict[str, Any]] = []
        for title, href in seed_links.get(platform, [])[:max_rows]:
            seeded_rows.append(
                {
                    "id": hashlib.sha1(f"{platform}:{href}".encode("utf-8")).hexdigest()[:16],
                    "title": title,
                    "text": f"Public {platform.upper()} source for: {q}",
                    "permalink": href,
                    "source": platform,
                    "author": "",
                    "authorUsername": "",
                    "createdAt": None,
                }
            )
        if seeded_rows:
            return seeded_rows, "Live social posts were unavailable; showing public search/profile fallback."

    if errors:
        return [], "; ".join(errors[:2])
    return [], "Web fallback search returned no matching posts."


def _fetch_reddit_social_posts(query: str, *, limit: int = 8) -> list[dict[str, Any]]:
    q = str(query or "").strip()
    if not q:
        return []
    fallback_posts, _ = _fetch_social_posts_via_web_search(q, platform="reddit", limit=limit)
    query_candidates = [q]
    compact = " ".join(_social_query_terms(q)[:4]).strip()
    if compact and compact.lower() != q.lower():
        query_candidates.append(compact)
    if "stock market" not in q.lower():
        query_candidates.append(f"{q} stock market")
    query_candidates.append("stock market investing")
    seen_queries: set[str] = set()
    for query_text in query_candidates:
        clean_query = str(query_text or "").strip()
        if not clean_query or clean_query.lower() in seen_queries:
            continue
        seen_queries.add(clean_query.lower())
        try:
            response = requests.get(
                "https://www.reddit.com/search.json",
                params={
                    "q": clean_query,
                    "sort": "hot",
                    "limit": max(1, min(int(limit or 8), 25)),
                    "restrict_sr": "false",
                },
                headers={"User-Agent": "quantura/1.0"},
                timeout=12,
            )
            if response.status_code >= 400:
                continue
            payload = response.json() if response.text else {}
        except Exception:
            continue

        out: list[dict[str, Any]] = []
        for child in (payload.get("data") or {}).get("children") or []:
            post = child.get("data") if isinstance(child, dict) else None
            if not isinstance(post, dict):
                continue
            permalink = str(post.get("permalink") or "").strip()
            out.append(
                {
                    "id": str(post.get("id") or "").strip(),
                    "title": str(post.get("title") or "").strip(),
                    "text": str(post.get("selftext") or "").strip()[:420],
                    "author": str(post.get("author") or "").strip(),
                    "subreddit": str(post.get("subreddit") or "").strip(),
                    "score": int(post.get("score") or 0),
                    "numComments": int(post.get("num_comments") or 0),
                    "createdAt": int(post.get("created_utc") or 0),
                    "permalink": f"https://www.reddit.com{permalink}" if permalink else "",
                    "source": "reddit",
                }
            )
            if len(out) >= limit:
                break
        if out:
            return out
    return fallback_posts


def _fetch_x_news_stories(
    query: str,
    *,
    limit: int = 6,
    max_age_hours: int = 72,
) -> tuple[list[dict[str, Any]], str]:
    """Fetch AI-generated news stories from X News Search API.

    Docs: https://docs.x.com/x-api/news/search-news
    """
    q = str(query or "").strip()
    if not q:
        return [], "Query is required."
    if not TWITTER_BEARER_TOKEN:
        return [], "X bearer token is not configured for news search."

    endpoint = "https://api.x.com/2/news/search"
    max_results = max(1, min(int(limit or 6), 20))
    params = {
        "query": q,
        "max_results": min(100, max_results),
        "max_age_hours": max(1, min(int(max_age_hours or 72), 720)),
        "news.fields": "category,contexts,hook,name,summary,updated_at,cluster_posts_results",
    }

    try:
        response = requests.get(
            endpoint,
            headers={"Authorization": f"Bearer {TWITTER_BEARER_TOKEN}"},
            params=params,
            timeout=15,
        )
        if response.status_code >= 400:
            return [], f"X News API returned HTTP {response.status_code}."
        payload = response.json() if response.text else {}
    except Exception as exc:
        return [], f"X News API error: {str(exc)[:160]}"

    stories: list[dict[str, Any]] = []
    for row in payload.get("data") or []:
        if not isinstance(row, dict):
            continue
        story_id = str(row.get("id") or row.get("rest_id") or "").strip()
        name = str(row.get("name") or "").strip()
        summary = str(row.get("summary") or "").strip()
        hook = str(row.get("hook") or "").strip()
        category = str(row.get("category") or "").strip()
        updated_at = row.get("updated_at") or row.get("updatedAt")

        tickers: list[str] = []
        contexts = row.get("contexts")
        if isinstance(contexts, dict):
            finance = contexts.get("finance")
            if isinstance(finance, dict):
                raw_tickers = finance.get("tickers") or []
                if isinstance(raw_tickers, list):
                    tickers = [str(item).strip().upper() for item in raw_tickers if str(item).strip()]

        post_ids: list[str] = []
        cluster = row.get("cluster_posts_results")
        if isinstance(cluster, list):
            for item in cluster:
                if not isinstance(item, dict):
                    continue
                post_id = str(item.get("post_id") or item.get("tweet_id") or "").strip()
                if post_id:
                    post_ids.append(post_id)

        search_term = name or q
        search_url = f"https://x.com/search?q={quote_plus(search_term)}&src=typed_query&f=live"

        stories.append(
            {
                "id": story_id,
                "name": name,
                "summary": summary,
                "hook": hook,
                "category": category,
                "updatedAt": updated_at,
                "tickers": tickers,
                "postIds": post_ids[:8],
                "searchUrl": search_url,
                "source": "x_news",
            }
        )
        if len(stories) >= max_results:
            break

    if stories:
        return stories, ""
    return [], "X News API returned no stories."


def _x_news_stories_as_pseudo_posts(stories: list[dict[str, Any]], *, limit: int = 8) -> list[dict[str, Any]]:
    out: list[dict[str, Any]] = []
    for story in stories[: max(1, min(int(limit or 8), 15))]:
        if not isinstance(story, dict):
            continue
        story_id = str(story.get("id") or "").strip()
        name = str(story.get("name") or "").strip() or "X News story"
        text = str(story.get("hook") or story.get("summary") or name).strip()
        out.append(
            {
                "id": f"news_{story_id}" if story_id else hashlib.sha1(name.encode("utf-8")).hexdigest()[:16],
                "text": text[:420],
                "authorName": "X News",
                "authorUsername": "",
                "createdAt": story.get("updatedAt"),
                "metrics": {"like_count": 0, "retweet_count": 0, "reply_count": 0},
                "permalink": str(story.get("searchUrl") or "").strip(),
                "source": "x",
                "kind": "news_story",
                "title": name,
            }
        )
    return out


def _fetch_x_social_posts(query: str, *, limit: int = 8) -> tuple[list[dict[str, Any]], str]:
    q = str(query or "").strip()
    if not q:
        return [], "Query is required."
    fallback_posts, fallback_warning = _fetch_social_posts_via_web_search(q, platform="x", limit=limit)

    endpoint = "https://api.twitter.com/2/tweets/search/recent"
    core_terms = _social_query_terms(q)
    query_text = " ".join(core_terms[:4]) if core_terms else q
    params = {
        "query": f"({query_text}) -is:retweet",
        "max_results": max(10, min(int(limit or 8) * 2, 40)),
        "tweet.fields": "created_at,public_metrics,author_id,lang",
        "expansions": "author_id",
        "user.fields": "username,name,verified",
    }

    auth_attempts: list[tuple[str, dict[str, str], Any]] = []
    if TWITTER_API_KEY and TWITTER_API_SECRET and TWITTER_ACCESS_TOKEN and TWITTER_ACCESS_TOKEN_SECRET:
        try:
            from requests_oauthlib import OAuth1  # type: ignore

            auth_attempts.append(
                (
                    "oauth1_user",
                    {},
                    OAuth1(
                        TWITTER_API_KEY,
                        TWITTER_API_SECRET,
                        TWITTER_ACCESS_TOKEN,
                        TWITTER_ACCESS_TOKEN_SECRET,
                    ),
                )
            )
        except Exception:
            pass
    if TWITTER_BEARER_TOKEN:
        auth_attempts.append(("bearer", {"Authorization": f"Bearer {TWITTER_BEARER_TOKEN}"}, None))
    if not auth_attempts:
        if fallback_posts:
            return fallback_posts, "Using web fallback for X posts because API credentials are missing."
        return [], "X API credentials are not configured."

    payload: dict[str, Any] = {}
    api_failure = ""
    try:
        for auth_mode, headers, auth_obj in auth_attempts:
            response = requests.get(
                endpoint,
                headers=headers,
                auth=auth_obj,
                params=params,
                timeout=15,
            )
            if response.status_code >= 400:
                api_failure = f"{auth_mode} HTTP {response.status_code}"
                continue
            payload = response.json() if response.text else {}
            api_failure = ""
            break
        if not payload:
            if fallback_posts:
                reason = f"API returned {api_failure}." if api_failure else "API was unavailable."
                return fallback_posts, f"Using web fallback for X posts because {reason}"
            if api_failure:
                return [], f"X API returned {api_failure}."
            return [], "X API returned no usable payload."
    except Exception as exc:
        if fallback_posts:
            return fallback_posts, f"Using web fallback for X posts because API errored: {str(exc)[:120]}"
        return [], f"X API error: {str(exc)[:160]}"

    users = {
        str(item.get("id") or "").strip(): item
        for item in (payload.get("includes") or {}).get("users") or []
        if isinstance(item, dict)
    }
    ranked: list[tuple[float, dict[str, Any]]] = []
    for row in payload.get("data") or []:
        if not isinstance(row, dict):
            continue
        post_id = str(row.get("id") or "").strip()
        if not post_id:
            continue
        text_body = str(row.get("text") or "").strip()
        if not _text_matches_social_query(text_body, q):
            continue
        author = users.get(str(row.get("author_id") or "").strip()) or {}
        metrics = row.get("public_metrics") if isinstance(row.get("public_metrics"), dict) else {}
        likes = int(metrics.get("like_count") or 0)
        reposts = int(metrics.get("retweet_count") or 0)
        replies = int(metrics.get("reply_count") or 0)
        score = likes + reposts * 2 + replies * 1.5
        username = str(author.get("username") or "").strip()
        ranked.append(
            (
                score,
                {
                    "id": post_id,
                    "text": text_body,
                    "authorName": str(author.get("name") or "").strip(),
                    "authorUsername": username,
                    "createdAt": row.get("created_at"),
                    "metrics": {
                        "like_count": likes,
                        "retweet_count": reposts,
                        "reply_count": replies,
                    },
                    "permalink": f"https://x.com/{username}/status/{post_id}" if username else f"https://x.com/i/web/status/{post_id}",
                    "source": "x",
                },
            )
        )
    ranked.sort(key=lambda item: item[0], reverse=True)
    top = [item[1] for item in ranked[: max(1, min(int(limit or 8), 15))]]
    if top:
        return top, ""
    if fallback_posts:
        warn = "Using web fallback for X posts because API returned no matching tweets."
        if fallback_warning:
            warn = f"{warn} {fallback_warning}"
        return fallback_posts, warn.strip()
    return [], "No matching X posts found."


def _fetch_meta_social_posts(query: str, *, platform: str, limit: int = 6) -> tuple[list[dict[str, Any]], str]:
    q = str(query or "").strip().lower()
    fallback_posts, fallback_warning = _fetch_social_posts_via_web_search(q, platform=platform, limit=limit)
    if platform == "facebook":
        if not FACEBOOK_PAGE_ID or not FACEBOOK_PAGE_ACCESS_TOKEN:
            if fallback_posts:
                return fallback_posts, "Using web fallback for Facebook because page credentials are missing."
            return [], "Facebook page credentials are not configured."
        try:
            response = requests.get(
                f"https://graph.facebook.com/{META_GRAPH_API_VERSION}/{FACEBOOK_PAGE_ID}/posts",
                params={
                    "access_token": FACEBOOK_PAGE_ACCESS_TOKEN,
                    "fields": "id,message,permalink_url,created_time,reactions.summary(true),comments.summary(true)",
                    "limit": 20,
                },
                timeout=14,
            )
            if response.status_code >= 400:
                if fallback_posts:
                    return fallback_posts, f"Using web fallback for Facebook because API returned HTTP {response.status_code}."
                return [], f"Facebook API returned HTTP {response.status_code}."
            payload = response.json() if response.text else {}
        except Exception as exc:
            if fallback_posts:
                return fallback_posts, f"Using web fallback for Facebook because API errored: {str(exc)[:120]}"
            return [], f"Facebook API error: {str(exc)[:160]}"

        out: list[dict[str, Any]] = []
        backup: list[dict[str, Any]] = []
        for row in payload.get("data") or []:
            if not isinstance(row, dict):
                continue
            message = str(row.get("message") or "").strip()
            item = {
                "id": str(row.get("id") or "").strip(),
                "text": message,
                "createdAt": row.get("created_time"),
                "permalink": str(row.get("permalink_url") or "").strip(),
                "reactions": int(((row.get("reactions") or {}).get("summary") or {}).get("total_count") or 0),
                "comments": int(((row.get("comments") or {}).get("summary") or {}).get("total_count") or 0),
                "source": "facebook",
            }
            backup.append(item)
            if _text_matches_social_query(message, q):
                out.append(item)
            if len(backup) >= max(20, limit * 2):
                break
        chosen = out[:limit] or backup[:limit]
        if chosen:
            return chosen, ""
        if fallback_posts:
            warn = "Using web fallback for Facebook because API returned no posts."
            if fallback_warning:
                warn = f"{warn} {fallback_warning}"
            return fallback_posts, warn.strip()
        return [], "No Facebook posts found."

    if platform == "instagram":
        if not INSTAGRAM_BUSINESS_ACCOUNT_ID or not INSTAGRAM_ACCESS_TOKEN:
            if fallback_posts:
                return fallback_posts, "Using web fallback for Instagram because business credentials are missing."
            return [], "Instagram business credentials are not configured."
        try:
            response = requests.get(
                f"https://graph.facebook.com/{META_GRAPH_API_VERSION}/{INSTAGRAM_BUSINESS_ACCOUNT_ID}/media",
                params={
                    "access_token": INSTAGRAM_ACCESS_TOKEN,
                    "fields": "id,caption,permalink,timestamp,media_type",
                    "limit": 25,
                },
                timeout=14,
            )
            if response.status_code >= 400:
                if fallback_posts:
                    return fallback_posts, f"Using web fallback for Instagram because API returned HTTP {response.status_code}."
                return [], f"Instagram API returned HTTP {response.status_code}."
            payload = response.json() if response.text else {}
        except Exception as exc:
            if fallback_posts:
                return fallback_posts, f"Using web fallback for Instagram because API errored: {str(exc)[:120]}"
            return [], f"Instagram API error: {str(exc)[:160]}"

        out: list[dict[str, Any]] = []
        backup: list[dict[str, Any]] = []
        for row in payload.get("data") or []:
            if not isinstance(row, dict):
                continue
            caption = str(row.get("caption") or "").strip()
            item = {
                "id": str(row.get("id") or "").strip(),
                "text": caption[:600],
                "createdAt": row.get("timestamp"),
                "permalink": str(row.get("permalink") or "").strip(),
                "mediaType": str(row.get("media_type") or "").strip(),
                "source": "instagram",
            }
            backup.append(item)
            if _text_matches_social_query(caption, q):
                out.append(item)
            if len(backup) >= max(25, limit * 2):
                break
        chosen = out[:limit] or backup[:limit]
        if chosen:
            return chosen, ""
        if fallback_posts:
            warn = "Using web fallback for Instagram because API returned no posts."
            if fallback_warning:
                warn = f"{warn} {fallback_warning}"
            return fallback_posts, warn.strip()
        return [], "No Instagram posts found."

    return [], "Unsupported platform."


def _resolve_screener_note_signals(
    notes: str,
    preferred_model: str = "gpt-5-mini",
    context: dict[str, Any] | None = None,
) -> dict[str, Any]:
    notes_text = str(notes or "").strip()
    if not notes_text:
        return {
            "tickers": [],
            "queries": [],
            "matchedHints": [],
            "usedWebSearch": False,
        }

    merged_tickers: list[str] = _extract_symbols_from_text(notes_text)
    queries: list[str] = [notes_text]
    matched_hints: list[str] = []
    lowered = notes_text.lower()
    for hint, symbols in SCREENER_NOTES_TICKER_HINTS.items():
        if hint in lowered:
            matched_hints.append(hint)
            for symbol in symbols:
                clean = _normalize_symbol_token(symbol)
                if clean and clean not in merged_tickers:
                    merged_tickers.append(clean)

    used_web_search = False
    used_provider = "none"
    allowed_models = _get_llm_allowed_models(context=context)
    allowed_set = set(allowed_models)
    model_id = _normalize_ai_model_id(preferred_model or "gpt-5-mini")
    if not _is_supported_llm_model(model_id):
        model_id = "gpt-5-mini"
    if allowed_set and model_id not in allowed_set:
        model_id = allowed_models[0] if allowed_models else "gpt-5-mini"

    system_prompt = (
        "You help convert natural-language investing requests into a stock watchlist. "
        "Use web search when useful and return strict JSON only with this shape: "
        '{"tickers":["AAPL"],"queries":["tech mega caps"],"reasoning":"short text"}. '
        "Use US-listed symbols whenever possible and keep tickers to 20 max."
    )
    user_prompt = (
        "Parse this request and infer likely stock symbols and search phrases for a screener:\\n"
        f"{notes_text}"
    )
    parsed: dict[str, Any] = {}

    provider = _model_provider_from_id(model_id)
    if provider == "openai" and OPENAI_API_KEY:
        responses_payload = {
            "model": model_id,
            "input": [
                {"role": "system", "content": [{"type": "input_text", "text": system_prompt}]},
                {"role": "user", "content": [{"type": "input_text", "text": user_prompt}]},
            ],
            "text": {"format": {"type": "json_object"}},
            "tools": [{"type": "web_search_preview"}],
            "max_output_tokens": 600,
        }
        try:
            response = requests.post(
                "https://api.openai.com/v1/responses",
                headers={
                    "Authorization": f"Bearer {OPENAI_API_KEY}",
                    "Content-Type": "application/json",
                },
                json=responses_payload,
                timeout=28,
            )
            response.raise_for_status()
            used_web_search = True
            used_provider = "openai"
            content = _extract_responses_output_text(response.json())
            parsed = _extract_json_object(content)
        except Exception:
            # Keep note parsing resilient if web search or model call is unavailable.
            parsed = {}
    elif provider == "amazon_nova":
        content = _invoke_amazon_nova_text(
            model_id=model_id,
            system_prompt=(
                f"{system_prompt} Return JSON only. "
                "If web search is unavailable, infer from the user prompt."
            ),
            user_prompt=user_prompt,
            max_tokens=700,
            temperature=0.2,
        )
        if content:
            parsed = _extract_json_object(content)
            used_provider = "amazon_nova"

    for raw_symbol in parsed.get("tickers", []) if isinstance(parsed.get("tickers"), list) else []:
        symbol = _normalize_symbol_token(raw_symbol)
        if symbol and symbol not in merged_tickers:
            merged_tickers.append(symbol)
    for raw_query in parsed.get("queries", []) if isinstance(parsed.get("queries"), list) else []:
        query = str(raw_query or "").strip()
        if query and query not in queries:
            queries.append(query)

    # Resolve natural-language queries through Yahoo symbol search as an extra safety net.
    for query in queries[:8]:
        for symbol in _yahoo_search_symbols(query, max_results=8):
            if symbol not in merged_tickers:
                merged_tickers.append(symbol)

    return {
        "tickers": merged_tickers[:60],
        "queries": queries[:12],
        "matchedHints": matched_hints[:24],
        "usedWebSearch": used_web_search,
        "provider": used_provider,
        "model": model_id,
    }


def _default_social_copy(
    topic: str,
    objective: str,
    audience: str,
    tone: str,
    channels: list[str],
    posts_per_channel: int,
    cta_url: str,
) -> dict[str, list[dict[str, Any]]]:
    base_objective = objective.strip() or "Drive qualified leads to Quantura"
    base_audience = audience.strip() or "active investors and data-driven operators"
    base_tone = tone.strip() or "confident, concise, practical"
    now_utc = datetime.now(timezone.utc)
    out: dict[str, list[dict[str, Any]]] = {}

    for channel in channels:
        channel_posts: list[dict[str, Any]] = []
        for idx in range(posts_per_channel):
            headline = f"{topic} insight #{idx + 1}"
            body = (
                f"{topic}: clear signal, clear execution. "
                f"We built this for {base_audience}. "
                f"Objective: {base_objective}. "
                f"Tone: {base_tone}."
            )
            channel_posts.append(
                {
                    "headline": headline,
                    "body": body,
                    "hashtags": ["#Quantura", "#Stocks", "#DataScience", f"#{channel.capitalize()}"],
                    "cta": "Explore the platform",
                    "ctaUrl": cta_url,
                    "suggestedPostTime": _strategic_suggested_post_time(channel, idx, now_utc=now_utc),
                }
            )
        out[channel] = channel_posts
    return out


def _generate_social_copy_with_openai(
    *,
    topic: str,
    objective: str,
    audience: str,
    tone: str,
    channels: list[str],
    posts_per_channel: int,
    cta_url: str,
) -> tuple[dict[str, list[dict[str, Any]]], str]:
    if not OPENAI_API_KEY:
        drafts = _default_social_copy(topic, objective, audience, tone, channels, posts_per_channel, cta_url)
        return drafts, "template_fallback"

    system_prompt = (
        "You are a social media strategist for a financial analytics company. "
        "Write concise, compliant, high-conversion posts. "
        "Avoid guaranteed-return language. Keep content useful and brand-safe. "
        "Return strict JSON only."
    )
    # Pull a tiny amount of live market context to make scheduled automation feel timely.
    # All failures should degrade silently to avoid blocking posting.
    market_context_lines: list[str] = []
    try:
        for item in _fetch_yahoo_news_query("US stock market top headlines today", limit=6)[:4]:
            if not isinstance(item, dict):
                continue
            title = str(item.get("title") or "").strip()
            publisher = str(item.get("publisher") or "").strip()
            if title:
                market_context_lines.append(f"- {title}{f' ({publisher})' if publisher else ''}")
    except Exception:
        pass
    try:
        x_stories, _ = _fetch_x_news_stories("US stock market top headlines today", limit=4, max_age_hours=48)
        for story in x_stories[:3]:
            if not isinstance(story, dict):
                continue
            name = str(story.get("name") or "").strip()
            hook = str(story.get("hook") or story.get("summary") or "").strip()
            if name:
                market_context_lines.append(f"- X News: {name}{f' — {hook}' if hook else ''}")
    except Exception:
        pass
    market_context = "\n".join(market_context_lines).strip()
    def _request_json_from_openai(user_prompt: str, max_output_tokens: int = 2200) -> dict[str, Any]:
        responses_payload = {
            "model": SOCIAL_CONTENT_MODEL,
            "input": [
                {"role": "system", "content": [{"type": "input_text", "text": system_prompt}]},
                {"role": "user", "content": [{"type": "input_text", "text": user_prompt}]},
            ],
            "text": {"format": {"type": "json_object"}},
            "max_output_tokens": max_output_tokens,
        }
        response = requests.post(
            "https://api.openai.com/v1/responses",
            headers={
                "Authorization": f"Bearer {OPENAI_API_KEY}",
                "Content-Type": "application/json",
            },
            json=responses_payload,
            timeout=35,
        )
        response.raise_for_status()
        content = _extract_responses_output_text(response.json())
        if content:
            parsed = _extract_json_object(content)
            if parsed:
                return parsed

        # Fallback path for compatibility across model/endpoint variants.
        chat_payload = {
            "model": SOCIAL_CONTENT_MODEL,
            "temperature": 0.7,
            "max_completion_tokens": 2200,
            "messages": [
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": user_prompt},
            ],
            "response_format": {"type": "json_object"},
        }
        chat_response = requests.post(
            "https://api.openai.com/v1/chat/completions",
            headers={
                "Authorization": f"Bearer {OPENAI_API_KEY}",
                "Content-Type": "application/json",
            },
            json=chat_payload,
            timeout=35,
        )
        chat_response.raise_for_status()
        chat_body = chat_response.json()
        chat_content = (
            chat_body.get("choices", [{}])[0]
            .get("message", {})
            .get("content", "{}")
        )
        parsed = _extract_json_object(chat_content)
        if not parsed:
            raise ValueError("OpenAI response did not return parseable JSON.")
        return parsed

    normalized: dict[str, list[dict[str, Any]]] = {}
    used_model = False

    for channel in channels:
        context_block = (
            f"\n\nCurrent market context (for inspiration only, do not quote verbatim):\n{market_context}\n"
            if market_context
            else ""
        )
        per_channel_prompt = (
            "Generate social drafts with this exact JSON shape: "
            '{"channels":{"x":[{"headline":"","body":"","hashtags":[],"cta":"","ctaUrl":"","suggestedPostTime":""}]}}. '
            f"Topic: {topic}. Objective: {objective}. Audience: {audience}. Tone: {tone}. "
            f"Generate only channel: {channel}. Posts per channel: {posts_per_channel}. "
            f"Use CTA URL: {cta_url}. Keep each body optimized for the channel and under typical limits."
            f"{context_block}"
        )
        rows: list[Any] = []
        try:
            parsed = _request_json_from_openai(per_channel_prompt)
            if isinstance(parsed, dict):
                channels_obj = parsed.get("channels")
                if isinstance(channels_obj, dict):
                    candidate = channels_obj.get(channel)
                    if isinstance(candidate, list):
                        rows = candidate
                elif isinstance(parsed.get(channel), list):
                    rows = parsed.get(channel)  # type: ignore[assignment]
            mapped: list[dict[str, Any]] = []
            for row in rows[:posts_per_channel]:
                if not isinstance(row, dict):
                    continue
                hashtags_raw = row.get("hashtags")
                if isinstance(hashtags_raw, list):
                    hashtags = [str(tag).strip() for tag in hashtags_raw if str(tag).strip()]
                elif isinstance(hashtags_raw, str):
                    hashtags = [tag.strip() for tag in hashtags_raw.split(",") if tag.strip()]
                else:
                    hashtags = []
                mapped.append(
                    {
                        "headline": str(row.get("headline") or "").strip(),
                        "body": str(row.get("body") or "").strip(),
                        "hashtags": hashtags,
                        "cta": str(row.get("cta") or "Learn more").strip(),
                        "ctaUrl": str(row.get("ctaUrl") or cta_url).strip() or cta_url,
                        "suggestedPostTime": str(row.get("suggestedPostTime") or "").strip(),
                    }
                )
            if mapped:
                normalized[channel] = mapped
                used_model = True
                continue
        except Exception:
            pass

        normalized[channel] = _default_social_copy(
            topic,
            objective,
            audience,
            tone,
            [channel],
            posts_per_channel,
            cta_url,
        ).get(channel, [])

    return normalized, (SOCIAL_CONTENT_MODEL if used_model else "template_fallback")


def _enqueue_social_posts(
    *,
    campaign_id: str,
    user_id: str,
    user_email: str,
    channels: list[str],
    drafts: dict[str, list[dict[str, Any]]],
    scheduled_for: datetime,
    meta: dict[str, Any],
) -> dict[str, Any]:
    created_ids: list[str] = []
    total = 0
    now_utc = datetime.now(timezone.utc)
    scheduled_utc = scheduled_for.astimezone(timezone.utc)
    for channel in channels:
        posts = drafts.get(channel) if isinstance(drafts, dict) else None
        if not isinstance(posts, list):
            continue
        for rank, post in enumerate(posts):
            if not isinstance(post, dict):
                continue
            body = str(post.get("body") or "").strip()
            headline = str(post.get("headline") or "").strip()
            if not body and not headline:
                continue
            suggested = _as_utc_datetime(post.get("suggestedPostTime"))
            publish_time = suggested or scheduled_utc
            if publish_time < now_utc - timedelta(minutes=2):
                while publish_time < now_utc + timedelta(minutes=2):
                    publish_time += timedelta(days=1)
            doc = {
                "campaignId": campaign_id,
                "userId": user_id,
                "userEmail": user_email,
                "platform": channel,
                "headline": headline,
                "body": body,
                "hashtags": post.get("hashtags") if isinstance(post.get("hashtags"), list) else [],
                "cta": str(post.get("cta") or "").strip(),
                "ctaUrl": str(post.get("ctaUrl") or SOCIAL_DEFAULT_CTA_URL).strip(),
                "scheduledFor": publish_time,
                "status": "queued",
                "attempts": 0,
                "orderIndex": rank,
                "createdAt": firestore.SERVER_TIMESTAMP,
                "updatedAt": firestore.SERVER_TIMESTAMP,
                "meta": meta or {},
            }
            ref = db.collection("social_queue").document()
            ref.set(doc)
            created_ids.append(ref.id)
            total += 1
    return {"count": total, "queueIds": created_ids}


def _post_to_social_channel(
    *,
    platform: str,
    body: str,
    headline: str,
    hashtags: list[str],
    cta: str,
    cta_url: str,
    campaign_id: str,
    queue_id: str,
) -> dict[str, Any]:
    rendered_text, tags = _render_social_text(
        body=body,
        headline=headline,
        hashtags=hashtags,
        cta=cta,
        cta_url=cta_url,
    )

    direct_dispatchers = {
        "x": _post_x_direct,
        "linkedin": _post_linkedin_direct,
        "facebook": _post_facebook_direct,
        "instagram": _post_instagram_direct,
        "tiktok": _post_tiktok_direct,
    }
    direct_result: dict[str, Any] | None = None
    dispatcher = direct_dispatchers.get(platform)
    if dispatcher is not None:
        direct_result = dispatcher(rendered_text)
        if direct_result.get("ok"):
            return direct_result
        # If direct credentials are missing or unsupported, fallback to webhook relay.
        if not direct_result.get("pendingCredentials"):
            webhook_candidate = str(_social_channel_webhooks().get(platform) or "").strip()
            if not webhook_candidate:
                return direct_result

    webhooks = _social_channel_webhooks()
    webhook = str(webhooks.get(platform) or "").strip()
    if not webhook:
        if direct_result:
            return direct_result
        if dispatcher is not None:
            return {
                "ok": False,
                "pendingCredentials": True,
                "error": f"Missing direct credentials and webhook for channel '{platform}'.",
            }
        return {"ok": False, "pendingCredentials": True, "error": f"Missing webhook for channel '{platform}'."}

    payload = {
        "platform": platform,
        "text": rendered_text,
        "headline": headline.strip(),
        "body": body.strip(),
        "hashtags": tags,
        "cta": cta.strip(),
        "ctaUrl": cta_url.strip(),
        "campaignId": campaign_id,
        "queueId": queue_id,
    }
    response = requests.post(webhook, json=payload, timeout=20)
    if response.status_code >= 400:
        return {"ok": False, "pendingCredentials": False, "error": f"{response.status_code} {response.text}"}
    external_id = ""
    try:
        response_body = response.json()
        external_id = str(response_body.get("id") or response_body.get("postId") or "").strip()
    except Exception:
        external_id = ""
    return {"ok": True, "externalId": external_id, "statusCode": response.status_code, "provider": "webhook"}


def _dispatch_due_social_posts(*, max_posts: int, trigger: str, actor_uid: str = "") -> dict[str, Any]:
    now_utc = datetime.now(timezone.utc)
    snapshot = (
        db.collection("social_queue")
        .where("status", "in", ["queued", "retry"])
        .limit(max(1, min(max_posts * 4, 200)))
        .stream()
    )

    due_docs: list[tuple[str, dict[str, Any]]] = []
    for doc in snapshot:
        data = doc.to_dict() or {}
        scheduled_for = _as_utc_datetime(data.get("scheduledFor"))
        if scheduled_for is None or scheduled_for <= now_utc:
            due_docs.append((doc.id, data))
        if len(due_docs) >= max_posts:
            break

    posted = 0
    waiting_credentials = 0
    failed = 0
    retried = 0
    processed_ids: list[str] = []

    for queue_id, item in due_docs:
        platform = str(item.get("platform") or "").strip().lower()
        attempts = int(item.get("attempts") or 0)
        result = _post_to_social_channel(
            platform=platform,
            body=str(item.get("body") or ""),
            headline=str(item.get("headline") or ""),
            hashtags=item.get("hashtags") if isinstance(item.get("hashtags"), list) else [],
            cta=str(item.get("cta") or ""),
            cta_url=str(item.get("ctaUrl") or SOCIAL_DEFAULT_CTA_URL),
            campaign_id=str(item.get("campaignId") or ""),
            queue_id=queue_id,
        )

        ref = db.collection("social_queue").document(queue_id)
        if result.get("ok"):
            ref.set(
                {
                    "status": "posted",
                    "attempts": attempts + 1,
                    "postedAt": firestore.SERVER_TIMESTAMP,
                    "updatedAt": firestore.SERVER_TIMESTAMP,
                    "externalId": result.get("externalId") or "",
                    "lastError": "",
                },
                merge=True,
            )
            posted += 1
        elif result.get("pendingCredentials"):
            ref.set(
                {
                    "status": "waiting_credentials",
                    "attempts": attempts + 1,
                    "updatedAt": firestore.SERVER_TIMESTAMP,
                    "lastError": str(result.get("error") or "Missing credentials"),
                },
                merge=True,
            )
            waiting_credentials += 1
        else:
            next_attempts = attempts + 1
            exhausted = next_attempts >= 3
            ref.set(
                {
                    "status": "failed" if exhausted else "retry",
                    "attempts": next_attempts,
                    "updatedAt": firestore.SERVER_TIMESTAMP,
                    "lastError": str(result.get("error") or "Publish failed"),
                },
                merge=True,
            )
            if exhausted:
                failed += 1
            else:
                retried += 1
        processed_ids.append(queue_id)

    db.collection("social_dispatch_logs").add(
        {
            "trigger": trigger,
            "actorUid": actor_uid,
            "scheduledAt": firestore.SERVER_TIMESTAMP,
            "processedCount": len(processed_ids),
            "postedCount": posted,
            "waitingCredentialsCount": waiting_credentials,
            "retryCount": retried,
            "failedCount": failed,
            "queueIds": processed_ids,
        }
    )

    return {
        "processed": len(processed_ids),
        "posted": posted,
        "waitingCredentials": waiting_credentials,
        "retry": retried,
        "failed": failed,
        "trigger": trigger,
    }


def _require_auth(req: https_fn.CallableRequest) -> dict[str, Any]:
    if req.auth is None:
        raise https_fn.HttpsError(
            https_fn.FunctionsErrorCode.UNAUTHENTICATED,
            "Sign in is required to perform this action.",
        )
    return req.auth.token or {}


def _require_admin(token: dict[str, Any]) -> None:
    if token.get("email") != ADMIN_EMAIL:
        raise https_fn.HttpsError(
            https_fn.FunctionsErrorCode.PERMISSION_DENIED,
            "Admin access required.",
        )


def _collaborator_snapshot(workspace_id: str, uid: str):
    return (
        db.collection("users")
        .document(workspace_id)
        .collection("collaborators")
        .document(uid)
        .get()
    )


def _can_read_workspace(workspace_id: str, uid: str, token: dict[str, Any]) -> bool:
    if token.get("email") == ADMIN_EMAIL:
        return True
    if workspace_id == uid:
        return True
    shared = _collaborator_snapshot(workspace_id, uid)
    return bool(shared.exists)


def _can_edit_workspace(workspace_id: str, uid: str, token: dict[str, Any]) -> bool:
    if token.get("email") == ADMIN_EMAIL:
        return True
    if workspace_id == uid:
        return True
    shared = _collaborator_snapshot(workspace_id, uid)
    if not shared.exists:
        return False
    role = str((shared.to_dict() or {}).get("role") or "viewer").strip().lower()
    return role == "editor"


def _require_workspace_access(workspace_id: str, uid: str, token: dict[str, Any]) -> None:
    if not _can_read_workspace(workspace_id, uid, token):
        raise https_fn.HttpsError(
            https_fn.FunctionsErrorCode.PERMISSION_DENIED,
            "Workspace access required.",
        )


def _require_workspace_editor(workspace_id: str, uid: str, token: dict[str, Any]) -> None:
    if not _can_edit_workspace(workspace_id, uid, token):
        raise https_fn.HttpsError(
            https_fn.FunctionsErrorCode.PERMISSION_DENIED,
            "Editor access required for this workspace.",
        )


def _alpaca_headers() -> dict[str, str]:
    if not ALPACA_API_KEY or not ALPACA_SECRET_KEY:
        raise https_fn.HttpsError(
            https_fn.FunctionsErrorCode.FAILED_PRECONDITION,
            "Missing Alpaca API credentials.",
        )
    return {
        "APCA-API-KEY-ID": ALPACA_API_KEY,
        "APCA-API-SECRET-KEY": ALPACA_SECRET_KEY,
        "Content-Type": "application/json",
    }


def _stripe_module():
    try:
        import stripe  # type: ignore
    except Exception as exc:  # pragma: no cover
        raise https_fn.HttpsError(
            https_fn.FunctionsErrorCode.FAILED_PRECONDITION,
            "Stripe SDK is not available in this deployment.",
        ) from exc

    if not STRIPE_SECRET_KEY:
        raise https_fn.HttpsError(
            https_fn.FunctionsErrorCode.FAILED_PRECONDITION,
            "Stripe is not configured (missing STRIPE_SECRET_KEY).",
        )
    stripe.api_key = STRIPE_SECRET_KEY
    return stripe


def _stripe_object_to_dict(obj: Any) -> dict[str, Any]:
    if obj is None:
        return {}
    if isinstance(obj, dict):
        return dict(obj)
    if hasattr(obj, "to_dict"):
        try:
            parsed = obj.to_dict()  # type: ignore[attr-defined]
            if isinstance(parsed, dict):
                return parsed
        except Exception:
            return {}
    try:
        return dict(obj)
    except Exception:
        return {}


def _stripe_id(value: Any) -> str:
    if isinstance(value, dict):
        return str(value.get("id") or "").strip()
    return str(value or "").strip()


def _persist_user_stripe_customer_id(user_id: str, customer_id: str) -> None:
    clean_customer_id = str(customer_id or "").strip()
    if not user_id or not clean_customer_id:
        return
    db.collection("users").document(user_id).set(
        {
            "stripeCustomerId": clean_customer_id,
            "profile": {
                "stripeCustomerId": clean_customer_id,
            },
            "updatedAt": firestore.SERVER_TIMESTAMP,
        },
        merge=True,
    )


def _resolve_stripe_customer_id(
    stripe: Any,
    user_id: str,
    email: str | None = None,
    *,
    create_if_missing: bool = False,
) -> str:
    user_ref = db.collection("users").document(user_id)
    user_snap = user_ref.get()
    user_doc = user_snap.to_dict() if user_snap.exists else {}
    profile = user_doc.get("profile") if isinstance(user_doc.get("profile"), dict) else {}

    existing_customer = str(
        user_doc.get("stripeCustomerId")
        or profile.get("stripeCustomerId")
        or ""
    ).strip()
    if existing_customer:
        return existing_customer

    normalized_email = _normalize_email(email)
    if normalized_email:
        try:
            existing = stripe.Customer.list(email=normalized_email, limit=1)
            existing_data = getattr(existing, "data", None)
            if not isinstance(existing_data, list):
                existing_data = _stripe_object_to_dict(existing).get("data") or []
            if isinstance(existing_data, list) and existing_data:
                customer_id = _stripe_id(existing_data[0])
                if customer_id:
                    _persist_user_stripe_customer_id(user_id, customer_id)
                    return customer_id
        except Exception:
            pass

    if not create_if_missing:
        return ""

    create_kwargs: dict[str, Any] = {"metadata": {"userId": user_id}}
    if normalized_email:
        create_kwargs["email"] = normalized_email
    customer = stripe.Customer.create(**create_kwargs)
    customer_id = _stripe_id(customer)
    if customer_id:
        _persist_user_stripe_customer_id(user_id, customer_id)
    return customer_id


def _extract_checkout_customer_id(session_obj: Any) -> str:
    session_dict = _stripe_object_to_dict(session_obj)
    return _stripe_id(session_dict.get("customer"))


def _sanitize_portal_return_url(value: Any) -> str:
    fallback = f"{PUBLIC_ORIGIN}/dashboard"
    raw = str(value or "").strip()
    if not raw:
        return fallback

    if raw.startswith("/"):
        return f"{PUBLIC_ORIGIN}{raw}"

    parsed = urlparse(raw)
    if parsed.scheme not in {"http", "https"}:
        return fallback

    public_host = urlparse(PUBLIC_ORIGIN).netloc
    if parsed.netloc != public_host:
        return fallback
    return raw


def _safe_float(value: Any) -> float | None:
    try:
        return float(value)
    except (TypeError, ValueError):
        return None


def _normalize_email(value: Any) -> str:
    return str(value or "").strip().lower()


def _meta_hash_sha256(value: Any) -> str:
    raw = str(value or "").strip().lower()
    if not raw:
        return ""
    return hashlib.sha256(raw.encode("utf-8")).hexdigest()


def _meta_normalize_phone(value: Any) -> str:
    digits = re.sub(r"\D+", "", str(value or ""))
    return digits if len(digits) >= 7 else ""


def _meta_header(req: https_fn.CallableRequest, key: str) -> str:
    raw_request = getattr(req, "raw_request", None)
    headers = getattr(raw_request, "headers", None)
    if headers is None:
        return ""
    try:
        return str(headers.get(key, "") or "").strip()
    except Exception:
        return ""


def _meta_client_ip(req: https_fn.CallableRequest, data: dict[str, Any]) -> str:
    explicit = str(data.get("clientIpAddress") or data.get("client_ip_address") or "").strip()
    if explicit:
        return explicit.split(",")[0].strip()
    forwarded = _meta_header(req, "x-forwarded-for")
    if forwarded:
        return forwarded.split(",")[0].strip()
    raw_request = getattr(req, "raw_request", None)
    remote_addr = getattr(raw_request, "remote_addr", "")
    return str(remote_addr or "").strip()


def _meta_build_user_data(req: https_fn.CallableRequest, data: dict[str, Any]) -> dict[str, Any]:
    user_data: dict[str, Any] = {}

    user_agent = str(data.get("userAgent") or data.get("user_agent") or _meta_header(req, "user-agent") or "").strip()
    client_ip = _meta_client_ip(req, data)
    if user_agent:
        user_data["client_user_agent"] = user_agent[:512]
    if client_ip:
        user_data["client_ip_address"] = client_ip

    email = str(data.get("email") or "").strip()
    if not email and getattr(req, "auth", None) and isinstance(req.auth.token, dict):
        email = str(req.auth.token.get("email") or "").strip()
    email_hash = _meta_hash_sha256(email)
    if email_hash:
        user_data["em"] = [email_hash]

    phone = _meta_normalize_phone(data.get("phone"))
    phone_hash = _meta_hash_sha256(phone)
    if phone_hash:
        user_data["ph"] = [phone_hash]

    external_id = str(data.get("externalId") or "").strip()
    if not external_id and getattr(req, "auth", None):
        external_id = str(getattr(req.auth, "uid", "") or "").strip()
    external_id_hash = _meta_hash_sha256(external_id)
    if external_id_hash:
        user_data["external_id"] = [external_id_hash]

    fbc = str(data.get("fbc") or "").strip()
    fbp = str(data.get("fbp") or "").strip()
    if fbc:
        user_data["fbc"] = fbc[:256]
    if fbp:
        user_data["fbp"] = fbp[:256]

    return user_data


def _meta_build_custom_data(params: Any) -> dict[str, Any]:
    if not isinstance(params, dict):
        return {}
    out: dict[str, Any] = {}
    allowed_keys = [
        "currency",
        "value",
        "content_name",
        "content_category",
        "search_string",
        "ticker",
        "workspace_id",
        "order_id",
        "status",
    ]
    for key in allowed_keys:
        if key not in params:
            continue
        value = params.get(key)
        if value is None:
            continue
        if key == "value":
            numeric = _safe_float(value)
            if numeric is None:
                continue
            out[key] = round(float(numeric), 6)
            continue
        if isinstance(value, bool):
            out[key] = value
            continue
        if isinstance(value, (int, float)):
            out[key] = value
            continue
        if isinstance(value, str):
            text = value.strip()
            if text:
                out[key] = text[:256]
    return out


def _audit_event(uid: str, email: str | None, event_type: str, payload: dict[str, Any]) -> None:
    db.collection("user_events").add(
        {
            "userId": uid,
            "userEmail": email,
            "eventType": event_type,
            "payload": payload,
            "createdAt": firestore.SERVER_TIMESTAMP,
        }
    )


def _token_doc_id(token: str) -> str:
    return hashlib.sha256(token.encode("utf-8")).hexdigest()


def _normalize_notification_token(raw: Any) -> str:
    token = str(raw or "").strip()
    if len(token) < 20:
        return ""
    return token


def _active_notification_tokens_for_user(user_id: str) -> list[str]:
    docs = (
        db.collection("notification_tokens")
        .where("userId", "==", user_id)
        .where("active", "==", True)
        .limit(500)
        .stream()
    )
    tokens: list[str] = []
    for doc in docs:
        token_value = (doc.to_dict() or {}).get("token")
        if isinstance(token_value, str) and token_value:
            tokens.append(token_value)
    return list(dict.fromkeys(tokens))


def _send_push_tokens(tokens: list[str], title: str, body: str, data: dict[str, Any] | None = None) -> dict[str, Any]:
    tokens = [str(token).strip() for token in tokens if str(token).strip()]
    if not tokens:
        return {"successCount": 0, "failureCount": 0, "failed": [], "staleTokenHashes": []}

    payload_data = {str(key): str(value) for key, value in (data or {}).items() if value is not None}
    target = payload_data.get("url") or "/dashboard"
    link = target if target.startswith("http") else f"{PUBLIC_ORIGIN}{target}"
    icon = f"{PUBLIC_ORIGIN}/assets/quantura-icon.svg"
    message = admin_messaging.MulticastMessage(
        tokens=tokens[:500],
        notification=admin_messaging.Notification(title=title, body=body),
        data=payload_data,
        webpush=admin_messaging.WebpushConfig(
            notification=admin_messaging.WebpushNotification(
                title=title,
                body=body,
                icon=icon,
                badge=icon,
            ),
            fcm_options=admin_messaging.WebpushFCMOptions(link=link),
            headers={"Urgency": "high"},
        ),
    )

    try:
        response = admin_messaging.send_each_for_multicast(message)
    except Exception as exc:
        # Avoid surfacing an opaque "internal error" to the caller. We still audit via _notify_user().
        return {
            "successCount": 0,
            "failureCount": len(tokens[:500]),
            "failed": [{"tokenHash": _token_doc_id(t), "error": str(exc)} for t in tokens[:5]],
            "staleTokenHashes": [],
            "error": str(exc),
        }
    failed: list[dict[str, str]] = []
    stale_token_hashes: list[str] = []

    for idx, resp in enumerate(response.responses):
        if resp.success:
            continue
        token = tokens[idx]
        error_text = str(resp.exception or "unknown_error")
        failed.append({"tokenHash": _token_doc_id(token), "error": error_text})
        normalized = error_text.lower()
        if (
            "unregistered" in normalized
            or "registration-token-not-registered" in normalized
            or "invalid registration token" in normalized
            or "requested entity was not found" in normalized
        ):
            stale_token_hashes.append(_token_doc_id(token))

    return {
        "successCount": response.success_count,
        "failureCount": response.failure_count,
        "failed": failed,
        "staleTokenHashes": stale_token_hashes,
    }


def _notify_user(
    user_id: str,
    user_email: str | None,
    title: str,
    body: str,
    data: dict[str, Any] | None = None,
    explicit_token: str | None = None,
) -> dict[str, Any]:
    tokens = _active_notification_tokens_for_user(user_id)
    fallback_token = _normalize_notification_token(explicit_token)
    used_fallback_token = False
    if fallback_token:
        if fallback_token not in tokens:
            used_fallback_token = True
            tokens.append(fallback_token)
        token_hash = _token_doc_id(fallback_token)
        db.collection("notification_tokens").document(token_hash).set(
            {
                "token": fallback_token,
                "tokenHash": token_hash,
                "active": True,
                "userId": user_id,
                "userEmail": user_email or "",
                "updatedAt": firestore.SERVER_TIMESTAMP,
                "lastSeenAt": firestore.SERVER_TIMESTAMP,
            },
            merge=True,
        )
    tokens = list(dict.fromkeys(tokens))
    result = _send_push_tokens(tokens, title=title, body=body, data=data)
    for stale_hash in result.get("staleTokenHashes", []):
        db.collection("notification_tokens").document(str(stale_hash)).set(
            {
                "active": False,
                "updatedAt": firestore.SERVER_TIMESTAMP,
            },
            merge=True,
        )
    _audit_event(
        user_id,
        user_email,
        "push_notification_attempted",
        {
            "title": title,
            "successCount": result.get("successCount", 0),
            "failureCount": result.get("failureCount", 0),
            "attemptedTokenCount": len(tokens),
            "usedFallbackToken": used_fallback_token,
        },
    )
    result["attemptedTokenCount"] = len(tokens)
    result["usedFallbackToken"] = used_fallback_token
    return result


def _serialize_for_firestore(value: Any) -> Any:
    try:
        import numpy as np  # type: ignore

        if isinstance(value, (np.floating, np.integer)):
            return value.item()
        if isinstance(value, np.ndarray):
            return [_serialize_for_firestore(item) for item in value.tolist()]
    except Exception:
        pass

    try:
        import pandas as pd  # type: ignore

        if isinstance(value, pd.Timestamp):
            return value.isoformat()
    except Exception:
        pass

    if isinstance(value, (datetime, date)):
        return value.isoformat()

    if isinstance(value, list):
        return [_serialize_for_firestore(item) for item in value]
    if isinstance(value, dict):
        return {key: _serialize_for_firestore(item) for key, item in value.items()}
    return value


def _trending_snapshots(tickers: list[str], max_rows: int = 18) -> dict[str, dict[str, Any]]:
    import pandas as pd  # type: ignore
    import yfinance as yf  # type: ignore

    tickers = [str(t).upper().strip() for t in (tickers or []) if str(t).strip()]
    tickers = list(dict.fromkeys([t for t in tickers if t]))[: max(1, int(max_rows or 18))]
    if not tickers:
        return {}

    symbols = " ".join(tickers)
    try:
        frame = yf.download(
            symbols,
            period="14d",
            interval="1d",
            group_by="ticker",
            progress=False,
            threads=True,
            auto_adjust=False,
        )
    except Exception:
        frame = pd.DataFrame()

    if frame.empty:
        return {}

    snapshots: dict[str, dict[str, Any]] = {}

    def _snap_from_close(series: pd.Series) -> dict[str, Any] | None:
        close = pd.to_numeric(series, errors="coerce").dropna()
        if close.empty:
            return None
        last = float(close.iloc[-1])
        prev = float(close.iloc[-2]) if len(close) >= 2 else None
        change = round(last - prev, 4) if prev else None
        change_pct = round((last - prev) / prev * 100.0, 4) if prev else None
        return {"lastClose": round(last, 4), "prevClose": round(prev, 4) if prev else None, "change": change, "changePct": change_pct}

    if isinstance(frame.columns, pd.MultiIndex):
        # shape: (field, ticker) or (ticker, field)
        level0 = list(frame.columns.get_level_values(0))
        is_ticker_level0 = any(t in set(level0) for t in tickers)
        for ticker in tickers:
            try:
                sub = frame[ticker] if is_ticker_level0 else frame.xs(ticker, axis=1, level=1)
            except Exception:
                continue
            if isinstance(sub.columns, pd.MultiIndex):
                sub.columns = sub.columns.get_level_values(-1)
            if "Close" not in sub.columns:
                continue
            snap = _snap_from_close(sub["Close"])
            if snap:
                snapshots[ticker] = snap
        return snapshots

    if "Close" in frame.columns and len(tickers) == 1:
        snap = _snap_from_close(frame["Close"])
        if snap:
            snapshots[tickers[0]] = snap
    return snapshots


def _parse_quantiles(raw: Any) -> list[float]:
    if raw is None:
        return [0.1, 0.5, 0.9]

    values: list[float] = []
    if isinstance(raw, str):
        parts = [p.strip() for p in raw.split(",") if p.strip()]
        if not parts:
            raise https_fn.HttpsError(https_fn.FunctionsErrorCode.INVALID_ARGUMENT, "Quantiles must be comma-delimited values between 0 and 1.")
        for part in parts:
            try:
                values.append(float(part))
            except Exception:
                raise https_fn.HttpsError(https_fn.FunctionsErrorCode.INVALID_ARGUMENT, f"Invalid quantile value: {part}")
    elif isinstance(raw, (list, tuple)):
        if not raw:
            raise https_fn.HttpsError(https_fn.FunctionsErrorCode.INVALID_ARGUMENT, "At least one quantile is required.")
        for item in raw:
            try:
                values.append(float(item))
            except Exception:
                raise https_fn.HttpsError(https_fn.FunctionsErrorCode.INVALID_ARGUMENT, f"Invalid quantile value: {item}")
    else:
        try:
            values = [float(raw)]
        except Exception:
            raise https_fn.HttpsError(https_fn.FunctionsErrorCode.INVALID_ARGUMENT, "Invalid quantile value.")

    cleaned: list[float] = []
    seen: set[int] = set()
    for q in values:
        if not (0.0 < float(q) < 1.0):
            raise https_fn.HttpsError(https_fn.FunctionsErrorCode.INVALID_ARGUMENT, "Quantiles must be between 0 and 1 (exclusive).")
        key = int(round(float(q) * 10000))
        if key in seen:
            continue
        seen.add(key)
        cleaned.append(float(q))

    if len(cleaned) > 12:
        raise https_fn.HttpsError(https_fn.FunctionsErrorCode.INVALID_ARGUMENT, "Too many quantiles requested (max 12).")

    return cleaned


def _load_history(ticker: str, start: str | None, interval: str) -> pd.DataFrame:
    import pandas as pd  # type: ignore
    import yfinance as yf  # type: ignore

    period = "730d" if interval == "1h" else "10y"
    frame = yf.download(
        ticker,
        start=start or None,
        period=None if start else period,
        interval=interval,
        progress=False,
        auto_adjust=False,
    )
    if frame.empty:
        raise https_fn.HttpsError(https_fn.FunctionsErrorCode.NOT_FOUND, "No market data found for ticker.")

    if isinstance(frame.columns, pd.MultiIndex):
        frame.columns = frame.columns.get_level_values(0)

    if "Close" not in frame.columns:
        raise https_fn.HttpsError(https_fn.FunctionsErrorCode.NOT_FOUND, "Close prices unavailable for ticker.")

    frame = frame.dropna(subset=["Close"])
    if frame.empty:
        raise https_fn.HttpsError(https_fn.FunctionsErrorCode.NOT_FOUND, "No valid close values for ticker.")

    return frame


def _latest_close_prices(tickers: list[str]) -> dict[str, float]:
    import pandas as pd  # type: ignore
    import yfinance as yf  # type: ignore

    tickers = [str(t).upper().strip() for t in (tickers or []) if str(t).strip()]
    tickers = list(dict.fromkeys([t for t in tickers if t]))
    if not tickers:
        return {}

    symbols = " ".join(tickers)
    try:
        frame = yf.download(
            symbols,
            period="7d",
            interval="1d",
            group_by="ticker",
            progress=False,
            threads=True,
            auto_adjust=False,
        )
    except Exception:
        frame = pd.DataFrame()

    prices: dict[str, float] = {}
    if frame.empty:
        return prices

    if isinstance(frame.columns, pd.MultiIndex):
        # shape: (field, ticker) or (ticker, field) depending on yfinance version
        level0 = list(frame.columns.get_level_values(0))
        level1 = list(frame.columns.get_level_values(1))
        is_ticker_level0 = any(t in set(level0) for t in tickers)
        for ticker in tickers:
            try:
                sub = frame[ticker] if is_ticker_level0 else frame.xs(ticker, axis=1, level=1)
            except Exception:
                continue
            if isinstance(sub.columns, pd.MultiIndex):
                sub.columns = sub.columns.get_level_values(-1)
            if "Close" not in sub.columns:
                continue
            close = pd.to_numeric(sub["Close"], errors="coerce").dropna()
            if close.empty:
                continue
            prices[ticker] = float(close.iloc[-1])
        return prices

    if "Close" in frame.columns and len(tickers) == 1:
        close = pd.to_numeric(frame["Close"], errors="coerce").dropna()
        if not close.empty:
            prices[tickers[0]] = float(close.iloc[-1])
    return prices


def _generate_quantile_forecast(
    close_series: pd.Series,
    horizon: int,
    quantiles: list[float],
    interval: str,
) -> dict[str, Any]:
    import numpy as np  # type: ignore
    import pandas as pd  # type: ignore

    quantiles = sorted({float(q) for q in (quantiles or []) if 0 < float(q) < 1})
    if not quantiles:
        quantiles = [0.5]

    horizon = max(1, min(horizon, 365 if interval == "1d" else 240))

    values = close_series.astype(float).dropna().values
    if len(values) < 15:
        raise https_fn.HttpsError(
            https_fn.FunctionsErrorCode.FAILED_PRECONDITION,
            "Not enough history to generate forecast.",
        )

    log_returns = np.diff(np.log(values))
    recent = log_returns[-min(len(log_returns), 252) :]
    drift = float(np.mean(recent[-min(len(recent), 40) :])) if len(recent) else 0.0
    vol = float(np.std(recent[-min(len(recent), 120) :])) if len(recent) else 0.01
    vol = max(vol, 0.0008)

    rng = np.random.default_rng(42)
    sims = rng.normal(loc=drift, scale=vol, size=(1500, horizon)).cumsum(axis=1)
    sim_prices = values[-1] * np.exp(sims)

    freq = "H" if interval == "1h" else "B"
    dates = pd.date_range(close_series.index[-1], periods=horizon + 1, freq=freq)[1:]

    forecast_rows: list[dict[str, Any]] = []
    for idx, ts in enumerate(dates):
        row = {"ds": ts.isoformat()}
        for quantile in quantiles:
            row[f"q{int(round(quantile * 100)):02d}"] = round(float(np.quantile(sim_prices[:, idx], quantile)), 4)
        forecast_rows.append(row)

    ref_q = min(quantiles, key=lambda q: abs(q - 0.5)) if quantiles else 0.5
    median_key = f"q{int(round(ref_q * 100)):02d}"
    median_path = np.array([row.get(median_key, values[-1]) for row in forecast_rows], dtype=float)
    last_actual = float(values[-1])
    mae_recent = float(np.mean(np.abs(np.diff(values[-min(len(values), 90) :])))) if len(values) > 2 else 0.0

    return {
        "engine": "statistical_fallback",
        "status": "completed",
        "serviceMessage": "Fallback Monte Carlo quantile model used.",
        "forecastRows": forecast_rows,
        "metrics": {
            "lastClose": round(last_actual, 4),
            "mae": round(mae_recent, 4),
            "horizon": horizon,
            "medianEnd": round(float(median_path[-1]), 4),
            "drift": round(drift, 6),
            "volatility": round(vol, 6),
            "historyPoints": int(len(values)),
            "historyStart": str(close_series.index[0].isoformat()) if len(close_series.index) else "",
            "historyEnd": str(close_series.index[-1].isoformat()) if len(close_series.index) else "",
        },
    }


def _run_prophet_engine(close_series: pd.Series, horizon: int, quantiles: list[float], interval: str) -> dict[str, Any]:
    try:
        from prophet import Prophet  # type: ignore
    except Exception:
        return _generate_quantile_forecast(close_series, horizon, quantiles, interval)

    import numpy as np  # type: ignore
    import pandas as pd  # type: ignore

    forecast_core = _generate_quantile_forecast(close_series, horizon, quantiles, interval)
    try:
        df = pd.DataFrame({"ds": close_series.index.tz_localize(None), "y": close_series.values.astype(float)})
        is_hourly = interval == "1h"
        model = Prophet(
            daily_seasonality=is_hourly,
            weekly_seasonality=True,
            yearly_seasonality=not is_hourly,
            changepoint_prior_scale=0.08,
            seasonality_prior_scale=12.0,
            interval_width=0.8,
        )
        model.add_country_holidays(country_name="US")
        if not is_hourly:
            model.add_seasonality(name="monthly", period=30.5, fourier_order=5)

        model.fit(df)

        # In-sample diagnostics (tail window) so users can compare engines on fit quality.
        try:
            in_sample = model.predict(df[["ds"]])
            tail = min(len(df), 90 if interval != "1h" else 240)
            if tail >= 5:
                actual = df["y"].to_numpy(dtype=float)
                yhat = in_sample["yhat"].to_numpy(dtype=float)
                abs_err = np.abs(actual - yhat)
                mae = float(np.mean(abs_err[-tail:]))
                rmse = float(np.sqrt(np.mean((actual[-tail:] - yhat[-tail:]) ** 2)))
                denom = np.maximum(np.abs(actual[-tail:]), 1e-9)
                mape = float(np.mean(np.abs((actual[-tail:] - yhat[-tail:]) / denom)))

                lower = in_sample["yhat_lower"].to_numpy(dtype=float)
                upper = in_sample["yhat_upper"].to_numpy(dtype=float)
                coverage = float(np.mean((actual[-tail:] >= lower[-tail:]) & (actual[-tail:] <= upper[-tail:])))

                forecast_core["metrics"].update(
                    {
                        "mae": round(mae, 4),
                        "rmse": round(rmse, 4),
                        "mape": round(mape, 6),
                        "coverage10_90": round(coverage, 4),
                        "historyPoints": int(len(df)),
                        "historyStart": str(df["ds"].iloc[0].isoformat()) if len(df) else "",
                        "historyEnd": str(df["ds"].iloc[-1].isoformat()) if len(df) else "",
                        "diagnosticWindow": int(tail),
                    }
                )
        except Exception:
            pass

        freq = "H" if is_hourly else "B"
        periods = max(1, min(horizon, 365 if not is_hourly else 240))
        future = model.make_future_dataframe(periods=periods, freq=freq, include_history=False)
        forecast = model.predict(future)

        normal = NormalDist()
        z_80 = normal.inv_cdf(0.9)
        if z_80 == 0:
            z_80 = 1.28155

        rows: list[dict[str, Any]] = []
        quantiles = sorted({float(q) for q in (quantiles or []) if 0 < float(q) < 1})
        if not quantiles:
            quantiles = [0.5]

        for _, row in forecast.iterrows():
            yhat = float(row["yhat"])
            lower = float(row["yhat_lower"])
            upper = float(row["yhat_upper"])
            sigma = max((upper - lower) / (2 * z_80), 1e-6)
            out_row: dict[str, Any] = {"ds": row["ds"].isoformat()}
            for q in quantiles:
                z = normal.inv_cdf(q)
                out_row[f"q{int(round(q * 100)):02d}"] = round(yhat + sigma * z, 4)
            rows.append(out_row)

        forecast_core["engine"] = "prophet"
        forecast_core["serviceMessage"] = "Forecast generated with Quantura Horizon, quantiles derived from model uncertainty."
        forecast_core["forecastRows"] = rows
        ref_q = min(quantiles, key=lambda q: abs(q - 0.5)) if quantiles else 0.5
        ref_key = f"q{int(round(ref_q * 100)):02d}"
        forecast_core["metrics"]["medianEnd"] = rows[-1].get(ref_key) if rows else forecast_core["metrics"].get("medianEnd")

        return forecast_core
    except Exception:
        # Prophet is best-effort: keep the product usable by returning the fallback model output.
        return forecast_core


def _run_timemixer_engine(close_series: pd.Series, horizon: int, quantiles: list[float], interval: str) -> dict[str, Any]:
    payload = {
        "model_id": IBM_TIMEMIXER_MODEL_ID,
        "history": [round(float(value), 6) for value in close_series.tail(2048).tolist()],
        "horizon": int(horizon),
        "quantiles": quantiles,
        "interval": interval,
    }

    headers = {"Content-Type": "application/json"}
    if IBM_TIMEMIXER_API_KEY:
        headers["Authorization"] = f"Bearer {IBM_TIMEMIXER_API_KEY}"

    if IBM_TIMEMIXER_ENDPOINT:
        try:
            response = requests.post(IBM_TIMEMIXER_ENDPOINT, headers=headers, json=payload, timeout=30)
            response.raise_for_status()
            body = response.json()
            if isinstance(body, dict) and isinstance(body.get("forecastRows"), list):
                result = _generate_quantile_forecast(close_series, horizon, quantiles, interval)
                result["engine"] = "ibm_timemixer_endpoint"
                result["serviceMessage"] = "Forecast generated by configured IBM TimeMixer endpoint."
                result["forecastRows"] = body["forecastRows"]
                metrics = body.get("metrics") or {}
                result["metrics"].update({k: _serialize_for_firestore(v) for k, v in metrics.items()})
                return result
        except Exception:
            pass

    if HUGGINGFACEHUB_API_TOKEN:
        try:
            hf_headers = {
                "Authorization": f"Bearer {HUGGINGFACEHUB_API_TOKEN}",
                "Content-Type": "application/json",
            }
            hf_response = requests.post(
                f"https://api-inference.huggingface.co/models/{IBM_TIMEMIXER_MODEL_ID}",
                headers=hf_headers,
                json={"inputs": payload},
                timeout=45,
            )
            if hf_response.status_code < 400:
                body = hf_response.json()
                if isinstance(body, dict) and isinstance(body.get("forecastRows"), list):
                    result = _generate_quantile_forecast(close_series, horizon, quantiles, interval)
                    result["engine"] = "ibm_timemixer_hf"
                    result["serviceMessage"] = "Forecast generated via IBM TimeMixer Hugging Face endpoint."
                    result["forecastRows"] = body["forecastRows"]
                    return result
        except Exception:
            pass

    result = _generate_quantile_forecast(close_series, horizon, quantiles, interval)
    result["engine"] = "ibm_timemixer_proxy"
    result["serviceMessage"] = (
        "IBM TimeMixer service selected (model ibm-granite/granite-timeseries-ttm-r2). "
        "Proxy fallback model executed because no external TimeMixer endpoint is configured."
    )
    return result


def _run_forecast_service(
    service: str,
    close_series: pd.Series,
    horizon: int,
    quantiles: list[float],
    interval: str,
) -> dict[str, Any]:
    if service == "ibm_timemixer":
        return _run_timemixer_engine(close_series, horizon, quantiles, interval)
    return _run_prophet_engine(close_series, horizon, quantiles, interval)


def _forecast_preview_rows(rows: list[dict[str, Any]], max_rows: int = 12) -> list[dict[str, Any]]:
    return rows[:max_rows]


def _forecast_quantile_entries(rows: list[dict[str, Any]]) -> list[tuple[float, str]]:
    if not rows:
        return []
    keys = set()
    for row in rows[: min(len(rows), 200)]:
        for key in (row or {}).keys():
            if re.match(r"^q\d\d$", str(key)):
                keys.add(str(key))
    entries: list[tuple[float, str]] = []
    for key in sorted(keys):
        try:
            q = int(key[1:]) / 100.0
            entries.append((q, key))
        except Exception:
            continue
    return sorted(entries, key=lambda item: item[0])


def _build_forecast_trade_rationale(
    service: str,
    horizon: int,
    metrics: dict[str, Any] | None,
    forecast_rows: list[dict[str, Any]] | None,
) -> str:
    metrics = metrics or {}
    rows = forecast_rows or []
    last_close = _safe_float(metrics.get("lastClose")) or 0.0
    median_end = _safe_float(metrics.get("medianEnd"))
    drift = _safe_float(metrics.get("drift")) or 0.0
    coverage = _safe_float(metrics.get("coverage10_90"))
    volatility = _safe_float(metrics.get("volatility"))

    implied_return = None
    if last_close > 0 and median_end is not None:
        implied_return = (median_end - last_close) / last_close

    trend_phrase = "balanced trend profile"
    if implied_return is not None and implied_return > 0.08:
        trend_phrase = "strong upward trend profile"
    elif implied_return is not None and implied_return > 0:
        trend_phrase = "moderately positive trend profile"
    elif implied_return is not None and implied_return < -0.05:
        trend_phrase = "defensive/negative trend profile"
    elif drift > 0:
        trend_phrase = "positive drift profile"

    service_label = "Quantura Horizon" if str(service).strip().lower() == "prophet" else "forecast model"
    sentence_one = (
        f"{service_label} indicates a {trend_phrase} over the next {int(horizon)} periods "
        f"with median-path projection anchored around recent price behavior."
    )

    coverage_text = "coverage unavailable"
    if coverage is not None:
        coverage_text = f"{coverage * 100:.1f}% coverage in the 10-90% band"
    vol_text = "volatility regime unavailable"
    if volatility is not None:
        vol_text = f"volatility near {volatility * 100:.2f}%"

    sentence_two = f"Confidence framing is based on {coverage_text} and {vol_text}."
    if rows:
        sentence_two += " Decision trigger should be tied to the median path versus lower-band support."
    return f"{sentence_one} {sentence_two}"


def _extract_forecast_key_levels(rows: list[dict[str, Any]]) -> dict[str, float | None]:
    entries = _forecast_quantile_entries(rows)
    if not rows or not entries:
        return {"support": None, "median": None, "resistance": None}
    last_row = rows[-1] or {}
    low_key = entries[0][1]
    high_key = entries[-1][1]
    mid_key = min(entries, key=lambda item: abs(item[0] - 0.5))[1]
    return {
        "support": _safe_float(last_row.get(low_key)),
        "median": _safe_float(last_row.get(mid_key)),
        "resistance": _safe_float(last_row.get(high_key)),
    }


def _generate_forecast_chart_png(forecast_doc: dict[str, Any]) -> bytes:
    import matplotlib

    matplotlib.use("Agg")
    import matplotlib.pyplot as plt  # type: ignore
    import pandas as pd  # type: ignore

    rows = forecast_doc.get("forecastRows") if isinstance(forecast_doc.get("forecastRows"), list) else []
    if not rows:
        raise ValueError("No forecast rows are available to chart.")

    entries = _forecast_quantile_entries(rows)
    if not entries:
        raise ValueError("No quantile columns detected in forecast rows.")

    low_key = entries[0][1]
    high_key = entries[-1][1]
    mid_key = min(entries, key=lambda item: abs(item[0] - 0.5))[1]

    x_vals = []
    for idx, row in enumerate(rows):
        ds = row.get("ds")
        if isinstance(ds, str):
            try:
                x_vals.append(pd.to_datetime(ds))
                continue
            except Exception:
                pass
        x_vals.append(pd.Timestamp.utcnow() + pd.Timedelta(days=idx))

    low = [float(row.get(low_key) or 0.0) for row in rows]
    high = [float(row.get(high_key) or 0.0) for row in rows]
    median = [float(row.get(mid_key) or 0.0) for row in rows]

    fig = plt.figure(figsize=(11.2, 5.2), dpi=170, facecolor="#0f172a")
    ax = fig.add_subplot(111)
    ax.set_facecolor("#0f172a")
    ax.plot(x_vals, median, color="#3b82f6", linewidth=2.2, label=f"Median ({mid_key})")
    ax.fill_between(x_vals, low, high, color="#10b981", alpha=0.18, label=f"Band {low_key}-{high_key}")
    ax.plot(x_vals, low, color="#ef4444", linewidth=1.2, alpha=0.9, label=f"Lower ({low_key})")
    ax.plot(x_vals, high, color="#10b981", linewidth=1.2, alpha=0.9, label=f"Upper ({high_key})")

    ticker = str(forecast_doc.get("ticker") or "Ticker")
    service = str(forecast_doc.get("service") or "Forecast")
    service_label = "Quantura Horizon" if service.lower() == "prophet" else service
    ax.set_title(f"{ticker} • {service_label}", color="#f8fafc", fontsize=13, fontweight="bold")
    ax.grid(alpha=0.22, color="#93c5fd", linewidth=0.5)
    ax.tick_params(colors="#e2e8f0")
    for spine in ax.spines.values():
        spine.set_color("#334155")
    legend = ax.legend(loc="upper left", frameon=False, fontsize=8)
    for txt in legend.get_texts():
        txt.set_color("#e2e8f0")

    fig.tight_layout()
    out = BytesIO()
    fig.savefig(out, format="png", bbox_inches="tight", transparent=True)
    plt.close(fig)
    return out.getvalue()


def _render_forecast_report_html(
    forecast_doc: dict[str, Any],
    rationale: str,
    chart_png_b64: str,
) -> str:
    ticker = str(forecast_doc.get("ticker") or "Ticker")
    horizon = int(forecast_doc.get("horizon") or 0)
    interval = str(forecast_doc.get("interval") or "1d")
    created = datetime.now(timezone.utc).strftime("%Y-%m-%d %H:%M UTC")
    metrics = forecast_doc.get("metrics") if isinstance(forecast_doc.get("metrics"), dict) else {}
    key_levels = _extract_forecast_key_levels(forecast_doc.get("forecastRows") or [])
    service = str(forecast_doc.get("service") or "prophet")
    service_label = "Quantura Horizon" if service.lower() == "prophet" else service

    rows = [
        ("Last Close", metrics.get("lastClose")),
        ("Median End", metrics.get("medianEnd")),
        ("Coverage 10-90", metrics.get("coverage10_90")),
        ("Support", key_levels.get("support")),
        ("Resistance", key_levels.get("resistance")),
    ]
    metric_rows = "".join(
        f"<tr><td>{label}</td><td>{value if value is not None else '—'}</td></tr>"
        for label, value in rows
    )

    return f"""
<!doctype html>
<html>
  <head>
    <meta charset="utf-8" />
    <style>
      body {{
        margin: 0;
        font-family: Inter, Roboto, sans-serif;
        background: #0f172a;
        color: #e2e8f0;
        padding: 28px;
      }}
      .wrap {{
        border: 1px solid rgba(148, 163, 184, 0.35);
        border-radius: 18px;
        padding: 22px;
        background: linear-gradient(180deg, rgba(15, 23, 42, 0.98) 0%, rgba(30, 41, 59, 0.92) 100%);
      }}
      .brand {{
        font-weight: 800;
        letter-spacing: 0.08em;
        color: #93c5fd;
        margin-bottom: 8px;
      }}
      h1 {{
        margin: 0 0 4px;
        font-size: 24px;
      }}
      .meta {{
        margin: 0 0 16px;
        color: #cbd5e1;
        font-size: 13px;
      }}
      .mono {{
        font-family: "JetBrains Mono", ui-monospace, SFMono-Regular, Menlo, monospace;
      }}
      img {{
        width: 100%;
        border-radius: 12px;
        border: 1px solid rgba(148, 163, 184, 0.35);
        background: rgba(15, 23, 42, 0.6);
      }}
      table {{
        width: 100%;
        border-collapse: collapse;
        margin-top: 14px;
      }}
      td {{
        padding: 8px 10px;
        border-bottom: 1px solid rgba(148, 163, 184, 0.25);
        font-size: 13px;
      }}
      td:first-child {{
        color: #93c5fd;
      }}
      .rationale {{
        margin-top: 16px;
        background: rgba(30, 41, 59, 0.7);
        border: 1px solid rgba(59, 130, 246, 0.28);
        border-radius: 12px;
        padding: 12px;
        line-height: 1.5;
      }}
    </style>
  </head>
  <body>
    <div class="wrap">
      <div class="brand">QUANTURA EXECUTIVE BRIEF</div>
      <h1>{ticker} • {service_label}</h1>
      <div class="meta">Generated: {created} • Horizon: {horizon} • Interval: {interval}</div>
      <img src="data:image/png;base64,{chart_png_b64}" alt="Forecast chart" />
      <table>
        {metric_rows}
      </table>
      <div class="rationale"><strong>AI Rationale:</strong> {rationale}</div>
    </div>
  </body>
</html>
""".strip()


def _generate_forecast_pdf_bytes(html: str, fallback_title: str, rationale: str) -> bytes:
    try:
        from weasyprint import HTML  # type: ignore

        return HTML(string=html, base_url=PUBLIC_ORIGIN).write_pdf()
    except Exception:
        import matplotlib

        matplotlib.use("Agg")
        import matplotlib.pyplot as plt  # type: ignore

        fig = plt.figure(figsize=(11.0, 8.5), dpi=170, facecolor="#ffffff")
        ax = fig.add_subplot(111)
        ax.axis("off")
        ax.text(0.02, 0.94, "Quantura Executive Brief", fontsize=16, fontweight="bold")
        ax.text(0.02, 0.89, fallback_title, fontsize=12)
        ax.text(0.02, 0.83, rationale, fontsize=10, wrap=True)
        out = BytesIO()
        fig.savefig(out, format="pdf", bbox_inches="tight")
        plt.close(fig)
        return out.getvalue()


def _generate_forecast_pptx_bytes(
    forecast_doc: dict[str, Any],
    chart_png: bytes,
    rationale: str,
) -> bytes:
    from pptx import Presentation  # type: ignore
    from pptx.util import Inches, Pt  # type: ignore

    ticker = str(forecast_doc.get("ticker") or "Ticker")
    created = datetime.now(timezone.utc).strftime("%Y-%m-%d")
    metrics = forecast_doc.get("metrics") if isinstance(forecast_doc.get("metrics"), dict) else {}
    rows = forecast_doc.get("forecastRows") if isinstance(forecast_doc.get("forecastRows"), list) else []
    key_levels = _extract_forecast_key_levels(rows)

    prs = Presentation()

    # Slide 1: Title.
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = f"{ticker} • Quantura Horizon Brief"
    subtitle = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.5), Inches(1.2))
    tf = subtitle.text_frame
    tf.text = f"Date: {created}\nCurrent Price: {metrics.get('lastClose', '—')}"
    tf.paragraphs[0].font.size = Pt(20)

    # Slide 2: Analysis chart + key levels.
    slide2 = prs.slides.add_slide(prs.slide_layouts[5])
    slide2.shapes.title.text = "Forecast Analysis"
    slide2.shapes.add_picture(BytesIO(chart_png), Inches(0.6), Inches(1.2), width=Inches(8.4))
    box = slide2.shapes.add_textbox(Inches(9.1), Inches(1.4), Inches(3.0), Inches(3.8))
    t2 = box.text_frame
    t2.text = "Key Levels"
    p = t2.add_paragraph()
    p.text = f"Support: {key_levels.get('support') if key_levels.get('support') is not None else '—'}"
    p.level = 1
    p = t2.add_paragraph()
    p.text = f"Median: {key_levels.get('median') if key_levels.get('median') is not None else '—'}"
    p.level = 1
    p = t2.add_paragraph()
    p.text = f"Resistance: {key_levels.get('resistance') if key_levels.get('resistance') is not None else '—'}"
    p.level = 1

    # Slide 3: Rationale.
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])
    slide3.shapes.title.text = "Trade Rationale"
    rationale_box = slide3.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(11.3), Inches(4.6))
    t3 = rationale_box.text_frame
    t3.word_wrap = True
    sentences = [s.strip() for s in re.split(r"(?<=[.!?])\s+", rationale) if s.strip()]
    if not sentences:
        sentences = [rationale.strip() or "No rationale available."]
    t3.text = sentences[0]
    for sentence in sentences[1:]:
        p = t3.add_paragraph()
        p.text = sentence
        p.level = 0

    out = BytesIO()
    prs.save(out)
    return out.getvalue()


def _generate_and_store_forecast_report_assets(
    forecast_id: str,
    forecast_doc: dict[str, Any],
) -> dict[str, Any]:
    workspace_id = str(forecast_doc.get("userId") or forecast_doc.get("createdByUid") or "").strip() or "workspace"
    ticker = str(forecast_doc.get("ticker") or "ticker").upper()
    service = str(forecast_doc.get("service") or "prophet")
    service_label = "Quantura Horizon" if service.lower() == "prophet" else service

    chart_png = _generate_forecast_chart_png(forecast_doc)
    rationale = str(forecast_doc.get("tradeRationale") or "").strip()
    if not rationale:
        rationale = _build_forecast_trade_rationale(
            service=service,
            horizon=int(forecast_doc.get("horizon") or 0),
            metrics=forecast_doc.get("metrics") if isinstance(forecast_doc.get("metrics"), dict) else {},
            forecast_rows=forecast_doc.get("forecastRows") if isinstance(forecast_doc.get("forecastRows"), list) else [],
        )

    chart_b64 = base64.b64encode(chart_png).decode("utf-8")
    html = _render_forecast_report_html(forecast_doc, rationale, chart_b64)
    pdf_bytes = _generate_forecast_pdf_bytes(html, f"{ticker} • {service_label}", rationale)
    pptx_bytes = _generate_forecast_pptx_bytes(forecast_doc, chart_png, rationale)

    safe_prefix = re.sub(r"[^A-Z0-9_-]+", "_", f"{ticker}_{forecast_id}".upper())
    base_path = f"forecast_reports/{workspace_id}"
    chart_path = f"{base_path}/{safe_prefix}_chart.png"
    pdf_path = f"{base_path}/{safe_prefix}_executive_brief.pdf"
    pptx_path = f"{base_path}/{safe_prefix}_slide_deck.pptx"

    bucket = admin_storage.bucket(STORAGE_BUCKET)
    bucket.blob(chart_path).upload_from_string(chart_png, content_type="image/png")
    bucket.blob(pdf_path).upload_from_string(pdf_bytes, content_type="application/pdf")
    bucket.blob(pptx_path).upload_from_string(
        pptx_bytes,
        content_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
    )

    return {
        "tradeRationale": rationale,
        "reportAssets": {
            "chartPath": chart_path,
            "pdfPath": pdf_path,
            "pptxPath": pptx_path,
        },
    }


@https_fn.on_call()
def create_order(req: https_fn.CallableRequest) -> dict[str, Any]:
    token = _require_auth(req)
    data = req.data or {}

    product = str(data.get("product") or "Deep Forecast")
    currency = str(data.get("currency") or "USD")
    price = data.get("price")
    if not isinstance(price, (int, float)):
        price = DEFAULT_FORECAST_PRICE
    meta = data.get("meta") or {}
    if not isinstance(meta, dict):
        meta = {}

    order = {
        "userId": req.auth.uid,
        "userEmail": token.get("email"),
        "product": product,
        "price": price,
        "currency": currency,
        "status": "pending",
        "paymentProvider": "stripe",
        "paymentStatus": "unpaid",
        "stripeCheckoutSessionId": "",
        "stripePaymentIntentId": "",
        "stripeSubscriptionId": "",
        "fulfillmentNotes": "",
        "meta": meta,
        "createdAt": firestore.SERVER_TIMESTAMP,
        "updatedAt": firestore.SERVER_TIMESTAMP,
    }

    doc_ref = db.collection("orders").document()
    doc_ref.set(order)
    _audit_event(req.auth.uid, token.get("email"), "order_created", {"orderId": doc_ref.id, "product": product})

    return {"orderId": doc_ref.id}


@https_fn.on_call()
def create_stripe_checkout_session(req: https_fn.CallableRequest) -> dict[str, Any]:
    token = _require_auth(req)
    data = req.data or {}
    order_id = str(data.get("orderId") or "").strip()
    meta = data.get("meta") if isinstance(data.get("meta"), dict) else {}

    if not order_id:
        raise https_fn.HttpsError(https_fn.FunctionsErrorCode.INVALID_ARGUMENT, "Order ID is required.")

    context = _remote_config_context(req, token, meta)
    if not _remote_config_bool("stripe_checkout_enabled", True, context=context):
        raise https_fn.HttpsError(
            https_fn.FunctionsErrorCode.FAILED_PRECONDITION,
            "Checkout is temporarily disabled.",
        )

    stripe = _stripe_module()
    order_ref = db.collection("orders").document(order_id)
    snapshot = order_ref.get()
    if not snapshot.exists:
        raise https_fn.HttpsError(https_fn.FunctionsErrorCode.NOT_FOUND, "Order not found.")

    order = snapshot.to_dict() or {}
    owner_id = str(order.get("userId") or "")
    if token.get("email") != ADMIN_EMAIL and owner_id != req.auth.uid:
        raise https_fn.HttpsError(https_fn.FunctionsErrorCode.PERMISSION_DENIED, "Order access denied.")

    payment_status = str(order.get("paymentStatus") or "unpaid").strip().lower()
    if payment_status in {"paid", "succeeded"}:
        return {"orderId": order_id, "status": "paid"}

    product = str(order.get("product") or "Deep Forecast").strip() or "Deep Forecast"
    currency = str(order.get("currency") or "USD").strip().lower() or "usd"
    price = _safe_float(order.get("price"))
    if price is None:
        price = float(DEFAULT_FORECAST_PRICE)
    amount_cents = max(50, int(round(float(price) * 100)))

    normalized_name = product.lower()
    mode = "payment" if "forecast" in normalized_name else "subscription"
    recurring = {"interval": "month"} if mode == "subscription" else None

    price_data: dict[str, Any] = {
        "currency": currency,
        "unit_amount": amount_cents,
        "product_data": {"name": product},
    }
    if recurring:
        price_data["recurring"] = recurring

    success_url = f"{PUBLIC_ORIGIN}/purchase?checkout=success&orderId={order_id}&session_id={{CHECKOUT_SESSION_ID}}"
    cancel_url = f"{PUBLIC_ORIGIN}/purchase?checkout=cancel&orderId={order_id}"

    session_kwargs: dict[str, Any] = {
        "mode": mode,
        "line_items": [{"price_data": price_data, "quantity": 1}],
        "success_url": success_url,
        "cancel_url": cancel_url,
        "client_reference_id": order_id,
        "metadata": {
            "orderId": order_id,
            "userId": owner_id,
            "product": product,
        },
    }

    email = str(token.get("email") or "").strip()
    customer_id = _resolve_stripe_customer_id(
        stripe,
        req.auth.uid,
        email,
        create_if_missing=True,
    )
    if customer_id:
        session_kwargs["customer"] = customer_id
    elif email:
        session_kwargs["customer_email"] = email

    try:
        session = stripe.checkout.Session.create(**session_kwargs)
    except Exception as exc:
        raise https_fn.HttpsError(
            https_fn.FunctionsErrorCode.INTERNAL,
            "Unable to start checkout session.",
            {"error": str(exc)},
        )

    session_customer_id = _extract_checkout_customer_id(session) or customer_id
    if session_customer_id:
        _persist_user_stripe_customer_id(req.auth.uid, session_customer_id)

    order_ref.set(
        {
            "paymentProvider": "stripe",
            "paymentStatus": "checkout_created",
            "stripeCheckoutSessionId": str(getattr(session, "id", "") or ""),
            "stripeCustomerId": session_customer_id or "",
            "stripeMode": mode,
            "updatedAt": firestore.SERVER_TIMESTAMP,
        },
        merge=True,
    )
    _audit_event(
        req.auth.uid,
        token.get("email"),
        "stripe_checkout_session_created",
        {"orderId": order_id, "mode": mode, "currency": currency, "amountCents": amount_cents},
    )

    return {
        "orderId": order_id,
        "status": "created",
        "mode": mode,
        "sessionId": str(getattr(session, "id", "") or ""),
        "url": str(getattr(session, "url", "") or ""),
    }


@https_fn.on_call()
def confirm_stripe_checkout(req: https_fn.CallableRequest) -> dict[str, Any]:
    token = _require_auth(req)
    data = req.data or {}
    session_id = str(data.get("sessionId") or data.get("session_id") or "").strip()
    order_hint = str(data.get("orderId") or "").strip()

    if not session_id:
        raise https_fn.HttpsError(https_fn.FunctionsErrorCode.INVALID_ARGUMENT, "Stripe session ID is required.")

    stripe = _stripe_module()
    try:
        session = stripe.checkout.Session.retrieve(session_id, expand=["payment_intent", "subscription"])
    except Exception as exc:
        raise https_fn.HttpsError(
            https_fn.FunctionsErrorCode.NOT_FOUND,
            "Stripe session could not be loaded.",
            {"error": str(exc)},
        )

    session_dict: dict[str, Any] = session.to_dict() if hasattr(session, "to_dict") else dict(session)
    order_id = (
        str(session_dict.get("client_reference_id") or "").strip()
        or str((session_dict.get("metadata") or {}).get("orderId") or "").strip()
        or order_hint
    )
    if not order_id:
        raise https_fn.HttpsError(https_fn.FunctionsErrorCode.NOT_FOUND, "Order ID missing from Stripe session.")

    order_ref = db.collection("orders").document(order_id)
    snapshot = order_ref.get()
    if not snapshot.exists:
        raise https_fn.HttpsError(https_fn.FunctionsErrorCode.NOT_FOUND, "Order not found.")

    order = snapshot.to_dict() or {}
    owner_id = str(order.get("userId") or "")
    if token.get("email") != ADMIN_EMAIL and owner_id != req.auth.uid:
        raise https_fn.HttpsError(https_fn.FunctionsErrorCode.PERMISSION_DENIED, "Order access denied.")

    mode = str(session_dict.get("mode") or "")
    payment_status = str(session_dict.get("payment_status") or "").lower()
    session_status = str(session_dict.get("status") or "").lower()
    subscription_id = session_dict.get("subscription")
    payment_intent = session_dict.get("payment_intent")
    customer_id = _stripe_id(session_dict.get("customer"))

    is_paid = False
    if mode == "payment":
        is_paid = payment_status == "paid"
    elif mode == "subscription":
        is_paid = session_status == "complete" or bool(subscription_id)
    else:
        is_paid = payment_status == "paid"

    payment_intent_id = ""
    if isinstance(payment_intent, dict):
        payment_intent_id = str(payment_intent.get("id") or "")
    elif isinstance(payment_intent, str):
        payment_intent_id = payment_intent

    subscription_id_str = ""
    if isinstance(subscription_id, dict):
        subscription_id_str = str(subscription_id.get("id") or "")
    elif isinstance(subscription_id, str):
        subscription_id_str = subscription_id

    update_payload: dict[str, Any] = {
        "paymentProvider": "stripe",
        "stripeCheckoutSessionId": session_id,
        "stripePaymentIntentId": payment_intent_id,
        "stripeSubscriptionId": subscription_id_str,
        "stripeCustomerId": customer_id,
        "stripeMode": mode,
        "stripePaymentStatus": payment_status,
        "stripeSessionStatus": session_status,
        "updatedAt": firestore.SERVER_TIMESTAMP,
    }

    if is_paid:
        update_payload["paymentStatus"] = "paid"
        update_payload["paidAt"] = firestore.SERVER_TIMESTAMP
    else:
        update_payload["paymentStatus"] = "unpaid"

    order_ref.set(update_payload, merge=True)
    if owner_id and customer_id:
        _persist_user_stripe_customer_id(owner_id, customer_id)
    _audit_event(
        req.auth.uid,
        token.get("email"),
        "stripe_checkout_confirmed",
        {"orderId": order_id, "paid": bool(is_paid), "mode": mode, "sessionId": session_id},
    )

    return {
        "orderId": order_id,
        "paid": bool(is_paid),
        "mode": mode,
        "paymentStatus": update_payload["paymentStatus"],
        "product": order.get("product") or "",
        "currency": order.get("currency") or "USD",
        "price": order.get("price") or 0,
    }


@https_fn.on_call()
def create_stripe_billing_portal_session(req: https_fn.CallableRequest) -> dict[str, Any]:
    token = _require_auth(req)
    data = req.data or {}
    stripe = _stripe_module()

    uid = req.auth.uid
    email = str(token.get("email") or "").strip()
    return_url = _sanitize_portal_return_url(data.get("returnUrl"))

    customer_id = _resolve_stripe_customer_id(
        stripe,
        uid,
        email,
        create_if_missing=False,
    )

    if not customer_id:
        order_docs = db.collection("orders").where("userId", "==", uid).limit(200).stream()
        checkout_sessions: list[tuple[str, Any]] = []

        for doc in order_docs:
            payload = doc.to_dict() or {}
            provider = str(payload.get("paymentProvider") or "").strip().lower()
            if provider and provider != "stripe":
                continue

            order_customer_id = str(payload.get("stripeCustomerId") or "").strip()
            if order_customer_id:
                customer_id = order_customer_id
                break

            session_id = str(payload.get("stripeCheckoutSessionId") or "").strip()
            if session_id:
                checkout_sessions.append((session_id, doc.reference))

        if not customer_id:
            for session_id, order_ref in checkout_sessions[:25]:
                try:
                    session = stripe.checkout.Session.retrieve(session_id)
                except Exception:
                    continue
                session_customer_id = _extract_checkout_customer_id(session)
                if not session_customer_id:
                    continue
                customer_id = session_customer_id
                order_ref.set(
                    {
                        "stripeCustomerId": session_customer_id,
                        "updatedAt": firestore.SERVER_TIMESTAMP,
                    },
                    merge=True,
                )
                break

        if customer_id:
            _persist_user_stripe_customer_id(uid, customer_id)

    if not customer_id:
        customer_id = _resolve_stripe_customer_id(
            stripe,
            uid,
            email,
            create_if_missing=True,
        )

    if not customer_id:
        raise https_fn.HttpsError(
            https_fn.FunctionsErrorCode.FAILED_PRECONDITION,
            "Unable to resolve Stripe customer for billing portal.",
        )

    try:
        portal_session = stripe.billing_portal.Session.create(
            customer=customer_id,
            return_url=return_url,
        )
    except Exception as exc:
        raise https_fn.HttpsError(
            https_fn.FunctionsErrorCode.INTERNAL,
            "Unable to create Stripe billing portal session.",
            {"error": str(exc)},
        )

    _audit_event(
        uid,
        token.get("email"),
        "stripe_billing_portal_session_created",
        {"customerId": customer_id},
    )

    return {
        "customerId": customer_id,
        "url": str(getattr(portal_session, "url", "") or ""),
        "returnUrl": return_url,
    }


@https_fn.on_call()
def create_stripe_connect_onboarding_link(req: https_fn.CallableRequest) -> dict[str, Any]:
    token = _require_auth(req)
    data = req.data or {}
    workspace_id = str(data.get("workspaceId") or req.auth.uid or "").strip() or req.auth.uid
    _require_workspace_editor(workspace_id, req.auth.uid, token)

    stripe = _stripe_module()
    user_ref = db.collection("users").document(workspace_id)
    user_snap = user_ref.get()
    user_doc = user_snap.to_dict() if user_snap.exists else {}
    profile = user_doc.get("profile") if isinstance(user_doc.get("profile"), dict) else {}

    account_id = str(
        user_doc.get("stripeConnectAccountId")
        or profile.get("stripeConnectAccountId")
        or ""
    ).strip()

    if not account_id:
        account_kwargs: dict[str, Any] = {
            "type": "express",
            "metadata": {"workspaceId": workspace_id, "ownerUid": req.auth.uid},
        }
        email = str(token.get("email") or "").strip()
        if email:
            account_kwargs["email"] = email
        try:
            account = stripe.Account.create(**account_kwargs)
            account_id = str(
                getattr(account, "id", "")
                or (account.get("id") if isinstance(account, dict) else "")
                or ""
            ).strip()
        except Exception as exc:
            raise https_fn.HttpsError(
                https_fn.FunctionsErrorCode.INTERNAL,
                "Unable to create Stripe Connect account.",
                {"error": str(exc)},
            )
        if not account_id:
            raise https_fn.HttpsError(https_fn.FunctionsErrorCode.INTERNAL, "Stripe Connect account ID is missing.")

        user_ref.set(
            {
                "stripeConnectAccountId": account_id,
                "profile": {
                    **(profile if isinstance(profile, dict) else {}),
                    "stripeConnectAccountId": account_id,
                },
                "updatedAt": firestore.SERVER_TIMESTAMP,
            },
            merge=True,
        )

    try:
        link = stripe.AccountLink.create(
            account=account_id,
            refresh_url=f"{PUBLIC_ORIGIN}/forecasting?connect=retry",
            return_url=f"{PUBLIC_ORIGIN}/forecasting?connect=done",
            type="account_onboarding",
        )
    except Exception as exc:
        raise https_fn.HttpsError(
            https_fn.FunctionsErrorCode.INTERNAL,
            "Unable to start Stripe Connect onboarding.",
            {"error": str(exc)},
        )

    _audit_event(
        req.auth.uid,
        token.get("email"),
        "stripe_connect_onboarding_link_created",
        {"workspaceId": workspace_id, "accountId": account_id},
    )

    return {
        "workspaceId": workspace_id,
        "accountId": account_id,
        "url": str(getattr(link, "url", "") or (link.get("url") if isinstance(link, dict) else "") or ""),
        "platformFeePercent": STRIPE_CONNECT_PLATFORM_FEE_PERCENT,
    }


@https_fn.on_call()
def create_creator_support_checkout(req: https_fn.CallableRequest) -> dict[str, Any]:
    token = _require_auth(req)
    data = req.data or {}

    creator_workspace_id = str(data.get("creatorWorkspaceId") or "").strip()
    if not creator_workspace_id:
        raise https_fn.HttpsError(https_fn.FunctionsErrorCode.INVALID_ARGUMENT, "creatorWorkspaceId is required.")

    mode_raw = str(data.get("mode") or "tip").strip().lower()
    mode = "subscription" if mode_raw in {"subscription", "subscribe", "sub"} else "payment"
    currency = str(data.get("currency") or "usd").strip().lower() or "usd"
    amount_default = CREATOR_DEFAULT_SUBSCRIBE_USD if mode == "subscription" else CREATOR_DEFAULT_THANKS_USD
    amount = _safe_float(data.get("amountUsd"))
    if amount is None or amount <= 0:
        amount = float(amount_default)
    amount_cents = max(100, int(round(float(amount) * 100)))

    target_type = str(data.get("targetType") or "profile").strip().lower() or "profile"
    target_id = str(data.get("targetId") or "").strip()

    creator_ref = db.collection("users").document(creator_workspace_id)
    creator_snap = creator_ref.get()
    if not creator_snap.exists:
        raise https_fn.HttpsError(https_fn.FunctionsErrorCode.NOT_FOUND, "Creator profile not found.")
    creator_doc = creator_snap.to_dict() or {}
    creator_profile = creator_doc.get("profile") if isinstance(creator_doc.get("profile"), dict) else {}
    creator_name = str(
        creator_profile.get("username")
        or creator_doc.get("name")
        or creator_doc.get("email")
        or creator_workspace_id
    ).strip()

    connect_account_id = str(
        creator_doc.get("stripeConnectAccountId")
        or creator_profile.get("stripeConnectAccountId")
        or ""
    ).strip()

    stripe = _stripe_module()
    platform_fee_percent = max(1.0, min(float(STRIPE_CONNECT_PLATFORM_FEE_PERCENT), 40.0))
    line_item_name = (
        f"Subscribe to @{creator_name}"
        if mode == "subscription"
        else f"Send thanks to @{creator_name}"
    )
    if target_type == "screener" and target_id:
        line_item_name = (
            f"Subscribe to @{creator_name}'s screener"
            if mode == "subscription"
            else f"Send thanks for screener {target_id}"
        )

    success_url = f"{PUBLIC_ORIGIN}/forecasting?creator_checkout=success&session_id={{CHECKOUT_SESSION_ID}}"
    cancel_url = f"{PUBLIC_ORIGIN}/forecasting?creator_checkout=cancel"
    session_kwargs: dict[str, Any] = {
        "mode": mode,
        "line_items": [
            {
                "price_data": {
                    "currency": currency,
                    "unit_amount": amount_cents,
                    "product_data": {"name": line_item_name},
                    **({"recurring": {"interval": "month"}} if mode == "subscription" else {}),
                },
                "quantity": 1,
            }
        ],
        "success_url": success_url,
        "cancel_url": cancel_url,
        "metadata": {
            "supporterUid": req.auth.uid,
            "supporterEmail": str(token.get("email") or ""),
            "creatorWorkspaceId": creator_workspace_id,
            "creatorName": creator_name,
            "mode": mode,
            "targetType": target_type,
            "targetId": target_id,
            "platformFeePercent": str(platform_fee_percent),
            "connectLinked": "1" if connect_account_id else "0",
        },
    }

    email = str(token.get("email") or "").strip()
    customer_id = _resolve_stripe_customer_id(
        stripe,
        req.auth.uid,
        email,
        create_if_missing=True,
    )
    if customer_id:
        session_kwargs["customer"] = customer_id
    elif email:
        session_kwargs["customer_email"] = email

    if connect_account_id:
        if mode == "payment":
            fee_amount = max(1, int(round(amount_cents * (platform_fee_percent / 100.0))))
            session_kwargs["payment_intent_data"] = {
                "application_fee_amount": fee_amount,
                "transfer_data": {"destination": connect_account_id},
            }
        else:
            session_kwargs["subscription_data"] = {
                "application_fee_percent": platform_fee_percent,
                "transfer_data": {"destination": connect_account_id},
            }

    try:
        session = stripe.checkout.Session.create(**session_kwargs)
    except Exception as exc:
        raise https_fn.HttpsError(
            https_fn.FunctionsErrorCode.INTERNAL,
            "Unable to create creator checkout session.",
            {"error": str(exc)},
        )

    session_customer_id = _extract_checkout_customer_id(session) or customer_id
    if session_customer_id:
        _persist_user_stripe_customer_id(req.auth.uid, session_customer_id)

    doc_ref = db.collection("creator_support_sessions").document()
    doc_ref.set(
        {
            "supporterUid": req.auth.uid,
            "supporterEmail": token.get("email"),
            "creatorWorkspaceId": creator_workspace_id,
            "creatorName": creator_name,
            "targetType": target_type,
            "targetId": target_id,
            "mode": mode,
            "currency": currency,
            "amountUsd": round(float(amount), 2),
            "amountCents": amount_cents,
            "platformFeePercent": platform_fee_percent,
            "connectAccountId": connect_account_id,
            "stripeCustomerId": session_customer_id or "",
            "stripeCheckoutSessionId": str(getattr(session, "id", "") or ""),
            "status": "created",
            "createdAt": firestore.SERVER_TIMESTAMP,
            "updatedAt": firestore.SERVER_TIMESTAMP,
            "meta": data.get("meta") if isinstance(data.get("meta"), dict) else {},
        }
    )

    _audit_event(
        req.auth.uid,
        token.get("email"),
        "creator_support_checkout_created",
        {
            "creatorWorkspaceId": creator_workspace_id,
            "mode": mode,
            "targetType": target_type,
            "targetId": target_id,
            "amountCents": amount_cents,
            "connectLinked": bool(connect_account_id),
        },
    )

    return {
        "sessionId": str(getattr(session, "id", "") or ""),
        "url": str(getattr(session, "url", "") or ""),
        "mode": mode,
        "creatorWorkspaceId": creator_workspace_id,
        "creatorName": creator_name,
        "platformFeePercent": platform_fee_percent,
        "connectLinked": bool(connect_account_id),
    }


@https_fn.on_request()
def stripe_webhook(req: https_fn.Request) -> tuple[str, int]:
    if req.method != "POST":
        return ("Method not allowed", 405)

    if not STRIPE_WEBHOOK_SECRET:
        return ("Stripe webhook secret is not configured.", 500)

    stripe = _stripe_module()
    payload = req.get_data(cache=False, as_text=False)
    sig = req.headers.get("Stripe-Signature", "")

    try:
        event = stripe.Webhook.construct_event(payload=payload, sig_header=sig, secret=STRIPE_WEBHOOK_SECRET)
    except Exception:
        return ("Invalid signature", 400)

    try:
        event_type = str(event.get("type") or "")
        if event_type == "checkout.session.completed":
            session_obj = (event.get("data") or {}).get("object") or {}
            order_id = str(session_obj.get("client_reference_id") or "").strip()
            if not order_id:
                order_id = str((session_obj.get("metadata") or {}).get("orderId") or "").strip()

            if order_id:
                payment_status = str(session_obj.get("payment_status") or "").lower()
                session_status = str(session_obj.get("status") or "").lower()
                mode = str(session_obj.get("mode") or "")

                is_paid = payment_status == "paid" or session_status == "complete"
                customer_id = _stripe_id(session_obj.get("customer"))
                db.collection("orders").document(order_id).set(
                    {
                        "paymentProvider": "stripe",
                        "paymentStatus": "paid" if is_paid else "unpaid",
                        "stripeCheckoutSessionId": str(session_obj.get("id") or ""),
                        "stripePaymentIntentId": str(session_obj.get("payment_intent") or ""),
                        "stripeSubscriptionId": str(session_obj.get("subscription") or ""),
                        "stripeCustomerId": customer_id,
                        "stripeMode": mode,
                        "stripePaymentStatus": payment_status,
                        "stripeSessionStatus": session_status,
                        "paidAt": firestore.SERVER_TIMESTAMP if is_paid else firestore.DELETE_FIELD,
                        "updatedAt": firestore.SERVER_TIMESTAMP,
                    },
                    merge=True,
                )
    except Exception:
        # Webhook processing is best-effort; Stripe will retry on non-2xx, so keep this 200 unless signature fails.
        pass

    return ("ok", 200)


@https_fn.on_call()
def update_order_status(req: https_fn.CallableRequest) -> dict[str, Any]:
    token = _require_auth(req)
    _require_admin(token)

    data = req.data or {}
    order_id = data.get("orderId")
    status = data.get("status")
    notes = data.get("notes", "")

    if not order_id:
        raise https_fn.HttpsError(
            https_fn.FunctionsErrorCode.INVALID_ARGUMENT,
            "Order ID is required.",
        )

    if status not in ALLOWED_STATUSES:
        raise https_fn.HttpsError(
            https_fn.FunctionsErrorCode.INVALID_ARGUMENT,
            "Invalid order status.",
        )

    order_ref = db.collection("orders").document(order_id)
    snapshot = order_ref.get()
    if not snapshot.exists:
        raise https_fn.HttpsError(
            https_fn.FunctionsErrorCode.NOT_FOUND,
            "Order not found.",
        )
    order_data = snapshot.to_dict() or {}

    update_payload: dict[str, Any] = {
        "status": status,
        "fulfillmentNotes": notes,
        "updatedAt": firestore.SERVER_TIMESTAMP,
    }

    if status == "fulfilled":
        update_payload["fulfilledAt"] = firestore.SERVER_TIMESTAMP

    order_ref.update(update_payload)
    _audit_event(req.auth.uid, token.get("email"), "order_status_updated", {"orderId": order_id, "status": status})

    order_user_id = order_data.get("userId")
    order_user_email = order_data.get("userEmail")
    if isinstance(order_user_id, str) and order_user_id:
        product = str(order_data.get("product") or "Deep Forecast")
        human_status = status.replace("_", " ").title()
        _notify_user(
            order_user_id,
            order_user_email if isinstance(order_user_email, str) else None,
            title="Quantura order update",
            body=f"{product} order {order_id} is now {human_status}.",
            data={
                "type": "order_status",
                "orderId": order_id,
                "status": status,
                "url": "/dashboard",
            },
        )

    return {"orderId": order_id, "status": status}


@https_fn.on_call()
def submit_contact(req: https_fn.CallableRequest) -> dict[str, Any]:
    data = req.data or {}
    payload = {
        "name": str(data.get("name") or "").strip(),
        "email": str(data.get("email") or "").strip(),
        "company": str(data.get("company") or "").strip(),
        "role": str(data.get("role") or "").strip(),
        "category": str(data.get("category") or "").strip(),
        "message": str(data.get("message") or "").strip(),
        "sourcePage": str(data.get("sourcePage") or "").strip(),
        "utm": data.get("utm") or {},
        "meta": data.get("meta") or {},
        "createdAt": firestore.SERVER_TIMESTAMP,
    }

    missing = [field for field in CONTACT_REQUIRED_FIELDS if not payload[field]]
    if missing:
        raise https_fn.HttpsError(
            https_fn.FunctionsErrorCode.INVALID_ARGUMENT,
            "Missing required contact fields.",
            {"missing": missing},
        )

    doc_ref = db.collection("contacts").document()
    doc_ref.set(payload)

    return {"contactId": doc_ref.id}


@https_fn.on_call()
def create_collab_invite(req: https_fn.CallableRequest) -> dict[str, Any]:
    token = _require_auth(req)
    data = req.data or {}
    to_email_raw = str(data.get("email") or "").strip()
    to_email_norm = _normalize_email(to_email_raw)
    if not to_email_norm or "@" not in to_email_norm:
        raise https_fn.HttpsError(https_fn.FunctionsErrorCode.INVALID_ARGUMENT, "Valid email is required.")

    role = str(data.get("role") or "viewer").strip().lower()
    if role not in {"viewer", "editor"}:
        raise https_fn.HttpsError(https_fn.FunctionsErrorCode.INVALID_ARGUMENT, "Role must be viewer or editor.")

    from_email = _normalize_email(token.get("email"))
    invite_doc = {
        "fromUserId": req.auth.uid,
        "fromEmail": from_email,
        "workspaceUserId": req.auth.uid,
        "workspaceEmail": from_email,
        "toEmail": to_email_raw,
        "toEmailNorm": to_email_norm,
        "toUserId": "",
        "role": role,
        "status": "pending",
        "createdAt": firestore.SERVER_TIMESTAMP,
        "updatedAt": firestore.SERVER_TIMESTAMP,
    }

    doc_ref = db.collection("collab_invites").document()
    doc_ref.set(invite_doc)
    _audit_event(
        req.auth.uid,
        token.get("email"),
        "collab_invite_created",
        {"inviteId": doc_ref.id, "toEmail": to_email_norm, "role": role},
    )
    return {"inviteId": doc_ref.id, "status": "pending"}


@https_fn.on_call()
def list_collab_invites(req: https_fn.CallableRequest) -> dict[str, Any]:
    token = _require_auth(req)
    email_norm = _normalize_email(token.get("email"))
    if not email_norm:
        return {"invites": []}

    docs = (
        db.collection("collab_invites")
        .where("toEmailNorm", "==", email_norm)
        .limit(100)
        .stream()
    )

    invites: list[dict[str, Any]] = []
    for doc in docs:
        data = doc.to_dict() or {}
        if data.get("status") != "pending":
            continue
        invites.append(
            {
                "inviteId": doc.id,
                "fromEmail": data.get("fromEmail"),
                "workspaceUserId": data.get("workspaceUserId"),
                "workspaceEmail": data.get("workspaceEmail"),
                "role": data.get("role", "viewer"),
                "status": data.get("status", "pending"),
                "createdAt": data.get("createdAt"),
            }
        )

    return {"invites": _serialize_for_firestore(invites)}


@https_fn.on_call()
def accept_collab_invite(req: https_fn.CallableRequest) -> dict[str, Any]:
    token = _require_auth(req)
    data = req.data or {}
    invite_id = str(data.get("inviteId") or "").strip()
    if not invite_id:
        raise https_fn.HttpsError(https_fn.FunctionsErrorCode.INVALID_ARGUMENT, "Invite ID is required.")

    email_norm = _normalize_email(token.get("email"))
    invite_ref = db.collection("collab_invites").document(invite_id)
    snapshot = invite_ref.get()
    if not snapshot.exists:
        raise https_fn.HttpsError(https_fn.FunctionsErrorCode.NOT_FOUND, "Invite not found.")

    invite = snapshot.to_dict() or {}
    if invite.get("status") != "pending":
        raise https_fn.HttpsError(https_fn.FunctionsErrorCode.FAILED_PRECONDITION, "Invite is not pending.")
    if _normalize_email(invite.get("toEmailNorm")) != email_norm:
        raise https_fn.HttpsError(https_fn.FunctionsErrorCode.PERMISSION_DENIED, "Invite is not addressed to you.")

    workspace_user_id = str(invite.get("workspaceUserId") or invite.get("fromUserId") or "").strip()
    workspace_email = _normalize_email(invite.get("workspaceEmail") or invite.get("fromEmail"))
    if not workspace_user_id:
        raise https_fn.HttpsError(https_fn.FunctionsErrorCode.INTERNAL, "Invite is missing workspace info.")

    role = str(invite.get("role") or "viewer").strip().lower()
    if role not in {"viewer", "editor"}:
        role = "viewer"

    collab_uid = req.auth.uid
    collab_email = _normalize_email(token.get("email"))

    batch = db.batch()
    batch.set(
        db.collection("users").document(workspace_user_id).collection("collaborators").document(collab_uid),
        {
            "userId": collab_uid,
            "email": collab_email,
            "role": role,
            "addedAt": firestore.SERVER_TIMESTAMP,
            "addedBy": workspace_user_id,
        },
        merge=True,
    )
    batch.set(
        db.collection("users").document(collab_uid).collection("shared_workspaces").document(workspace_user_id),
        {
            "workspaceUserId": workspace_user_id,
            "workspaceEmail": workspace_email,
            "role": role,
            "addedAt": firestore.SERVER_TIMESTAMP,
        },
        merge=True,
    )
    batch.update(
        invite_ref,
        {
            "status": "accepted",
            "toUserId": collab_uid,
            "acceptedAt": firestore.SERVER_TIMESTAMP,
            "updatedAt": firestore.SERVER_TIMESTAMP,
        },
    )
    batch.commit()

    _audit_event(
        collab_uid,
        token.get("email"),
        "collab_invite_accepted",
        {"inviteId": invite_id, "workspaceUserId": workspace_user_id, "role": role},
    )
    return {"status": "accepted", "workspaceUserId": workspace_user_id, "workspaceEmail": workspace_email, "role": role}


@https_fn.on_call()
def revoke_collab_invite(req: https_fn.CallableRequest) -> dict[str, Any]:
    token = _require_auth(req)
    data = req.data or {}
    invite_id = str(data.get("inviteId") or "").strip()
    if not invite_id:
        raise https_fn.HttpsError(https_fn.FunctionsErrorCode.INVALID_ARGUMENT, "Invite ID is required.")

    invite_ref = db.collection("collab_invites").document(invite_id)
    snapshot = invite_ref.get()
    if not snapshot.exists:
        raise https_fn.HttpsError(https_fn.FunctionsErrorCode.NOT_FOUND, "Invite not found.")
    invite = snapshot.to_dict() or {}
    if invite.get("fromUserId") != req.auth.uid:
        raise https_fn.HttpsError(https_fn.FunctionsErrorCode.PERMISSION_DENIED, "Only the inviter can revoke.")
    if invite.get("status") != "pending":
        return {"inviteId": invite_id, "status": invite.get("status")}

    invite_ref.update({"status": "revoked", "updatedAt": firestore.SERVER_TIMESTAMP})
    _audit_event(req.auth.uid, token.get("email"), "collab_invite_revoked", {"inviteId": invite_id})
    return {"inviteId": invite_id, "status": "revoked"}


@https_fn.on_call()
def list_collaborators(req: https_fn.CallableRequest) -> dict[str, Any]:
    token = _require_auth(req)
    docs = db.collection("users").document(req.auth.uid).collection("collaborators").limit(200).stream()
    collaborators: list[dict[str, Any]] = []
    for doc in docs:
        data = doc.to_dict() or {}
        collaborators.append(
            {
                "userId": doc.id,
                "email": data.get("email"),
                "role": data.get("role", "viewer"),
                "addedAt": data.get("addedAt"),
            }
        )
    return {"collaborators": _serialize_for_firestore(collaborators)}


@https_fn.on_call()
def remove_collaborator(req: https_fn.CallableRequest) -> dict[str, Any]:
    token = _require_auth(req)
    data = req.data or {}
    collaborator_user_id = str(data.get("collaboratorUserId") or "").strip()
    if not collaborator_user_id:
        raise https_fn.HttpsError(https_fn.FunctionsErrorCode.INVALID_ARGUMENT, "Collaborator user ID is required.")

    owner_id = req.auth.uid
    batch = db.batch()
    batch.delete(db.collection("users").document(owner_id).collection("collaborators").document(collaborator_user_id))
    batch.delete(db.collection("users").document(collaborator_user_id).collection("shared_workspaces").document(owner_id))
    batch.commit()
    _audit_event(owner_id, token.get("email"), "collaborator_removed", {"collaboratorUserId": collaborator_user_id})
    return {"status": "removed", "collaboratorUserId": collaborator_user_id}


@https_fn.on_call()
def track_meta_conversion_event(req: https_fn.CallableRequest) -> dict[str, Any]:
    data = req.data or {}
    if not isinstance(data, dict):
        raise https_fn.HttpsError(
            https_fn.FunctionsErrorCode.INVALID_ARGUMENT,
            "Event payload must be an object.",
        )

    if not META_PIXEL_ID or not META_CAPI_ACCESS_TOKEN:
        return {"ok": False, "skipped": "meta_capi_not_configured"}

    origin = _meta_header(req, "origin")
    if origin and origin not in META_ALLOWED_ORIGINS:
        return {"ok": False, "skipped": "origin_not_allowed"}

    source_event_name = str(data.get("sourceEventName") or data.get("eventName") or "").strip()
    if not source_event_name:
        raise https_fn.HttpsError(
            https_fn.FunctionsErrorCode.INVALID_ARGUMENT,
            "eventName is required.",
        )
    source_event_name = re.sub(r"[^A-Za-z0-9_ \-]+", "", source_event_name)[:80] or "custom_event"

    event_name = str(data.get("eventName") or source_event_name).strip()
    event_name = "PageView" if event_name == "page_view" else event_name
    event_name = re.sub(r"[^A-Za-z0-9_ \-]+", "", event_name)[:80] or "custom_event"

    now_ts = int(time.time())
    event_time = int(_safe_float(data.get("eventTime")) or now_ts)
    if event_time < 946684800 or event_time > now_ts + 3600:
        event_time = now_ts

    event_id = str(data.get("eventId") or "").strip()
    if not event_id:
        seed = f"{source_event_name}:{event_time}:{time.time_ns()}"
        event_id = f"q_{hashlib.sha256(seed.encode('utf-8')).hexdigest()[:24]}"
    event_id = re.sub(r"[^A-Za-z0-9_-]+", "", event_id)[:100] or f"q_{int(time.time() * 1000)}"

    action_source = str(data.get("actionSource") or "website").strip().lower()
    valid_action_sources = {
        "website",
        "app",
        "phone_call",
        "chat",
        "physical_store",
        "system_generated",
        "business_messaging",
        "other",
    }
    if action_source not in valid_action_sources:
        action_source = "website"

    event_source_url = str(data.get("eventSourceUrl") or "").strip()
    if event_source_url and not event_source_url.startswith("http"):
        event_source_url = f"{PUBLIC_ORIGIN}{event_source_url if event_source_url.startswith('/') else '/' + event_source_url}"
    if not event_source_url:
        event_source_url = PUBLIC_ORIGIN

    user_data = _meta_build_user_data(req, data)
    custom_data = _meta_build_custom_data(data.get("params"))

    event_payload: dict[str, Any] = {
        "event_name": event_name,
        "event_time": event_time,
        "event_id": event_id,
        "event_source_url": event_source_url,
        "action_source": action_source,
        "user_data": user_data,
        "original_event_data": {
            "event_name": source_event_name,
            "event_time": event_time,
        },
    }
    if custom_data:
        event_payload["custom_data"] = custom_data

    payload = {"data": [event_payload]}
    endpoint = f"https://graph.facebook.com/{META_GRAPH_API_VERSION}/{META_PIXEL_ID}/events"

    try:
        response = requests.post(
            endpoint,
            params={"access_token": META_CAPI_ACCESS_TOKEN},
            json=payload,
            timeout=10,
        )
    except Exception as exc:
        return {"ok": False, "eventId": event_id, "error": f"request_failed: {exc}"}

    response_text = response.text[:1200]
    try:
        response_json = response.json()
    except Exception:
        response_json = {"raw": response_text}

    if response.status_code >= 400:
        return {
            "ok": False,
            "eventId": event_id,
            "status": response.status_code,
            "error": response_json,
        }

    return {
        "ok": True,
        "eventId": event_id,
        "status": response.status_code,
        "meta": _serialize_for_firestore(response_json),
    }


@https_fn.on_call()
def get_web_push_config(req: https_fn.CallableRequest) -> dict[str, Any]:
    token = _require_auth(req)
    meta = req.data.get("meta") if isinstance(req.data, dict) else {}
    context = _remote_config_context(req, token, meta if isinstance(meta, dict) else None)
    vapid_key = FCM_WEB_VAPID_KEY or _remote_config_param("webpush_vapid_key", "", context=context)
    if not vapid_key:
        raise https_fn.HttpsError(
            https_fn.FunctionsErrorCode.FAILED_PRECONDITION,
            "FCM_WEB_VAPID_KEY is not configured.",
        )
    return {"vapidKey": vapid_key}


@https_fn.on_call()
def get_feature_flags(req: https_fn.CallableRequest) -> dict[str, Any]:
    token = _require_auth(req)
    meta = req.data.get("meta") if isinstance(req.data, dict) else {}
    context = _remote_config_context(req, token, meta if isinstance(meta, dict) else None)
    tier_key, tier = _resolve_ai_tier(req.auth.uid, token, context=context)
    usage_tiers = _get_ai_usage_tiers(context=context)
    return {
        "forecastProphetEnabled": _remote_config_bool("forecast_prophet_enabled", True, context=context),
        "forecastTimeMixerEnabled": _remote_config_bool("forecast_timemixer_enabled", True, context=context),
        "watchlistEnabled": _remote_config_bool("watchlist_enabled", True, context=context),
        "enableSocialLeaderboard": _remote_config_bool("enable_social_leaderboard", True, context=context),
        "pushEnabled": _remote_config_bool("push_notifications_enabled", True, context=context),
        "webPushVapidKey": _remote_config_param("webpush_vapid_key", "", context=context),
        "forecastModelPrimary": _remote_config_param("forecast_model_primary", "Quantura Horizon", context=context),
        "promoBannerText": _remote_config_param("promo_banner_text", "", context=context),
        "maintenanceMode": _remote_config_bool("maintenance_mode", False, context=context),
        "volatilityThreshold": _remote_config_float("volatility_threshold", 0.05, context=context),
        "llmAllowedModels": _serialize_for_firestore(_get_llm_allowed_models(context=context)),
        "aiUsageTiers": _serialize_for_firestore(usage_tiers),
        "aiUsageTierKey": tier_key,
        "aiUsageTier": _serialize_for_firestore(tier),
    }


@https_fn.on_call()
def register_notification_token(req: https_fn.CallableRequest) -> dict[str, Any]:
    token = _require_auth(req)
    data = req.data or {}
    registration_token = str(data.get("token") or "").strip()
    meta = data.get("meta") if isinstance(data.get("meta"), dict) else {}
    if len(registration_token) < 20:
        raise https_fn.HttpsError(
            https_fn.FunctionsErrorCode.INVALID_ARGUMENT,
            "Valid notification token is required.",
        )

    doc_id = _token_doc_id(registration_token)
    doc_ref = db.collection("notification_tokens").document(doc_id)
    doc_ref.set(
        {
            "token": registration_token,
            "tokenHash": doc_id,
            "active": True,
            "userId": req.auth.uid,
            "userEmail": token.get("email"),
            "userAgent": str(meta.get("userAgent") or ""),
            "meta": meta,
            "forceRefresh": bool(data.get("forceRefresh")),
            "updatedAt": firestore.SERVER_TIMESTAMP,
            "lastSeenAt": firestore.SERVER_TIMESTAMP,
            "createdAt": firestore.SERVER_TIMESTAMP,
        },
        merge=True,
    )

    _audit_event(req.auth.uid, token.get("email"), "notification_token_registered", {"tokenHash": doc_id})
    return {"tokenHash": doc_id, "active": True}


@https_fn.on_call()
def unregister_notification_token(req: https_fn.CallableRequest) -> dict[str, Any]:
    token = _require_auth(req)
    data = req.data or {}
    registration_token = str(data.get("token") or "").strip()
    if len(registration_token) < 20:
        raise https_fn.HttpsError(
            https_fn.FunctionsErrorCode.INVALID_ARGUMENT,
            "Valid notification token is required.",
        )

    doc_id = _token_doc_id(registration_token)
    db.collection("notification_tokens").document(doc_id).set(
        {
            "active": False,
            "updatedAt": firestore.SERVER_TIMESTAMP,
        },
        merge=True,
    )
    _audit_event(req.auth.uid, token.get("email"), "notification_token_unregistered", {"tokenHash": doc_id})
    return {"tokenHash": doc_id, "active": False}


@https_fn.on_call()
def send_test_notification(req: https_fn.CallableRequest) -> dict[str, Any]:
    token = _require_auth(req)
    data = req.data or {}
    title = str(data.get("title") or "Quantura notification test")
    body = str(data.get("body") or "Your web push notifications are active.")
    payload_data = data.get("data") if isinstance(data.get("data"), dict) else {}
    explicit_token = _normalize_notification_token(data.get("token"))

    result = _notify_user(
        req.auth.uid,
        token.get("email"),
        title=title,
        body=body,
        data={
            **payload_data,
            "type": "test",
            "url": "/dashboard",
        },
        explicit_token=explicit_token or None,
    )
    result["requestedTokenProvided"] = bool(explicit_token)
    return result


@https_fn.on_call()
def check_price_alerts(req: https_fn.CallableRequest) -> dict[str, Any]:
    token = _require_auth(req)
    data = req.data or {}
    workspace_id = str(data.get("workspaceId") or req.auth.uid or "").strip()
    if not workspace_id:
        workspace_id = req.auth.uid

    # Allow scanning your own workspace, or a shared workspace you have access to.
    if workspace_id != req.auth.uid and token.get("email") != ADMIN_EMAIL:
        shared = (
            db.collection("users")
            .document(workspace_id)
            .collection("collaborators")
            .document(req.auth.uid)
            .get()
        )
        if not shared.exists:
            raise https_fn.HttpsError(
                https_fn.FunctionsErrorCode.PERMISSION_DENIED,
                "Workspace access required to scan alerts.",
            )

    alerts_ref = db.collection("users").document(workspace_id).collection("price_alerts")
    docs = (
        alerts_ref.where("createdByUid", "==", req.auth.uid)
        .where("active", "==", True)
        .limit(200)
        .stream()
    )
    alert_docs = [(doc, doc.to_dict() or {}) for doc in docs]
    if not alert_docs:
        return {"workspaceId": workspace_id, "checked": 0, "triggered": 0, "triggeredAlerts": []}

    alerts: list[dict[str, Any]] = []
    tickers: list[str] = []
    for doc, payload in alert_docs:
        ticker = str(payload.get("ticker") or "").upper().strip()
        if ticker:
            tickers.append(ticker)
        alerts.append({"ref": doc.reference, "id": doc.id, "data": payload, "ticker": ticker})

    price_map = _latest_close_prices(tickers)
    batch = db.batch()
    triggered: list[dict[str, Any]] = []

    for alert in alerts:
        ref = alert["ref"]
        payload = alert["data"] or {}
        ticker = alert["ticker"]
        condition = str(payload.get("condition") or "above").strip().lower()
        target = _safe_float(payload.get("targetPrice") or payload.get("target") or payload.get("price"))
        current = price_map.get(ticker)

        update_payload: dict[str, Any] = {
            "lastCheckedAt": firestore.SERVER_TIMESTAMP,
            "lastPrice": current,
            "updatedAt": firestore.SERVER_TIMESTAMP,
        }

        fired = False
        if current is not None and target is not None:
            if condition == "below" and current <= target:
                fired = True
            if condition != "below" and current >= target:
                fired = True

        if fired:
            update_payload.update(
                {
                    "active": False,
                    "status": "triggered",
                    "triggeredAt": firestore.SERVER_TIMESTAMP,
                    "triggeredPrice": current,
                }
            )
            triggered.append(
                {
                    "alertId": alert["id"],
                    "workspaceId": workspace_id,
                    "ticker": ticker,
                    "condition": condition,
                    "targetPrice": target,
                    "currentPrice": current,
                }
            )
            event_ref = db.collection("price_alert_events").document()
            batch.set(
                event_ref,
                {
                    "alertId": alert["id"],
                    "workspaceId": workspace_id,
                    "userId": req.auth.uid,
                    "userEmail": token.get("email"),
                    "ticker": ticker,
                    "condition": condition,
                    "targetPrice": target,
                    "currentPrice": current,
                    "createdAt": firestore.SERVER_TIMESTAMP,
                    "meta": data.get("meta") or {},
                },
            )

        batch.set(ref, update_payload, merge=True)

    batch.commit()

    if triggered:
        lines = []
        for item in triggered[:6]:
            t = item.get("ticker") or "?"
            cond = ">= " if item.get("condition") != "below" else "<= "
            lines.append(f"{t} {cond}{item.get('targetPrice')} (now {item.get('currentPrice')})")
        body = "Triggered: " + "; ".join(lines)
        _notify_user(
            req.auth.uid,
            token.get("email"),
            title="Quantura price alert",
            body=body,
            data={"type": "price_alert", "count": len(triggered), "url": "/dashboard#watchlist"},
        )

    _audit_event(
        req.auth.uid,
        token.get("email"),
        "price_alerts_checked",
        {"workspaceId": workspace_id, "checked": len(alerts), "triggered": len(triggered)},
    )

    return {
        "workspaceId": workspace_id,
        "checked": len(alerts),
        "triggered": len(triggered),
        "triggeredAlerts": _serialize_for_firestore(triggered),
        "prices": _serialize_for_firestore(price_map),
    }


def _handle_forecast_request(req: https_fn.CallableRequest, forced_service: str | None = None) -> dict[str, Any]:
    token = _require_auth(req)
    data = dict(req.data or {})
    if forced_service:
        data["service"] = forced_service

    workspace_id = str(data.get("workspaceId") or req.auth.uid or "").strip()
    if not workspace_id:
        workspace_id = req.auth.uid
    if workspace_id != req.auth.uid:
        _require_workspace_editor(workspace_id, req.auth.uid, token)

    ticker = str(data.get("ticker") or "").upper().strip()
    if not ticker:
        raise https_fn.HttpsError(https_fn.FunctionsErrorCode.INVALID_ARGUMENT, "Ticker is required.")

    service = str(data.get("service") or "prophet").strip().lower()
    if service not in FORECAST_SERVICES:
        raise https_fn.HttpsError(
            https_fn.FunctionsErrorCode.INVALID_ARGUMENT,
            "Invalid forecast service.",
            {"allowed": sorted(FORECAST_SERVICES)},
        )
    context = _remote_config_context(req, token, data.get("meta") if isinstance(data.get("meta"), dict) else None)
    if _remote_config_bool("maintenance_mode", False, context=context):
        _raise_structured_error(
            https_fn.FunctionsErrorCode.FAILED_PRECONDITION,
            "maintenance_mode",
            "Quantura is temporarily in maintenance mode.",
        )
    if service == "prophet" and not _remote_config_bool("forecast_prophet_enabled", True, context=context):
        raise https_fn.HttpsError(
            https_fn.FunctionsErrorCode.FAILED_PRECONDITION,
            "Quantura Horizon forecasting is currently disabled.",
        )
    if service == "ibm_timemixer" and not _remote_config_bool("forecast_timemixer_enabled", True, context=context):
        raise https_fn.HttpsError(
            https_fn.FunctionsErrorCode.FAILED_PRECONDITION,
            "IBM TimeMixer forecasting is currently disabled.",
        )

    interval = str(data.get("interval") or "1d")
    if interval not in {"1d", "1h"}:
        raise https_fn.HttpsError(https_fn.FunctionsErrorCode.INVALID_ARGUMENT, "Interval must be 1d or 1h.")

    try:
        horizon = int(data.get("horizon") or 90)
    except Exception:
        raise https_fn.HttpsError(https_fn.FunctionsErrorCode.INVALID_ARGUMENT, "Horizon must be an integer.")
    if horizon <= 0:
        raise https_fn.HttpsError(https_fn.FunctionsErrorCode.INVALID_ARGUMENT, "Horizon must be greater than 0.")
    quantiles = _parse_quantiles(data.get("quantiles"))
    start = data.get("start")

    try:
        history = _load_history(ticker=ticker, start=start, interval=interval)
    except https_fn.HttpsError:
        raise
    except Exception as exc:
        _raise_structured_error(
            https_fn.FunctionsErrorCode.NOT_FOUND,
            "ticker_not_found",
            "Unable to load market data for ticker.",
            {"ticker": ticker, "raw": str(exc)},
        )
    close_series = history["Close"].copy()

    try:
        result = _run_forecast_service(service, close_series, horizon, quantiles, interval)
    except https_fn.HttpsError:
        raise
    except Exception as service_exc:
        try:
            result = _generate_quantile_forecast(close_series, horizon, quantiles, interval)
            result["serviceMessage"] = "Forecast service failed; fallback model executed."
        except https_fn.HttpsError:
            raise
        except Exception as exc:
            _raise_structured_error(
                https_fn.FunctionsErrorCode.INTERNAL,
                "forecast_failed",
                "Forecast generation failed.",
                {"ticker": ticker, "service": service, "serviceRaw": str(service_exc), "raw": str(exc)},
            )

    trade_rationale = _build_forecast_trade_rationale(
        service=service,
        horizon=horizon,
        metrics=result.get("metrics") if isinstance(result.get("metrics"), dict) else {},
        forecast_rows=result.get("forecastRows") if isinstance(result.get("forecastRows"), list) else [],
    )

    request_doc = {
        "userId": workspace_id,
        "userEmail": token.get("email"),
        "createdByUid": req.auth.uid,
        "createdByEmail": token.get("email"),
        "ticker": ticker,
        "interval": interval,
        "horizon": horizon,
        "start": start,
        "quantiles": quantiles,
        "service": service,
        "engine": result.get("engine"),
        "status": result.get("status", "completed"),
        "serviceMessage": result.get("serviceMessage"),
        "metrics": _serialize_for_firestore(result.get("metrics") or {}),
        "forecastPreview": _serialize_for_firestore(_forecast_preview_rows(result.get("forecastRows") or [])),
        "forecastRows": _serialize_for_firestore(result.get("forecastRows") or []),
        "tradeRationale": trade_rationale,
        "reportStatus": "queued",
        "reportAssets": {},
        "meta": data.get("meta") or {},
        "utm": data.get("utm") or {},
        "createdAt": firestore.SERVER_TIMESTAMP,
        "updatedAt": firestore.SERVER_TIMESTAMP,
    }

    doc_ref = db.collection("forecast_requests").document()
    doc_ref.set(request_doc)

    _audit_event(
        req.auth.uid,
        token.get("email"),
        "forecast_requested",
        {
            "requestId": doc_ref.id,
            "ticker": ticker,
            "service": service,
            "engine": result.get("engine"),
            "workspaceId": workspace_id,
        },
    )

    metrics = result.get("metrics") or {}
    return {
        "requestId": doc_ref.id,
        "status": request_doc["status"],
        "service": service,
        "engine": request_doc["engine"],
        "serviceMessage": request_doc["serviceMessage"],
        "lastClose": metrics.get("lastClose"),
        "mae": metrics.get("mae"),
        "coverage10_90": metrics.get("coverage10_90", "n/a"),
        "forecastPreview": request_doc["forecastPreview"],
        "tradeRationale": trade_rationale,
        "reportStatus": "queued",
    }


@https_fn.on_call(memory=MemoryOption.GB_1, timeout_sec=180)
def run_timeseries_forecast(req: https_fn.CallableRequest) -> dict[str, Any]:
    return _handle_forecast_request(req, forced_service="prophet")


@https_fn.on_call(memory=MemoryOption.GB_1, timeout_sec=180)
def run_prophet_forecast(req: https_fn.CallableRequest) -> dict[str, Any]:
    return _handle_forecast_request(req, forced_service="prophet")


@https_fn.on_call()
def delete_forecast_request(req: https_fn.CallableRequest) -> dict[str, Any]:
    token = _require_auth(req)
    data = req.data or {}
    forecast_id = str(data.get("forecastId") or data.get("id") or "").strip()
    if not forecast_id:
        raise https_fn.HttpsError(https_fn.FunctionsErrorCode.INVALID_ARGUMENT, "Forecast ID is required.")

    ref = db.collection("forecast_requests").document(forecast_id)
    snap = ref.get()
    if not snap.exists:
        raise https_fn.HttpsError(https_fn.FunctionsErrorCode.NOT_FOUND, "Forecast not found.")

    doc = snap.to_dict() or {}
    workspace_id = str(doc.get("userId") or "").strip()
    if not workspace_id:
        if token.get("email") != ADMIN_EMAIL:
            raise https_fn.HttpsError(https_fn.FunctionsErrorCode.PERMISSION_DENIED, "Access denied.")
    else:
        _require_workspace_editor(workspace_id, req.auth.uid, token)

    ref.delete()
    _audit_event(req.auth.uid, token.get("email"), "forecast_deleted", {"forecastId": forecast_id, "workspaceId": workspace_id})
    return {"deleted": True, "forecastId": forecast_id}


@https_fn.on_call()
def generate_forecast_report_assets(req: https_fn.CallableRequest) -> dict[str, Any]:
    token = _require_auth(req)
    data = req.data or {}
    forecast_id = str(data.get("forecastId") or data.get("id") or "").strip()
    force = bool(data.get("force"))
    if not forecast_id:
        raise https_fn.HttpsError(https_fn.FunctionsErrorCode.INVALID_ARGUMENT, "Forecast ID is required.")

    ref = db.collection("forecast_requests").document(forecast_id)
    snap = ref.get()
    if not snap.exists:
        raise https_fn.HttpsError(https_fn.FunctionsErrorCode.NOT_FOUND, "Forecast not found.")

    doc = snap.to_dict() or {}
    workspace_id = str(doc.get("userId") or "").strip()
    if workspace_id:
        _require_workspace_editor(workspace_id, req.auth.uid, token)
    elif token.get("email") != ADMIN_EMAIL:
        raise https_fn.HttpsError(https_fn.FunctionsErrorCode.PERMISSION_DENIED, "Access denied.")

    existing_assets = doc.get("reportAssets") if isinstance(doc.get("reportAssets"), dict) else {}
    existing_status = str(doc.get("reportStatus") or "").strip().lower()
    if not force and existing_status == "ready" and existing_assets:
        return {
            "forecastId": forecast_id,
            "reportStatus": "ready",
            "reportAssets": _serialize_for_firestore(existing_assets),
            "tradeRationale": str(doc.get("tradeRationale") or "").strip(),
        }

    ref.set(
        {
            "reportStatus": "generating",
            "reportError": "",
            "reportUpdatedAt": firestore.SERVER_TIMESTAMP,
        },
        merge=True,
    )

    try:
        generated = _generate_and_store_forecast_report_assets(forecast_id, doc)
        ref.set(
            {
                "reportStatus": "ready",
                "reportAssets": generated.get("reportAssets") or {},
                "tradeRationale": generated.get("tradeRationale") or str(doc.get("tradeRationale") or "").strip(),
                "reportUpdatedAt": firestore.SERVER_TIMESTAMP,
            },
            merge=True,
        )
        _audit_event(req.auth.uid, token.get("email"), "forecast_report_generated", {"forecastId": forecast_id, "workspaceId": workspace_id})
        return {
            "forecastId": forecast_id,
            "reportStatus": "ready",
            "reportAssets": _serialize_for_firestore(generated.get("reportAssets") or {}),
            "tradeRationale": str(generated.get("tradeRationale") or "").strip(),
        }
    except Exception as error:
        ref.set(
            {
                "reportStatus": "failed",
                "reportError": str(error)[:700],
                "reportUpdatedAt": firestore.SERVER_TIMESTAMP,
            },
            merge=True,
        )
        raise https_fn.HttpsError(https_fn.FunctionsErrorCode.INTERNAL, f"Unable to generate report assets: {error}")


@scheduler_fn.on_schedule(
    schedule="*/30 * * * *",
    timezone=scheduler_fn.Timezone(SOCIAL_AUTOMATION_TIMEZONE),
)
def forecast_report_agent_scheduler(event: scheduler_fn.ScheduledEvent) -> None:
    del event
    docs = list(db.collection("forecast_requests").where("reportStatus", "==", "queued").limit(REPORT_AGENT_BATCH_SIZE).stream())
    for item in docs:
        ref = db.collection("forecast_requests").document(item.id)
        data = item.to_dict() or {}
        try:
            ref.set(
                {
                    "reportStatus": "generating",
                    "reportError": "",
                    "reportUpdatedAt": firestore.SERVER_TIMESTAMP,
                },
                merge=True,
            )
            generated = _generate_and_store_forecast_report_assets(item.id, data)
            ref.set(
                {
                    "reportStatus": "ready",
                    "reportAssets": generated.get("reportAssets") or {},
                    "tradeRationale": generated.get("tradeRationale") or str(data.get("tradeRationale") or "").strip(),
                    "reportUpdatedAt": firestore.SERVER_TIMESTAMP,
                },
                merge=True,
            )
        except Exception as error:
            ref.set(
                {
                    "reportStatus": "failed",
                    "reportError": str(error)[:700],
                    "reportUpdatedAt": firestore.SERVER_TIMESTAMP,
                },
                merge=True,
            )


@https_fn.on_call()
def delete_screener_run(req: https_fn.CallableRequest) -> dict[str, Any]:
    token = _require_auth(req)
    data = req.data or {}
    run_id = str(data.get("runId") or data.get("id") or "").strip()
    if not run_id:
        raise https_fn.HttpsError(https_fn.FunctionsErrorCode.INVALID_ARGUMENT, "Run ID is required.")

    ref = db.collection("screener_runs").document(run_id)
    snap = ref.get()
    if not snap.exists:
        raise https_fn.HttpsError(https_fn.FunctionsErrorCode.NOT_FOUND, "Screener run not found.")

    doc = snap.to_dict() or {}
    workspace_id = str(doc.get("userId") or "").strip()
    if not workspace_id:
        if token.get("email") != ADMIN_EMAIL:
            raise https_fn.HttpsError(https_fn.FunctionsErrorCode.PERMISSION_DENIED, "Access denied.")
    else:
        _require_workspace_editor(workspace_id, req.auth.uid, token)

    ref.delete()
    _audit_event(req.auth.uid, token.get("email"), "screener_deleted", {"runId": run_id, "workspaceId": workspace_id})
    return {"deleted": True, "runId": run_id}


@https_fn.on_call()
def rename_screener_run(req: https_fn.CallableRequest) -> dict[str, Any]:
    token = _require_auth(req)
    data = req.data or {}
    run_id = str(data.get("runId") or data.get("id") or "").strip()
    title = str(data.get("title") or "").strip()
    if not run_id:
        raise https_fn.HttpsError(https_fn.FunctionsErrorCode.INVALID_ARGUMENT, "Run ID is required.")
    if not title:
        raise https_fn.HttpsError(https_fn.FunctionsErrorCode.INVALID_ARGUMENT, "Title is required.")
    if len(title) > 180:
        raise https_fn.HttpsError(https_fn.FunctionsErrorCode.INVALID_ARGUMENT, "Title is too long.")

    ref = db.collection("screener_runs").document(run_id)
    snap = ref.get()
    if not snap.exists:
        raise https_fn.HttpsError(https_fn.FunctionsErrorCode.NOT_FOUND, "Screener run not found.")

    doc = snap.to_dict() or {}
    workspace_id = str(doc.get("userId") or "").strip()
    if not workspace_id:
        if token.get("email") != ADMIN_EMAIL:
            raise https_fn.HttpsError(https_fn.FunctionsErrorCode.PERMISSION_DENIED, "Access denied.")
    else:
        _require_workspace_editor(workspace_id, req.auth.uid, token)

    ref.set({"title": title, "updatedAt": firestore.SERVER_TIMESTAMP}, merge=True)
    _audit_event(req.auth.uid, token.get("email"), "screener_renamed", {"runId": run_id, "workspaceId": workspace_id})
    return {"updated": True, "runId": run_id, "title": title}


@https_fn.on_call()
def set_screener_public_visibility(req: https_fn.CallableRequest) -> dict[str, Any]:
    token = _require_auth(req)
    data = req.data or {}
    run_id = str(data.get("runId") or data.get("id") or "").strip()
    is_public = bool(data.get("isPublic"))
    if not run_id:
        raise https_fn.HttpsError(https_fn.FunctionsErrorCode.INVALID_ARGUMENT, "Run ID is required.")

    ref = db.collection("screener_runs").document(run_id)
    snap = ref.get()
    if not snap.exists:
        raise https_fn.HttpsError(https_fn.FunctionsErrorCode.NOT_FOUND, "Screener run not found.")

    doc = snap.to_dict() or {}
    workspace_id = str(doc.get("userId") or "").strip()
    if not workspace_id:
        raise https_fn.HttpsError(https_fn.FunctionsErrorCode.PERMISSION_DENIED, "Workspace owner is missing.")
    _require_workspace_editor(workspace_id, req.auth.uid, token)

    payload: dict[str, Any] = {
        "isPublic": is_public,
        "updatedAt": firestore.SERVER_TIMESTAMP,
    }
    if is_public:
        payload["publishedAt"] = firestore.SERVER_TIMESTAMP
        payload["publishedByUid"] = req.auth.uid
    else:
        payload["publishedAt"] = firestore.DELETE_FIELD
        payload["publishedByUid"] = firestore.DELETE_FIELD

    ref.set(payload, merge=True)
    _audit_event(
        req.auth.uid,
        token.get("email"),
        "screener_public_visibility_updated",
        {"runId": run_id, "workspaceId": workspace_id, "isPublic": is_public},
    )
    return {"runId": run_id, "isPublic": is_public, "workspaceId": workspace_id, "updated": True}


@https_fn.on_call()
def upsert_ai_agent_social_action(req: https_fn.CallableRequest) -> dict[str, Any]:
    token = _require_auth(req)
    data = req.data or {}

    workspace_id = str(data.get("workspaceId") or req.auth.uid or "").strip() or req.auth.uid
    if workspace_id != req.auth.uid:
        _require_workspace_editor(workspace_id, req.auth.uid, token)

    agent_id = str(data.get("agentId") or "").strip()
    action = str(data.get("action") or "").strip().lower()
    active = bool(data.get("active"))
    if not agent_id:
        raise https_fn.HttpsError(https_fn.FunctionsErrorCode.INVALID_ARGUMENT, "Agent ID is required.")
    if action not in {"follow", "like"}:
        raise https_fn.HttpsError(https_fn.FunctionsErrorCode.INVALID_ARGUMENT, "Action must be follow or like.")

    social_collection = "ai_agent_followers" if action == "follow" else "ai_agent_likes"
    count_field = "followersCount" if action == "follow" else "likesCount"
    user_doc_id = f"{agent_id}__{req.auth.uid}"

    agent_ref = db.collection("users").document(workspace_id).collection("ai_agents").document(agent_id)
    social_ref = db.collection("users").document(workspace_id).collection(social_collection).document(user_doc_id)

    txn = db.transaction()
    agent_snap = agent_ref.get(transaction=txn)
    if not agent_snap.exists:
        raise https_fn.HttpsError(https_fn.FunctionsErrorCode.NOT_FOUND, "AI Agent not found.")
    social_snap = social_ref.get(transaction=txn)

    agent_doc = agent_snap.to_dict() or {}
    current_count = int(agent_doc.get(count_field) or 0)
    now = firestore.SERVER_TIMESTAMP

    changed = False
    next_active = social_snap.exists
    if active and not social_snap.exists:
        txn.set(
            social_ref,
            {
                "agentId": agent_id,
                "workspaceId": workspace_id,
                "userId": req.auth.uid,
                "userEmail": token.get("email"),
                "action": action,
                "createdAt": now,
                "updatedAt": now,
                "meta": data.get("meta") if isinstance(data.get("meta"), dict) else {},
            },
            merge=True,
        )
        next_active = True
        changed = True
    elif not active and social_snap.exists:
        txn.delete(social_ref)
        next_active = False
        changed = True

    next_count = current_count
    if changed:
        next_count = current_count + 1 if next_active else max(0, current_count - 1)
        txn.set(agent_ref, {count_field: next_count, "updatedAt": now}, merge=True)

    txn.commit()

    # Follow action auto-enables default volatility alerts on holdings.
    if action == "follow" and next_active:
        holdings = agent_doc.get("holdings") if isinstance(agent_doc.get("holdings"), list) else []
        symbols: list[str] = []
        for item in holdings[:20]:
            if isinstance(item, str):
                symbol = str(item).upper().strip()
            elif isinstance(item, dict):
                symbol = str(item.get("symbol") or "").upper().strip()
            else:
                symbol = ""
            if symbol and re.match(r"^[A-Z0-9.\-]{1,12}$", symbol):
                symbols.append(symbol)
        symbols = list(dict.fromkeys(symbols))
        for symbol in symbols:
            alert_ref = db.collection("users").document(workspace_id).collection("price_alerts").document(
                f"volatility_follow_{agent_id}_{symbol}"
            )
            alert_ref.set(
                {
                    "workspaceId": workspace_id,
                    "userId": req.auth.uid,
                    "createdByUid": req.auth.uid,
                    "userEmail": token.get("email"),
                    "ticker": symbol,
                    "condition": "volatility",
                    "thresholdPercent": 0.05,
                    "notes": "Auto-enabled from AI Agent follow action.",
                    "active": True,
                    "triggered": False,
                    "createdAt": now,
                    "updatedAt": now,
                    "source": "ai_agent_follow",
                    "meta": data.get("meta") if isinstance(data.get("meta"), dict) else {},
                },
                merge=True,
            )

    _audit_event(
        req.auth.uid,
        token.get("email"),
        "ai_agent_social_action",
        {"workspaceId": workspace_id, "agentId": agent_id, "action": action, "active": next_active, "changed": changed},
    )

    return {
        "workspaceId": workspace_id,
        "agentId": agent_id,
        "action": action,
        "active": next_active,
        "countField": count_field,
        "count": next_count,
        "changed": changed,
    }


@https_fn.on_call()
def rename_prediction_upload(req: https_fn.CallableRequest) -> dict[str, Any]:
    token = _require_auth(req)
    _require_admin(token)
    data = req.data or {}
    upload_id = str(data.get("uploadId") or data.get("id") or "").strip()
    title = str(data.get("title") or "").strip()
    if not upload_id:
        raise https_fn.HttpsError(https_fn.FunctionsErrorCode.INVALID_ARGUMENT, "Upload ID is required.")
    if not title:
        raise https_fn.HttpsError(https_fn.FunctionsErrorCode.INVALID_ARGUMENT, "Title is required.")
    if len(title) > 180:
        raise https_fn.HttpsError(https_fn.FunctionsErrorCode.INVALID_ARGUMENT, "Title is too long.")

    ref = db.collection("prediction_uploads").document(upload_id)
    snap = ref.get()
    if not snap.exists:
        raise https_fn.HttpsError(https_fn.FunctionsErrorCode.NOT_FOUND, "Upload not found.")

    ref.set({"title": title, "updatedAt": firestore.SERVER_TIMESTAMP}, merge=True)
    _audit_event(req.auth.uid, token.get("email"), "predictions_renamed", {"uploadId": upload_id})
    return {"updated": True, "uploadId": upload_id, "title": title}


@https_fn.on_call()
def delete_prediction_upload(req: https_fn.CallableRequest) -> dict[str, Any]:
    token = _require_auth(req)
    _require_admin(token)
    data = req.data or {}
    upload_id = str(data.get("uploadId") or data.get("id") or "").strip()
    if not upload_id:
        raise https_fn.HttpsError(https_fn.FunctionsErrorCode.INVALID_ARGUMENT, "Upload ID is required.")

    ref = db.collection("prediction_uploads").document(upload_id)
    snap = ref.get()
    if not snap.exists:
        raise https_fn.HttpsError(https_fn.FunctionsErrorCode.NOT_FOUND, "Upload not found.")

    doc = snap.to_dict() or {}

    file_path = str(doc.get("filePath") or "").strip()
    if file_path:
        try:
            bucket = admin_storage.bucket(STORAGE_BUCKET)
            bucket.blob(file_path).delete()
        except Exception:
            # Best-effort cleanup; still delete Firestore metadata.
            pass

    ref.delete()
    _audit_event(req.auth.uid, token.get("email"), "predictions_deleted", {"uploadId": upload_id})
    return {"deleted": True, "uploadId": upload_id}


@https_fn.on_call()
def get_prediction_upload_csv(req: https_fn.CallableRequest) -> dict[str, Any]:
    token = _require_auth(req)
    _require_admin(token)
    data = req.data or {}
    upload_id = str(data.get("uploadId") or data.get("id") or "").strip()
    if not upload_id:
        raise https_fn.HttpsError(https_fn.FunctionsErrorCode.INVALID_ARGUMENT, "Upload ID is required.")

    max_bytes_raw = data.get("maxBytes")
    max_bytes = 2_000_000
    if max_bytes_raw is not None:
        try:
            max_bytes = int(max_bytes_raw)
        except Exception:
            max_bytes = 2_000_000
    max_bytes = max(25_000, min(max_bytes, 5_000_000))

    ref = db.collection("prediction_uploads").document(upload_id)
    snap = ref.get()
    if not snap.exists:
        raise https_fn.HttpsError(https_fn.FunctionsErrorCode.NOT_FOUND, "Upload not found.")

    doc = snap.to_dict() or {}

    file_path = str(doc.get("filePath") or "").strip()
    if not file_path:
        raise https_fn.HttpsError(https_fn.FunctionsErrorCode.FAILED_PRECONDITION, "Upload is missing a storage path.")

    try:
        bucket = admin_storage.bucket(STORAGE_BUCKET)
        blob = bucket.blob(file_path)
        if not blob.exists():
            raise https_fn.HttpsError(https_fn.FunctionsErrorCode.NOT_FOUND, "CSV file not found in storage.")

        truncated = False
        size = 0
        try:
            blob.reload()
            size = int(blob.size or 0)
        except Exception:
            size = 0

        data_bytes = b""
        if size and size > max_bytes:
            truncated = True
            try:
                with blob.open("rb") as handle:
                    data_bytes = handle.read(max_bytes)
            except Exception:
                data_bytes = blob.download_as_bytes()[:max_bytes]
        else:
            data_bytes = blob.download_as_bytes()
            if len(data_bytes) > max_bytes:
                truncated = True
                data_bytes = data_bytes[:max_bytes]

        text = data_bytes.decode("utf-8", errors="replace")
        _audit_event(req.auth.uid, token.get("email"), "predictions_csv_fetched", {"uploadId": upload_id, "truncated": truncated})
        return {
            "uploadId": upload_id,
            "title": doc.get("title") or "",
            "filePath": file_path,
            "bytes": len(data_bytes),
            "truncated": truncated,
            "csv": text,
        }
    except https_fn.HttpsError:
        raise
    except Exception:
        raise https_fn.HttpsError(https_fn.FunctionsErrorCode.INTERNAL, "Unable to fetch CSV.")


def _prediction_agent_fallback(summary: dict[str, Any], ticker: str) -> str:
    quantile = str(summary.get("selectedQuantileLabel") or summary.get("selectedQuantile") or "N/A")
    point = summary.get("pointForecastValue")
    first_date = str(summary.get("firstRowDate") or "")
    last_date = str(summary.get("lastRowDate") or "")
    relation = str(summary.get("relation") or "").strip()
    warning = str(summary.get("warningText") or "").strip()
    pieces = [
        f"Ticker: {ticker or 'N/A'}.",
        f"Selected quantile: {quantile}.",
        f"Point forecast on last weekday row: {point}.",
    ]
    if first_date and last_date:
        pieces.append(f"Prediction window: {first_date} -> {last_date}.")
    if relation:
        pieces.append(relation)
    if warning:
        pieces.append(f"Weekday warning: {warning}")
    pieces.append("OpenAI commentary is unavailable; this is the deterministic fallback summary.")
    return " ".join(pieces)


@https_fn.on_call()
def run_prediction_upload_agent(req: https_fn.CallableRequest) -> dict[str, Any]:
    token = _require_auth(req)
    _require_admin(token)
    data = req.data or {}
    upload_id = str(data.get("uploadId") or data.get("id") or "").strip()
    if not upload_id:
        raise https_fn.HttpsError(https_fn.FunctionsErrorCode.INVALID_ARGUMENT, "Upload ID is required.")

    ref = db.collection("prediction_uploads").document(upload_id)
    snap = ref.get()
    if not snap.exists:
        raise https_fn.HttpsError(https_fn.FunctionsErrorCode.NOT_FOUND, "Upload not found.")

    doc = snap.to_dict() or {}

    mapping_summary = data.get("mappingSummary") if isinstance(data.get("mappingSummary"), dict) else {}
    ticker = str(data.get("ticker") or mapping_summary.get("ticker") or doc.get("ticker") or "").upper().strip()
    if not ticker:
        ticker = "UNKNOWN"

    # Trim very long cells before sending to model.
    compact_summary = dict(mapping_summary or {})
    for key in ("firstRowText", "lastRowText", "relation", "warningText"):
        if key in compact_summary and compact_summary[key] is not None:
            compact_summary[key] = str(compact_summary[key])[:1200]

    fallback_text = _prediction_agent_fallback(compact_summary, ticker)
    analysis = fallback_text
    model_used = "template_fallback"
    provider = "fallback"

    if OPENAI_API_KEY:
        model_used = str(os.environ.get("PREDICTION_AGENT_MODEL") or "gpt-5-mini").strip() or "gpt-5-mini"
        system_prompt = (
            "You are a market data QA assistant for Quantura. "
            "Given CSV-derived quantile mapping metadata, write a concise analyst note. "
            "Do not provide financial advice, guarantees, or return promises. "
            "Output plain text only."
        )
        user_prompt = (
            "Write a compact report with these headings:\n"
            "1) Mapping result\n"
            "2) Weekday checks\n"
            "3) Risk notes\n"
            "4) Next validation step\n\n"
            f"Ticker: {ticker}\n"
            f"Upload title: {doc.get('title') or 'predictions.csv'}\n"
            f"Summary JSON: {json.dumps(compact_summary, ensure_ascii=True)}"
        )
        payload = {
            "model": model_used,
            "temperature": 0.2,
            "max_tokens": 450,
            "messages": [
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": user_prompt},
            ],
        }
        try:
            response = requests.post(
                "https://api.openai.com/v1/chat/completions",
                headers={
                    "Authorization": f"Bearer {OPENAI_API_KEY}",
                    "Content-Type": "application/json",
                },
                json=payload,
                timeout=30,
            )
            response.raise_for_status()
            body = response.json()
            message = body.get("choices", [{}])[0].get("message", {})
            content = str(message.get("content") or "").strip()
            if content:
                analysis = content
                provider = "openai"
            else:
                analysis = fallback_text
                provider = "fallback"
        except Exception:
            analysis = fallback_text
            provider = "fallback"

    _audit_event(
        req.auth.uid,
        token.get("email"),
        "predictions_agent_run",
        {"uploadId": upload_id, "provider": provider, "model": model_used, "ticker": ticker},
    )
    return {
        "uploadId": upload_id,
        "ticker": ticker,
        "analysis": analysis,
        "provider": provider,
        "model": model_used,
    }


def _is_paid_user(uid: str) -> bool:
    try:
        docs = db.collection("orders").where("userId", "==", uid).limit(25).stream()
        for doc in docs:
            data = doc.to_dict() or {}
            status = str(data.get("paymentStatus") or "").strip().lower()
            if status in {"paid", "succeeded"}:
                return True
    except Exception:
        return False
    return False


def _enforce_daily_usage(uid: str, feature_key: str, limit: int, limit_message: str | None = None) -> None:
    if limit <= 0:
        return
    day_key = datetime.now(timezone.utc).date().isoformat()
    ref = db.collection("usage_daily").document(f"{uid}_{day_key}")
    transaction = db.transaction()
    snap = ref.get(transaction=transaction)
    existing = snap.to_dict() if snap.exists else {}
    count = int(existing.get(feature_key) or 0)
    if count >= limit:
        raise https_fn.HttpsError(
            https_fn.FunctionsErrorCode.RESOURCE_EXHAUSTED,
            limit_message or "Daily usage limit reached.",
        )
    transaction.set(
        ref,
        {
            "userId": uid,
            "date": day_key,
            feature_key: count + 1,
            "updatedAt": firestore.SERVER_TIMESTAMP,
        },
        merge=True,
    )
    transaction.commit()


def _iso_week_key(now_utc: datetime | None = None) -> str:
    dt = now_utc or datetime.now(timezone.utc)
    iso = dt.isocalendar()
    return f"{iso.year}-W{iso.week:02d}"


def _enforce_weekly_usage(uid: str, feature_key: str, limit: int, limit_message: str | None = None) -> None:
    if limit <= 0:
        return
    week_key = _iso_week_key()
    ref = db.collection("usage_weekly").document(f"{uid}_{week_key}")
    transaction = db.transaction()
    snap = ref.get(transaction=transaction)
    existing = snap.to_dict() if snap.exists else {}
    count = int(existing.get(feature_key) or 0)
    if count >= limit:
        raise https_fn.HttpsError(
            https_fn.FunctionsErrorCode.RESOURCE_EXHAUSTED,
            limit_message or "Weekly usage limit reached.",
        )
    transaction.set(
        ref,
        {
            "userId": uid,
            "week": week_key,
            feature_key: count + 1,
            "updatedAt": firestore.SERVER_TIMESTAMP,
        },
        merge=True,
    )
    transaction.commit()


def _generate_backtest_code(payload: dict[str, Any]) -> str:
    ticker = str(payload.get("ticker") or "").upper().strip()
    interval = str(payload.get("interval") or "1d").strip()
    lookback_days = int(payload.get("lookbackDays") or 730)
    strategy = str(payload.get("strategy") or "sma_cross").strip().lower()
    params = payload.get("params") or {}
    cash = float(payload.get("cash") or 10_000)
    commission = float(payload.get("commission") or 0.0)

    # Keep emitted code self-contained so users can run it locally.
    fast = int(params.get("fast") or 20)
    slow = int(params.get("slow") or 50)
    rsi_period = int(params.get("rsiPeriod") or 14)
    oversold = float(params.get("oversold") or 30)
    exit_above = float(params.get("exitAbove") or 55)
    sl_pct = max(0.0, min(float(params.get("slPct") or params.get("sl") or 0.0), 0.5))
    tp_pct = max(0.0, min(float(params.get("tpPct") or params.get("tp") or 0.0), 0.5))
    strategy_class = "RsiReversion" if strategy == "rsi_reversion" else "SmaCross"

    if strategy == "rsi_reversion":
        strategy_code = f"""
class RsiReversion(Strategy):
    rsi_period = {rsi_period}
    oversold = {oversold}
    exit_above = {exit_above}
    sl_pct = {sl_pct}
    tp_pct = {tp_pct}

    def init(self):
        self.rsi = self.I(calc_rsi, self.data.Close, self.rsi_period)

    def _risk_args(self):
        last_close = float(self.data.Close[-1])
        out = {{}}
        if self.sl_pct > 0:
            out["sl"] = last_close * (1 - self.sl_pct)
        if self.tp_pct > 0:
            out["tp"] = last_close * (1 + self.tp_pct)
        return out

    def next(self):
        if not self.position and self.rsi[-1] < self.oversold:
            self.buy(**self._risk_args())
        elif self.position and self.rsi[-1] > self.exit_above:
            self.position.close()
"""
    else:
        strategy_code = f"""
class SmaCross(Strategy):
    fast = {fast}
    slow = {slow}
    sl_pct = {sl_pct}
    tp_pct = {tp_pct}

    def init(self):
        price = self.data.Close
        self.sma_fast = self.I(SMA, price, self.fast)
        self.sma_slow = self.I(SMA, price, self.slow)

    def _risk_args(self):
        last_close = float(self.data.Close[-1])
        out = {{}}
        if self.sl_pct > 0:
            out["sl"] = last_close * (1 - self.sl_pct)
        if self.tp_pct > 0:
            out["tp"] = last_close * (1 + self.tp_pct)
        return out

    def next(self):
        if crossover(self.sma_fast, self.sma_slow):
            self.position.close()
            self.buy(**self._risk_args())
        elif crossover(self.sma_slow, self.sma_fast):
            self.position.close()
"""

    return (
        f"""# Quantura backtest export (generated)
# pip install yfinance backtesting pandas matplotlib

from datetime import datetime, timedelta, timezone

import pandas as pd
import yfinance as yf
from backtesting import Backtest, Strategy
from backtesting.lib import crossover
from backtesting.test import SMA


def calc_rsi(series, period=14):
    s = pd.Series(series)
    delta = s.diff()
    up = delta.clip(lower=0).rolling(period).mean()
    down = (-delta.clip(upper=0)).rolling(period).mean()
    rs = up / down
    return 100 - (100 / (1 + rs))


{strategy_code.strip()}


def main():
    ticker = {ticker!r}
    interval = {interval!r}
    lookback_days = {lookback_days}
    end = datetime.now(timezone.utc)
    start = end - timedelta(days=lookback_days)

    df = yf.download(
        ticker,
        start=start.strftime("%Y-%m-%d"),
        end=end.strftime("%Y-%m-%d"),
        interval=interval,
        progress=False,
        auto_adjust=False,
        threads=False,
    )
    if isinstance(df.columns, pd.MultiIndex):
        df.columns = df.columns.get_level_values(0)
    df = df.dropna()
    df.index = pd.to_datetime(df.index, errors="coerce")
    if getattr(df.index, "tz", None) is not None:
        df.index = df.index.tz_localize(None)

    bt = Backtest(df, {strategy_class}, cash={cash}, commission={commission})
    stats = bt.run()
    print(stats)
    bt.plot()


if __name__ == "__main__":
    main()
"""
    ).strip()


def _render_backtest_tradingview_code(payload: dict[str, Any]) -> str:
    ticker = str(payload.get("ticker") or "").upper().strip()
    strategy = str(payload.get("strategy") or "sma_cross").strip().lower()
    params = payload.get("params") or {}
    cash = float(payload.get("cash") or 10_000)
    commission = float(payload.get("commission") or 0.0)

    fast = int(params.get("fast") or 20)
    slow = int(params.get("slow") or 50)
    rsi_period = int(params.get("rsiPeriod") or 14)
    oversold = float(params.get("oversold") or 30)
    exit_above = float(params.get("exitAbove") or 55)
    sl_pct = max(0.0, min(float(params.get("slPct") or params.get("sl") or 0.0), 0.5))
    tp_pct = max(0.0, min(float(params.get("tpPct") or params.get("tp") or 0.0), 0.5))

    label = "RSI Reversion" if strategy == "rsi_reversion" else "SMA Crossover"
    lines = [
        "//@version=5",
        f"strategy(\"Quantura {label} ({ticker})\", overlay=true, initial_capital={int(cash)}, commission_type=strategy.commission.percent, commission_value={commission * 100.0:.6f})",
        "",
        f"slPct = input.float({sl_pct:.6f}, \"Stop loss (decimal)\", step=0.001, minval=0.0)",
        f"tpPct = input.float({tp_pct:.6f}, \"Take profit (decimal)\", step=0.001, minval=0.0)",
    ]

    if strategy == "rsi_reversion":
        lines.extend(
            [
                f"rsiPeriod = input.int({rsi_period}, \"RSI Period\", minval=2, maxval=60)",
                f"oversold = input.float({oversold:.2f}, \"Oversold\", minval=1, maxval=49)",
                f"exitAbove = input.float({exit_above:.2f}, \"Exit Above\", minval=50, maxval=95)",
                "rsiValue = ta.rsi(close, rsiPeriod)",
                "longCondition = rsiValue < oversold",
                "exitCondition = rsiValue > exitAbove",
            ]
        )
    else:
        lines.extend(
            [
                f"fastSmaLen = input.int({fast}, \"Fast SMA\", minval=2, maxval=200)",
                f"slowSmaLen = input.int({slow}, \"Slow SMA\", minval=5, maxval=400)",
                "fastSma = ta.sma(close, fastSmaLen)",
                "slowSma = ta.sma(close, slowSmaLen)",
                "longCondition = ta.crossover(fastSma, slowSma)",
                "exitCondition = ta.crossunder(fastSma, slowSma)",
                "plot(fastSma, color=color.new(color.aqua, 0), title=\"Fast SMA\")",
                "plot(slowSma, color=color.new(color.orange, 0), title=\"Slow SMA\")",
            ]
        )

    lines.extend(
        [
            "",
            "if longCondition and strategy.position_size <= 0",
            "    strategy.entry(\"Long\", strategy.long)",
            "",
            "if strategy.position_size > 0",
            "    stopPrice = slPct > 0 ? strategy.position_avg_price * (1 - slPct) : na",
            "    takePrice = tpPct > 0 ? strategy.position_avg_price * (1 + tpPct) : na",
            "    strategy.exit(\"Risk Exit\", from_entry=\"Long\", stop=stopPrice, limit=takePrice)",
            "",
            "if exitCondition and strategy.position_size > 0",
            "    strategy.close(\"Long\")",
            "",
            "// Generated by Quantura Export Strategy Logic",
        ]
    )
    return "\n".join(lines).strip()


def _render_backtest_mq5_code(payload: dict[str, Any]) -> str:
    ticker = str(payload.get("ticker") or "").upper().strip()
    strategy = str(payload.get("strategy") or "sma_cross").strip().lower()
    params = payload.get("params") or {}
    commission = float(payload.get("commission") or 0.0)
    lookback_days = int(payload.get("lookbackDays") or 730)

    fast = int(params.get("fast") or 20)
    slow = int(params.get("slow") or 50)
    rsi_period = int(params.get("rsiPeriod") or 14)
    oversold = float(params.get("oversold") or 30)
    exit_above = float(params.get("exitAbove") or 55)
    sl_pct = max(0.0, min(float(params.get("slPct") or params.get("sl") or 0.0), 0.5))
    tp_pct = max(0.0, min(float(params.get("tpPct") or params.get("tp") or 0.0), 0.5))

    lines = [
        "#property strict",
        "#include <Trade/Trade.mqh>",
        "",
        f"// Quantura template for {ticker} ({strategy})",
        f"// Backtest lookback used in Quantura: {lookback_days} days",
        f"// Commission in Quantura run: {commission:.6f}",
        "input double LotSize = 0.10;",
        f"input double StopLossPct = {sl_pct:.6f};",
        f"input double TakeProfitPct = {tp_pct:.6f};",
        "",
        "CTrade trade;",
        "double ReadBuffer(int handle, int shift)",
        "{",
        "  double values[];",
        "  if(CopyBuffer(handle, 0, shift, 1, values) <= 0) return EMPTY_VALUE;",
        "  return values[0];",
        "}",
        "",
        "void ApplyRisk()",
        "{",
        "  if(!PositionSelect(_Symbol)) return;",
        "  double openPrice = PositionGetDouble(POSITION_PRICE_OPEN);",
        "  double sl = 0.0;",
        "  double tp = 0.0;",
        "  if(StopLossPct > 0.0) sl = NormalizeDouble(openPrice * (1.0 - StopLossPct), _Digits);",
        "  if(TakeProfitPct > 0.0) tp = NormalizeDouble(openPrice * (1.0 + TakeProfitPct), _Digits);",
        "  trade.PositionModify(_Symbol, sl, tp);",
        "}",
        "",
    ]

    if strategy == "rsi_reversion":
        lines.extend(
            [
                f"input int RsiPeriod = {rsi_period};",
                f"input double Oversold = {oversold:.2f};",
                f"input double ExitAbove = {exit_above:.2f};",
                "int rsiHandle = INVALID_HANDLE;",
                "",
                "int OnInit()",
                "{",
                "  rsiHandle = iRSI(_Symbol, PERIOD_CURRENT, RsiPeriod, PRICE_CLOSE);",
                "  if(rsiHandle == INVALID_HANDLE) return(INIT_FAILED);",
                "  return(INIT_SUCCEEDED);",
                "}",
                "",
                "void OnDeinit(const int reason)",
                "{",
                "  if(rsiHandle != INVALID_HANDLE) IndicatorRelease(rsiHandle);",
                "}",
                "",
                "void OnTick()",
                "{",
                "  double rsiNow = ReadBuffer(rsiHandle, 0);",
                "  if(rsiNow == EMPTY_VALUE) return;",
                "  bool hasPosition = PositionSelect(_Symbol);",
                "  if(!hasPosition && rsiNow < Oversold)",
                "  {",
                "    trade.Buy(LotSize, _Symbol);",
                "    ApplyRisk();",
                "  }",
                "  if(hasPosition && rsiNow > ExitAbove)",
                "  {",
                "    trade.PositionClose(_Symbol);",
                "  }",
                "}",
            ]
        )
    else:
        lines.extend(
            [
                f"input int FastSma = {fast};",
                f"input int SlowSma = {slow};",
                "int fastHandle = INVALID_HANDLE;",
                "int slowHandle = INVALID_HANDLE;",
                "",
                "int OnInit()",
                "{",
                "  fastHandle = iMA(_Symbol, PERIOD_CURRENT, FastSma, 0, MODE_SMA, PRICE_CLOSE);",
                "  slowHandle = iMA(_Symbol, PERIOD_CURRENT, SlowSma, 0, MODE_SMA, PRICE_CLOSE);",
                "  if(fastHandle == INVALID_HANDLE || slowHandle == INVALID_HANDLE) return(INIT_FAILED);",
                "  return(INIT_SUCCEEDED);",
                "}",
                "",
                "void OnDeinit(const int reason)",
                "{",
                "  if(fastHandle != INVALID_HANDLE) IndicatorRelease(fastHandle);",
                "  if(slowHandle != INVALID_HANDLE) IndicatorRelease(slowHandle);",
                "}",
                "",
                "void OnTick()",
                "{",
                "  double fastNow = ReadBuffer(fastHandle, 0);",
                "  double fastPrev = ReadBuffer(fastHandle, 1);",
                "  double slowNow = ReadBuffer(slowHandle, 0);",
                "  double slowPrev = ReadBuffer(slowHandle, 1);",
                "  if(fastNow == EMPTY_VALUE || fastPrev == EMPTY_VALUE || slowNow == EMPTY_VALUE || slowPrev == EMPTY_VALUE) return;",
                "  bool crossUp = fastPrev <= slowPrev && fastNow > slowNow;",
                "  bool crossDown = fastPrev >= slowPrev && fastNow < slowNow;",
                "  bool hasPosition = PositionSelect(_Symbol);",
                "  if(crossUp && !hasPosition)",
                "  {",
                "    trade.Buy(LotSize, _Symbol);",
                "    ApplyRisk();",
                "  }",
                "  if(crossDown && hasPosition)",
                "  {",
                "    trade.PositionClose(_Symbol);",
                "  }",
                "}",
            ]
        )

    return "\n".join(lines).strip()


def _render_tradelocker_payload(payload: dict[str, Any]) -> str:
    ticker = str(payload.get("ticker") or "").upper().strip()
    interval = str(payload.get("interval") or "1d").strip()
    strategy = str(payload.get("strategy") or "sma_cross").strip().lower()
    params = payload.get("params") or {}
    cash = float(payload.get("cash") or 10_000)
    commission = float(payload.get("commission") or 0.0)
    sl_pct = max(0.0, min(float(params.get("slPct") or params.get("sl") or 0.0), 0.5))
    tp_pct = max(0.0, min(float(params.get("tpPct") or params.get("tp") or 0.0), 0.5))

    if strategy == "rsi_reversion":
        entry_rules = [f"RSI({int(params.get('rsiPeriod') or 14)}) < {float(params.get('oversold') or 30):.2f}"]
        exit_rules = [f"RSI({int(params.get('rsiPeriod') or 14)}) > {float(params.get('exitAbove') or 55):.2f}"]
    else:
        entry_rules = [f"SMA({int(params.get('fast') or 20)}) crosses above SMA({int(params.get('slow') or 50)})"]
        exit_rules = [f"SMA({int(params.get('fast') or 20)}) crosses below SMA({int(params.get('slow') or 50)})"]

    doc = {
        "provider": "tradelocker",
        "version": "1.0",
        "strategyName": f"quantura_{strategy}",
        "symbol": ticker,
        "interval": interval,
        "execution": {
            "side": "buy",
            "orderType": "market",
            "quantityType": "notional",
            "notionalUsd": round(cash, 2),
            "commissionRate": round(commission, 6),
        },
        "risk": {
            "stopLossPct": round(sl_pct, 6),
            "takeProfitPct": round(tp_pct, 6),
        },
        "entryRules": entry_rules,
        "exitRules": exit_rules,
        "notes": [
            "Generated by Quantura Export Strategy Logic.",
            "Map these rules into your TradeLocker webhook bot or strategy runner.",
        ],
    }
    return json.dumps(doc, indent=2, sort_keys=False)


def _build_backtest_export_sources(payload: dict[str, Any]) -> dict[str, dict[str, str]]:
    ticker = str(payload.get("ticker") or "backtest").upper().strip() or "backtest"
    strategy = str(payload.get("strategy") or "strategy").strip().lower()
    safe_prefix = re.sub(r"[^A-Z0-9_-]+", "_", f"{ticker}_{strategy}")
    python_code = _generate_backtest_code(payload)

    return {
        "python": {
            "label": "Python (.py)",
            "filename": f"{safe_prefix}.py",
            "mimeType": "text/x-python",
            "content": python_code,
        },
        "tradingview": {
            "label": "TradingView (.pine)",
            "filename": f"{safe_prefix}.pine",
            "mimeType": "text/plain",
            "content": _render_backtest_tradingview_code(payload),
        },
        "metatrader5": {
            "label": "MetaTrader 5 (.mq5)",
            "filename": f"{safe_prefix}.mq5",
            "mimeType": "text/plain",
            "content": _render_backtest_mq5_code(payload),
        },
        "tradelocker": {
            "label": "TradeLocker (JSON)",
            "filename": f"{safe_prefix}_tradelocker.json",
            "mimeType": "application/json",
            "content": _render_tradelocker_payload(payload),
        },
    }


@https_fn.on_call(memory=MemoryOption.GB_1, timeout_sec=180)
def run_backtest(req: https_fn.CallableRequest) -> dict[str, Any]:
    token = _require_auth(req)
    data = req.data or {}
    meta = data.get("meta") if isinstance(data.get("meta"), dict) else {}
    context = _remote_config_context(req, token, meta)

    if not _remote_config_bool("backtesting_enabled", True, context=context):
        raise https_fn.HttpsError(https_fn.FunctionsErrorCode.FAILED_PRECONDITION, "Backtesting is temporarily disabled.")

    ticker = str(data.get("ticker") or "").upper().strip()
    if not ticker or len(ticker) > 12 or not all(ch.isalnum() or ch in {".", "-"} for ch in ticker):
        raise https_fn.HttpsError(https_fn.FunctionsErrorCode.INVALID_ARGUMENT, "Ticker is required.")

    interval = str(data.get("interval") or "1d").strip()
    if interval not in {"1d", "1h"}:
        raise https_fn.HttpsError(https_fn.FunctionsErrorCode.INVALID_ARGUMENT, "Invalid interval.")

    strategy = str(data.get("strategy") or "sma_cross").strip().lower()
    if strategy not in {"sma_cross", "rsi_reversion"}:
        raise https_fn.HttpsError(https_fn.FunctionsErrorCode.INVALID_ARGUMENT, "Invalid strategy.")

    lookback_days = int(data.get("lookbackDays") or data.get("lookback") or 730)
    if lookback_days < 30 or lookback_days > 5000:
        raise https_fn.HttpsError(https_fn.FunctionsErrorCode.INVALID_ARGUMENT, "Lookback must be between 30 and 5000 days.")

    cash = float(data.get("cash") or 10_000)
    commission = float(data.get("commission") or 0.0)
    if cash <= 0 or cash > 10_000_000:
        raise https_fn.HttpsError(https_fn.FunctionsErrorCode.INVALID_ARGUMENT, "Invalid cash amount.")
    if commission < 0 or commission > 0.05:
        raise https_fn.HttpsError(https_fn.FunctionsErrorCode.INVALID_ARGUMENT, "Invalid commission.")

    params_raw = data.get("params") if isinstance(data.get("params"), dict) else {}
    sl_pct = float(params_raw.get("slPct") or params_raw.get("sl") or 0.0)
    tp_pct = float(params_raw.get("tpPct") or params_raw.get("tp") or 0.0)
    if sl_pct < 0 or sl_pct > 0.5:
        raise https_fn.HttpsError(https_fn.FunctionsErrorCode.INVALID_ARGUMENT, "Stop-loss percent must be between 0 and 0.5.")
    if tp_pct < 0 or tp_pct > 0.5:
        raise https_fn.HttpsError(https_fn.FunctionsErrorCode.INVALID_ARGUMENT, "Take-profit percent must be between 0 and 0.5.")

    params: dict[str, Any] = {}
    if strategy == "sma_cross":
        fast = int(params_raw.get("fast") or 20)
        slow = int(params_raw.get("slow") or 50)
        if fast < 2 or slow < 5 or fast >= slow or slow > 400:
            raise https_fn.HttpsError(https_fn.FunctionsErrorCode.INVALID_ARGUMENT, "Invalid SMA parameters.")
        params = {"fast": fast, "slow": slow, "slPct": sl_pct, "tpPct": tp_pct}
    else:
        rsi_period = int(params_raw.get("rsiPeriod") or params_raw.get("period") or 14)
        oversold = float(params_raw.get("oversold") or 30)
        exit_above = float(params_raw.get("exitAbove") or 55)
        if rsi_period < 2 or rsi_period > 60:
            raise https_fn.HttpsError(https_fn.FunctionsErrorCode.INVALID_ARGUMENT, "Invalid RSI period.")
        if oversold <= 0 or oversold >= 50:
            raise https_fn.HttpsError(https_fn.FunctionsErrorCode.INVALID_ARGUMENT, "Invalid RSI oversold threshold.")
        if exit_above <= 50 or exit_above >= 95:
            raise https_fn.HttpsError(https_fn.FunctionsErrorCode.INVALID_ARGUMENT, "Invalid RSI exit threshold.")
        params = {"rsiPeriod": rsi_period, "oversold": oversold, "exitAbove": exit_above, "slPct": sl_pct, "tpPct": tp_pct}

    is_pro = token.get("email") == ADMIN_EMAIL or _is_paid_user(req.auth.uid)
    free_limit = _remote_config_int("backtesting_free_daily_limit", 1, context=context)
    pro_limit = _remote_config_int("backtesting_pro_daily_limit", 25, context=context)
    _enforce_daily_usage(
        req.auth.uid,
        "backtestRuns",
        pro_limit if is_pro else free_limit,
        "Daily usage limit reached. Upgrade your plan to run more backtests.",
    )

    from datetime import timedelta  # local import to reduce cold start
    import io

    try:
        import pandas as pd  # type: ignore
        import yfinance as yf  # type: ignore
        from backtesting import Backtest, Strategy  # type: ignore
        from backtesting.lib import crossover  # type: ignore
        from backtesting.test import SMA  # type: ignore
        import matplotlib

        matplotlib.use("Agg")
        import matplotlib.pyplot as plt  # type: ignore
    except Exception as exc:
        _raise_structured_error(
            https_fn.FunctionsErrorCode.FAILED_PRECONDITION,
            "backtest_dependency_error",
            "Backtest dependencies are unavailable.",
            {"raw": str(exc)},
        )

    def calc_rsi(series, period: int = 14):
        s = pd.Series(series)
        delta = s.diff()
        up = delta.clip(lower=0).rolling(period).mean()
        down = (-delta.clip(upper=0)).rolling(period).mean()
        rs = up / down
        return 100 - (100 / (1 + rs))

    if strategy == "rsi_reversion":

        class RsiReversion(Strategy):  # type: ignore
            rsi_period = int(params["rsiPeriod"])
            oversold = float(params["oversold"])
            exit_above = float(params["exitAbove"])
            sl_pct = float(params.get("slPct") or 0.0)
            tp_pct = float(params.get("tpPct") or 0.0)

            def _risk_args(self):
                last_close = float(self.data.Close[-1])
                out: dict[str, float] = {}
                if self.sl_pct > 0:
                    out["sl"] = last_close * (1 - self.sl_pct)
                if self.tp_pct > 0:
                    out["tp"] = last_close * (1 + self.tp_pct)
                return out

            def init(self):
                self.rsi = self.I(calc_rsi, self.data.Close, self.rsi_period)

            def next(self):
                if not self.position and self.rsi[-1] < self.oversold:
                    self.buy(**self._risk_args())
                elif self.position and self.rsi[-1] > self.exit_above:
                    self.position.close()

        strategy_cls = RsiReversion
        strategy_name = "RSI reversion"
    else:

        class SmaCross(Strategy):  # type: ignore
            fast = int(params["fast"])
            slow = int(params["slow"])
            sl_pct = float(params.get("slPct") or 0.0)
            tp_pct = float(params.get("tpPct") or 0.0)

            def _risk_args(self):
                last_close = float(self.data.Close[-1])
                out: dict[str, float] = {}
                if self.sl_pct > 0:
                    out["sl"] = last_close * (1 - self.sl_pct)
                if self.tp_pct > 0:
                    out["tp"] = last_close * (1 + self.tp_pct)
                return out

            def init(self):
                price = self.data.Close
                self.sma_fast = self.I(SMA, price, self.fast)
                self.sma_slow = self.I(SMA, price, self.slow)

            def next(self):
                if crossover(self.sma_fast, self.sma_slow):
                    self.position.close()
                    self.buy(**self._risk_args())
                elif crossover(self.sma_slow, self.sma_fast):
                    self.position.close()

        strategy_cls = SmaCross
        strategy_name = "SMA crossover"

    end = datetime.now(timezone.utc)
    start = end - timedelta(days=lookback_days)
    try:
        df = yf.download(
            ticker,
            start=start.strftime("%Y-%m-%d"),
            end=end.strftime("%Y-%m-%d"),
            interval=interval,
            progress=False,
            auto_adjust=False,
            threads=False,
        )
    except Exception as exc:
        _raise_structured_error(
            https_fn.FunctionsErrorCode.NOT_FOUND,
            "ticker_not_found",
            "Unable to load ticker history for backtest.",
            {"ticker": ticker, "raw": str(exc)},
        )
    if df is None or getattr(df, "empty", True):
        raise https_fn.HttpsError(https_fn.FunctionsErrorCode.FAILED_PRECONDITION, "No price history returned.")
    if isinstance(df.columns, pd.MultiIndex):
        df.columns = df.columns.get_level_values(0)
    df = df.dropna()
    df.index = pd.to_datetime(df.index, errors="coerce")
    if getattr(df.index, "tz", None) is not None:
        df.index = df.index.tz_localize(None)

    for col in ["Open", "High", "Low", "Close"]:
        if col not in df.columns:
            raise https_fn.HttpsError(https_fn.FunctionsErrorCode.FAILED_PRECONDITION, f"Missing column: {col}")

    try:
        bt = Backtest(df, strategy_cls, cash=cash, commission=commission)
        stats = bt.run()
        equity = None
        try:
            equity = stats.get("_equity_curve") if hasattr(stats, "get") else None
        except Exception:
            equity = None

        fig = plt.figure(figsize=(10.5, 4.4), dpi=160)
        ax = fig.add_subplot(111)
        ax.set_title(f"{ticker} backtest: {strategy_name}")
        ax.set_xlabel("Date")
        ax.set_ylabel("Equity (USD)")
        ax.grid(True, alpha=0.22)
        if equity is not None and hasattr(equity, "__getitem__") and "Equity" in equity:
            ax.plot(equity.index, equity["Equity"], linewidth=1.8)
        else:
            # Fall back to close price if equity curve is unavailable.
            ax.plot(df.index, df["Close"], linewidth=1.6)
            ax.set_ylabel("Close (USD)")
        fig.tight_layout()
        buf = io.BytesIO()
        fig.savefig(buf, format="png", bbox_inches="tight")
        plt.close(fig)
        png_bytes = buf.getvalue()
    except https_fn.HttpsError:
        raise
    except Exception as exc:
        _raise_structured_error(
            https_fn.FunctionsErrorCode.INTERNAL,
            "backtest_execution_error",
            "Backtest execution failed.",
            {"ticker": ticker, "strategy": strategy, "raw": str(exc)},
        )

    export_payload = {
        "ticker": ticker,
        "interval": interval,
        "lookbackDays": lookback_days,
        "strategy": strategy,
        "params": params,
        "cash": cash,
        "commission": commission,
    }
    export_sources = _build_backtest_export_sources(export_payload)
    code = export_sources.get("python", {}).get("content") or ""

    doc_ref = db.collection("backtests").document()
    backtest_id = doc_ref.id
    image_path = f"backtests/{req.auth.uid}/{backtest_id}.png"

    try:
        bucket = admin_storage.bucket(STORAGE_BUCKET)
        bucket.blob(image_path).upload_from_string(png_bytes, content_type="image/png")
    except Exception:
        raise https_fn.HttpsError(https_fn.FunctionsErrorCode.INTERNAL, "Unable to store backtest chart.")

    metrics = {
        "ReturnPct": float(stats.get("Return [%]", 0.0)) if hasattr(stats, "get") else 0.0,
        "Sharpe": float(stats.get("Sharpe Ratio", 0.0)) if hasattr(stats, "get") else 0.0,
        "MaxDrawdownPct": float(stats.get("Max. Drawdown [%]", 0.0)) if hasattr(stats, "get") else 0.0,
        "Trades": int(stats.get("# Trades", 0)) if hasattr(stats, "get") else 0,
        "WinRatePct": float(stats.get("Win Rate [%]", 0.0)) if hasattr(stats, "get") else 0.0,
    }

    title = str(data.get("title") or f"{ticker} · {strategy_name}").strip()
    if len(title) > 180:
        title = title[:180]

    doc_ref.set(
        {
            "userId": req.auth.uid,
            "title": title,
            "ticker": ticker,
            "interval": interval,
            "lookbackDays": lookback_days,
            "strategy": strategy,
            "params": params,
            "cash": cash,
            "commission": commission,
            "status": "completed",
            "metrics": metrics,
            "imagePath": image_path,
            "code": code,
            "exportSources": export_sources,
            "createdAt": firestore.SERVER_TIMESTAMP,
            "updatedAt": firestore.SERVER_TIMESTAMP,
            "meta": meta,
        }
    )

    _audit_event(req.auth.uid, token.get("email"), "backtest_completed", {"backtestId": backtest_id, "ticker": ticker, "strategy": strategy})
    return {
        "backtestId": backtest_id,
        "metrics": metrics,
        "imagePath": image_path,
        "title": title,
        "exportFormats": sorted(BACKTEST_SOURCE_FORMATS),
    }


@https_fn.on_call()
def rename_backtest(req: https_fn.CallableRequest) -> dict[str, Any]:
    token = _require_auth(req)
    data = req.data or {}
    backtest_id = str(data.get("backtestId") or data.get("id") or "").strip()
    title = str(data.get("title") or "").strip()
    if not backtest_id:
        raise https_fn.HttpsError(https_fn.FunctionsErrorCode.INVALID_ARGUMENT, "Backtest ID is required.")
    if not title:
        raise https_fn.HttpsError(https_fn.FunctionsErrorCode.INVALID_ARGUMENT, "Title is required.")
    if len(title) > 180:
        raise https_fn.HttpsError(https_fn.FunctionsErrorCode.INVALID_ARGUMENT, "Title is too long.")

    ref = db.collection("backtests").document(backtest_id)
    snap = ref.get()
    if not snap.exists:
        raise https_fn.HttpsError(https_fn.FunctionsErrorCode.NOT_FOUND, "Backtest not found.")
    doc = snap.to_dict() or {}
    owner_id = str(doc.get("userId") or "").strip()
    if token.get("email") != ADMIN_EMAIL and owner_id != req.auth.uid:
        raise https_fn.HttpsError(https_fn.FunctionsErrorCode.PERMISSION_DENIED, "Access denied.")

    ref.set({"title": title, "updatedAt": firestore.SERVER_TIMESTAMP}, merge=True)
    _audit_event(req.auth.uid, token.get("email"), "backtest_renamed", {"backtestId": backtest_id})
    return {"updated": True, "backtestId": backtest_id, "title": title}


@https_fn.on_call()
def delete_backtest(req: https_fn.CallableRequest) -> dict[str, Any]:
    token = _require_auth(req)
    data = req.data or {}
    backtest_id = str(data.get("backtestId") or data.get("id") or "").strip()
    if not backtest_id:
        raise https_fn.HttpsError(https_fn.FunctionsErrorCode.INVALID_ARGUMENT, "Backtest ID is required.")

    ref = db.collection("backtests").document(backtest_id)
    snap = ref.get()
    if not snap.exists:
        raise https_fn.HttpsError(https_fn.FunctionsErrorCode.NOT_FOUND, "Backtest not found.")
    doc = snap.to_dict() or {}
    owner_id = str(doc.get("userId") or "").strip()
    if token.get("email") != ADMIN_EMAIL and owner_id != req.auth.uid:
        raise https_fn.HttpsError(https_fn.FunctionsErrorCode.PERMISSION_DENIED, "Access denied.")

    image_path = str(doc.get("imagePath") or "").strip()
    if image_path:
        try:
            bucket = admin_storage.bucket(STORAGE_BUCKET)
            bucket.blob(image_path).delete()
        except Exception:
            pass

    ref.delete()
    _audit_event(req.auth.uid, token.get("email"), "backtest_deleted", {"backtestId": backtest_id})
    return {"deleted": True, "backtestId": backtest_id}


@https_fn.on_call()
def create_share_link(req: https_fn.CallableRequest) -> dict[str, Any]:
    token = _require_auth(req)
    data = req.data or {}
    kind = str(data.get("kind") or "").strip().lower()
    source_id = str(data.get("id") or "").strip()
    if kind not in {"forecast", "screener", "upload"}:
        raise https_fn.HttpsError(https_fn.FunctionsErrorCode.INVALID_ARGUMENT, "Invalid share kind.")
    if not source_id:
        raise https_fn.HttpsError(https_fn.FunctionsErrorCode.INVALID_ARGUMENT, "Source ID is required.")

    collection = "forecast_requests" if kind == "forecast" else "screener_runs" if kind == "screener" else "prediction_uploads"
    snap = db.collection(collection).document(source_id).get()
    if not snap.exists:
        raise https_fn.HttpsError(https_fn.FunctionsErrorCode.NOT_FOUND, "Source document not found.")

    doc = snap.to_dict() or {}
    owner_id = str(doc.get("userId") or "").strip()
    if kind in {"forecast", "screener"}:
        if not owner_id:
            _require_admin(token)
        else:
            _require_workspace_access(owner_id, req.auth.uid, token)
    else:
        _require_admin(token)

    share_ref = db.collection("shares").document()
    share_doc = {
        "kind": kind,
        "sourceCollection": collection,
        "sourceId": source_id,
        "sourceUserId": owner_id,
        "createdByUid": req.auth.uid,
        "createdByEmail": token.get("email"),
        "title": doc.get("title") or doc.get("ticker") or "",
        "ticker": doc.get("ticker") or "",
        "createdAt": firestore.SERVER_TIMESTAMP,
    }
    share_ref.set(share_doc)

    path = "/forecasting" if kind == "forecast" else "/screener" if kind == "screener" else "/uploads"
    share_url = f"{PUBLIC_ORIGIN}{path}?share={share_ref.id}"
    _audit_event(req.auth.uid, token.get("email"), "share_created", {"shareId": share_ref.id, "kind": kind})
    return {"shareId": share_ref.id, "shareUrl": share_url, "kind": kind}


@https_fn.on_call()
def import_shared_item(req: https_fn.CallableRequest) -> dict[str, Any]:
    token = _require_auth(req)
    data = req.data or {}
    share_id = str(data.get("shareId") or "").strip()
    if not share_id:
        raise https_fn.HttpsError(https_fn.FunctionsErrorCode.INVALID_ARGUMENT, "Share ID is required.")

    share_snap = db.collection("shares").document(share_id).get()
    if not share_snap.exists:
        raise https_fn.HttpsError(https_fn.FunctionsErrorCode.NOT_FOUND, "Share link not found.")

    share_doc = share_snap.to_dict() or {}
    kind = str(share_doc.get("kind") or "").strip().lower()
    source_collection = str(share_doc.get("sourceCollection") or "").strip()
    source_id = str(share_doc.get("sourceId") or "").strip()
    if kind not in {"forecast", "screener", "upload"} or not source_collection or not source_id:
        raise https_fn.HttpsError(https_fn.FunctionsErrorCode.FAILED_PRECONDITION, "Share link is invalid.")

    destination_workspace_id = str(data.get("workspaceId") or req.auth.uid or "").strip() or req.auth.uid
    if destination_workspace_id != req.auth.uid:
        _require_workspace_editor(destination_workspace_id, req.auth.uid, token)

    import_key = f"{destination_workspace_id}_{share_id}"
    import_ref = db.collection("share_imports").document(import_key)
    existing = import_ref.get()
    if existing.exists:
        payload = existing.to_dict() or {}
        imported_id = str(payload.get("importedId") or "").strip()
        if imported_id:
            return {"kind": kind, "importedId": imported_id, "shareId": share_id}

    source_snap = db.collection(source_collection).document(source_id).get()
    if not source_snap.exists:
        raise https_fn.HttpsError(https_fn.FunctionsErrorCode.NOT_FOUND, "Source document no longer exists.")

    source = source_snap.to_dict() or {}
    imported_ref = db.collection(source_collection).document()
    imported_doc: dict[str, Any] = {
        "userId": destination_workspace_id,
        "userEmail": token.get("email"),
        "createdByUid": req.auth.uid,
        "createdByEmail": token.get("email"),
        "status": source.get("status") or "completed",
        "createdAt": firestore.SERVER_TIMESTAMP,
        "updatedAt": firestore.SERVER_TIMESTAMP,
        "sourceShareId": share_id,
        "sourceId": source_id,
        "sourceUserId": str(source.get("userId") or ""),
        "meta": source.get("meta") or {},
    }

    if kind == "forecast":
        for key in [
            "ticker",
            "interval",
            "horizon",
            "start",
            "quantiles",
            "service",
            "engine",
            "serviceMessage",
            "metrics",
            "forecastPreview",
            "forecastRows",
        ]:
            if key in source:
                imported_doc[key] = source.get(key)
        imported_doc["title"] = source.get("title") or f"{source.get('ticker') or ''} forecast".strip()
    elif kind == "screener":
        for key in ["market", "universe", "maxNames", "results", "notes"]:
            if key in source:
                imported_doc[key] = source.get(key)
        imported_doc["title"] = source.get("title") or "Screener run"
    else:
        title = str(source.get("title") or "predictions.csv")
        notes = str(source.get("notes") or "")
        file_path = str(source.get("filePath") or "").strip()
        new_path = ""
        if file_path:
            try:
                bucket = admin_storage.bucket(STORAGE_BUCKET)
                src_blob = bucket.blob(file_path)
                if src_blob.exists():
                    base = os.path.basename(file_path)
                    new_path = f"predictions/{destination_workspace_id}/{int(time.time()*1000)}_{base}"
                    bucket.copy_blob(src_blob, bucket, new_path)
            except Exception:
                new_path = ""

        imported_doc.update(
            {
                "title": title,
                "notes": notes,
                "ticker": str(source.get("ticker") or "").upper(),
                "status": "uploaded",
                "filePath": new_path or file_path,
                "fileUrl": "",
            }
        )

    imported_ref.set(imported_doc)
    import_ref.set(
        {
            "shareId": share_id,
            "kind": kind,
            "importedCollection": source_collection,
            "importedId": imported_ref.id,
            "workspaceId": destination_workspace_id,
            "createdByUid": req.auth.uid,
            "createdAt": firestore.SERVER_TIMESTAMP,
        }
    )

    _audit_event(req.auth.uid, token.get("email"), "share_imported", {"shareId": share_id, "kind": kind})
    return {"kind": kind, "importedId": imported_ref.id, "shareId": share_id}


@https_fn.on_call(memory=MemoryOption.MB_512, timeout_sec=120)
def get_ticker_history(req: https_fn.CallableRequest) -> dict[str, Any]:
    import pandas as pd  # type: ignore
    import yfinance as yf  # type: ignore

    data = req.data or {}
    ticker = str(data.get("ticker") or "").upper()
    if not ticker:
        raise https_fn.HttpsError(https_fn.FunctionsErrorCode.INVALID_ARGUMENT, "Ticker is required.")

    interval = str(data.get("interval") or "1d").strip().lower()
    if interval not in {"1d", "1h"}:
        raise https_fn.HttpsError(https_fn.FunctionsErrorCode.INVALID_ARGUMENT, "Interval must be 1d or 1h.")
    start = data.get("start")
    end = data.get("end")

    try:
        history = yf.download(
            ticker,
            start=start,
            end=end,
            interval=interval,
            progress=False,
            auto_adjust=False,
            threads=False,
        )
    except Exception as exc:
        _raise_structured_error(
            https_fn.FunctionsErrorCode.NOT_FOUND,
            "ticker_not_found",
            "Unable to load market history for ticker.",
            {"ticker": ticker, "raw": str(exc)},
        )
    if isinstance(history.columns, pd.MultiIndex):
        history.columns = history.columns.get_level_values(0)
    # Keep only the columns the client chart needs to reduce payload/memory.
    keep_cols = [col for col in ["Open", "High", "Low", "Close", "Adj Close", "Volume"] if col in history.columns]
    if keep_cols:
        history = history[keep_cols]
    history = history.dropna().reset_index()
    if history.empty:
        _raise_structured_error(
            https_fn.FunctionsErrorCode.NOT_FOUND,
            "ticker_not_found",
            "No history returned for ticker.",
            {"ticker": ticker},
        )

    date_col = "Datetime" if "Datetime" in history.columns else "Date"
    if date_col in history.columns:
        history[date_col] = history[date_col].astype(str)

    rows = history.to_dict(orient="records")
    return {"rows": _serialize_for_firestore(rows)}


@https_fn.on_call(memory=MemoryOption.MB_512, timeout_sec=180)
def download_price_csv(req: https_fn.CallableRequest) -> dict[str, Any]:
    import pandas as pd  # type: ignore
    import yfinance as yf  # type: ignore

    _require_auth(req)
    data = req.data or {}
    ticker = str(data.get("ticker") or "").upper().strip()
    interval = str(data.get("interval") or "1d").strip().lower()
    if not ticker:
        raise https_fn.HttpsError(https_fn.FunctionsErrorCode.INVALID_ARGUMENT, "Ticker is required.")
    if interval not in {"1d", "1h"}:
        raise https_fn.HttpsError(https_fn.FunctionsErrorCode.INVALID_ARGUMENT, "Interval must be 1d or 1h.")

    end_raw = str(data.get("end") or "").strip()
    start_raw = str(data.get("start") or "").strip()

    def _parse_iso_date(raw: str, fallback: date) -> date:
        raw = (raw or "").strip()
        if not raw:
            return fallback
        try:
            return datetime.fromisoformat(raw[:10]).date()
        except Exception:
            return fallback

    end_date = _parse_iso_date(end_raw, date.today())
    if interval == "1h":
        min_start = end_date - timedelta(days=729)
        start_date = _parse_iso_date(start_raw, min_start)
        if start_date < min_start:
            start_date = min_start
    else:
        start_date = _parse_iso_date(start_raw, date(1900, 1, 1))

    # yfinance treats end as exclusive; bump one day so user-selected end is included.
    end_exclusive = end_date + timedelta(days=1)

    history = yf.download(
        ticker,
        start=start_date.isoformat(),
        end=end_exclusive.isoformat(),
        interval=interval,
        progress=False,
    )
    if isinstance(history.columns, pd.MultiIndex):
        history.columns = history.columns.get_level_values(0)
    history = history.dropna()
    if history.empty or "Close" not in history.columns:
        raise https_fn.HttpsError(https_fn.FunctionsErrorCode.NOT_FOUND, "No data returned for the requested range.")

    out_df = history[["Close"]].rename(columns={"Close": "Price"}).copy()
    out_df.reset_index(inplace=True)
    date_col = "Datetime" if "Datetime" in out_df.columns else "Date"
    if date_col != "Date" and date_col in out_df.columns:
        out_df.rename(columns={date_col: "Date"}, inplace=True)
    if "Date" in out_df.columns:
        out_df["Date"] = out_df["Date"].astype(str)
    out_df.insert(0, "Item_Id", ticker.lower())

    buffer = StringIO()
    out_df.to_csv(buffer, index=False)
    buffer.seek(0)
    csv_text = _force_remove_second_line(buffer.getvalue())

    filename = f"{ticker}_{start_date.isoformat()}_{end_date.isoformat()}_{interval}.csv"
    return {
        "ticker": ticker,
        "interval": interval,
        "start": start_date.isoformat(),
        "end": end_date.isoformat(),
        "rowCount": int(len(out_df)),
        "filename": filename,
        "csv": csv_text,
    }


@https_fn.on_call()
def get_unsplash_gallery(req: https_fn.CallableRequest) -> dict[str, Any]:
    data = req.data or {}
    query = str(data.get("query") or data.get("q") or "stock market, trading desk").strip()
    count = max(1, min(int(data.get("count") or data.get("limit") or 1), 8))

    photos, warning = _fetch_unsplash_random_photos(query, count=count)
    return {
        "query": query,
        "count": count,
        "photos": _serialize_for_firestore(photos),
        "warning": str(warning or "").strip(),
    }


@https_fn.on_call()
def get_trending_tickers(req: https_fn.CallableRequest) -> dict[str, Any]:
    data = req.data or {}
    force = bool(data.get("force"))
    now = time.time()

    cache_ref = db.collection("cache").document("trending_us")
    ttl_seconds = 24 * 60 * 60

    if not force:
        try:
            cached_snap = cache_ref.get()
            if cached_snap.exists:
                cached = cached_snap.to_dict() or {}
                updated_epoch = float(cached.get("updatedAtEpoch") or 0.0)
                tickers_cached = cached.get("tickers") or []
                items_cached = cached.get("items") or []
                if updated_epoch and (now - updated_epoch) < ttl_seconds and (tickers_cached or items_cached):
                    return {
                        "tickers": tickers_cached,
                        "items": items_cached,
                        "cached": True,
                        "updatedAtEpoch": updated_epoch,
                    }
        except Exception:
            # Cache is best-effort; fall back to live fetch.
            pass

    tickers: list[str] = []
    try:
        response = requests.get(TRENDING_URL, headers=_yahoo_headers(), timeout=10)
        response.raise_for_status()
        data = response.json()
        quotes = data.get("finance", {}).get("result", [{}])[0].get("quotes", [])
        tickers = [item.get("symbol") for item in quotes if item.get("symbol")]
    except Exception:
        # Keep the UI functional even when Yahoo throttles.
        tickers = ["AAPL", "MSFT", "NVDA", "AMZN", "GOOGL", "TSLA", "META"]

    tickers = [str(t).upper().strip() for t in (tickers or []) if str(t).strip()]
    tickers = list(dict.fromkeys(tickers))

    snapshots = _trending_snapshots(tickers, max_rows=18)
    items: list[dict[str, Any]] = []
    for symbol in tickers[:18]:
        snap = snapshots.get(symbol) or {}
        items.append({"symbol": symbol, **snap})

    payload = {
        "tickers": tickers,
        "items": _serialize_for_firestore(items),
        "cached": False,
        "updatedAtEpoch": int(now),
    }

    try:
        cache_ref.set(
            {
                "tickers": tickers,
                "items": payload["items"],
                "updatedAtEpoch": payload["updatedAtEpoch"],
                "updatedAt": firestore.SERVER_TIMESTAMP,
            },
            merge=True,
        )
    except Exception:
        pass

    return payload


@https_fn.on_call(memory=MemoryOption.GB_1, timeout_sec=180)
def get_ticker_intel(req: https_fn.CallableRequest) -> dict[str, Any]:
    import numpy as np  # type: ignore
    import pandas as pd  # type: ignore
    import yfinance as yf  # type: ignore

    data = req.data or {}
    ticker = str(data.get("ticker") or "").upper().strip()
    if not ticker:
        raise https_fn.HttpsError(https_fn.FunctionsErrorCode.INVALID_ARGUMENT, "Ticker is required.")

    def _as_float(value: Any) -> float | None:
        try:
            num = float(value)
            return num if math.isfinite(num) else None
        except Exception:
            return None

    def _statement_value(frame: Any, labels: list[str]) -> float | None:
        if frame is None or getattr(frame, "empty", True):
            return None
        try:
            for label in labels:
                if label in frame.index:
                    series = frame.loc[label]
                    if hasattr(series, "dropna"):
                        series = series.dropna()
                    if getattr(series, "empty", True):
                        continue
                    return _as_float(series.iloc[0])
        except Exception:
            return None
        return None

    def _score_by_range(value: float | None, low: float, high: float, invert: bool = False) -> float | None:
        if value is None or high <= low:
            return None
        clamped = max(low, min(high, value))
        pct = (clamped - low) / (high - low)
        if invert:
            pct = 1 - pct
        return round(max(0.0, min(1.0, pct)) * 100, 1)

    ticker_obj = yf.Ticker(ticker)

    info: dict[str, Any] = {}
    try:
        info = ticker_obj.info or {}
        if not isinstance(info, dict):
            info = {}
    except Exception:
        info = {}

    summary = str(info.get("longBusinessSummary") or "").strip()
    if len(summary) > 900:
        summary = summary[:897].rstrip() + "..."

    profile = {
        "name": info.get("longName") or info.get("shortName") or ticker,
        "sector": info.get("sector") or "",
        "industry": info.get("industry") or "",
        "website": info.get("website") or "",
        "country": info.get("country") or "",
        "exchange": info.get("fullExchangeName") or info.get("exchange") or "",
        "currency": info.get("currency") or "",
        "summary": summary,
        "marketCap": info.get("marketCap"),
        "beta": info.get("beta"),
        "trailingPE": info.get("trailingPE"),
        "forwardPE": info.get("forwardPE"),
        "dividendYield": info.get("dividendYield"),
        "fiftyTwoWeekLow": info.get("fiftyTwoWeekLow"),
        "fiftyTwoWeekHigh": info.get("fiftyTwoWeekHigh"),
        "debtToEquity": info.get("debtToEquity"),
        "currentRatio": info.get("currentRatio"),
        "quickRatio": info.get("quickRatio"),
        "returnOnEquity": info.get("returnOnEquity"),
        "returnOnAssets": info.get("returnOnAssets"),
    }

    income_stmt = None
    balance_sheet = None
    cashflow = None
    try:
        income_stmt = ticker_obj.quarterly_income_stmt
    except Exception:
        income_stmt = None
    try:
        balance_sheet = ticker_obj.quarterly_balance_sheet
    except Exception:
        balance_sheet = None
    try:
        cashflow = ticker_obj.quarterly_cashflow
    except Exception:
        cashflow = None

    total_revenue = _statement_value(income_stmt, ["Total Revenue", "Revenue"])
    gross_profit = _statement_value(income_stmt, ["Gross Profit"])
    operating_income = _statement_value(income_stmt, ["Operating Income"])
    net_income = _statement_value(income_stmt, ["Net Income", "Net Income Common Stockholders"])
    free_cash_flow = _statement_value(cashflow, ["Free Cash Flow"])

    net_margin = (net_income / total_revenue) if net_income is not None and total_revenue and total_revenue > 0 else None
    roi = _as_float(info.get("returnOnEquity"))
    if roi is None:
        roi = _as_float(info.get("returnOnAssets"))
    debt_to_equity = _as_float(info.get("debtToEquity"))
    if debt_to_equity is not None and debt_to_equity > 10:
        debt_to_equity = debt_to_equity / 100.0
    current_ratio = _as_float(info.get("currentRatio"))
    beta = _as_float(info.get("beta"))

    dividend_yield = _as_float(info.get("dividendYield"))
    payout_ratio = _as_float(info.get("payoutRatio"))
    shares_outstanding = _as_float(info.get("sharesOutstanding"))
    float_shares = _as_float(info.get("floatShares"))
    target_12m = _as_float(info.get("targetMeanPrice"))
    liquidity_cash = _as_float(info.get("totalCash"))
    liquidity_debt = _as_float(info.get("totalDebt"))

    dividend_policy = "No regular dividend policy disclosed."
    if dividend_yield is not None and dividend_yield > 0:
        payout_text = f"{payout_ratio * 100:.1f}%" if payout_ratio is not None and payout_ratio > 0 else "n/a"
        dividend_policy = f"Dividend yield near {dividend_yield * 100:.2f}% with payout ratio {payout_text}."

    buyback_summary = "Share buyback signal not explicit in current data."
    if shares_outstanding and float_shares and shares_outstanding > 0 and float_shares > 0:
        ratio = float_shares / shares_outstanding
        if ratio < 0.9:
            buyback_summary = "Share count structure suggests periodic repurchases versus total outstanding base."
        elif ratio > 0.98:
            buyback_summary = "Limited buyback footprint detected in the current share structure."

    risk_mitigation = []
    if current_ratio is not None:
        risk_mitigation.append(f"Current ratio {current_ratio:.2f}")
    if debt_to_equity is not None:
        risk_mitigation.append(f"Debt-to-equity {debt_to_equity:.2f}")
    if liquidity_cash is not None and liquidity_cash > 0:
        risk_mitigation.append(f"Cash buffer ${liquidity_cash:,.0f}")
    risk_mitigation_text = ", ".join(risk_mitigation) if risk_mitigation else "Liquidity and hedging data is limited for this symbol."

    env_score = _as_float(info.get("environmentScore"))
    social_score = _as_float(info.get("socialScore"))
    governance_score = _as_float(info.get("governanceScore"))
    total_esg = _as_float(info.get("totalEsg"))
    if total_esg is None:
        esg_parts = [x for x in [env_score, social_score, governance_score] if x is not None]
        total_esg = float(np.mean(esg_parts)) if esg_parts else None

    heatmap = [
        {
            "label": "Liquidity",
            "value": current_ratio,
            "score": _score_by_range(current_ratio, 0.5, 3.0),
            "hint": "Current ratio",
        },
        {
            "label": "Leverage",
            "value": debt_to_equity,
            "score": _score_by_range(debt_to_equity, 0.0, 3.0, invert=True),
            "hint": "Debt / equity",
        },
        {
            "label": "Net margin",
            "value": None if net_margin is None else net_margin * 100,
            "score": _score_by_range(net_margin, -0.05, 0.30),
            "hint": "Profitability",
        },
        {
            "label": "ROI",
            "value": None if roi is None else roi * 100,
            "score": _score_by_range(roi, -0.05, 0.25),
            "hint": "Return on capital",
        },
        {
            "label": "Cash conversion",
            "value": (free_cash_flow / total_revenue) * 100
            if free_cash_flow is not None and total_revenue and total_revenue > 0
            else None,
            "score": _score_by_range(
                (free_cash_flow / total_revenue) if free_cash_flow is not None and total_revenue and total_revenue > 0 else None,
                -0.05,
                0.25,
            ),
            "hint": "FCF / revenue",
        },
        {
            "label": "Volatility control",
            "value": beta,
            "score": _score_by_range(beta, 0.6, 2.0, invert=True),
            "hint": "Beta stability",
        },
    ]

    sector_peer_map: dict[str, list[str]] = {
        "Technology": ["MSFT", "AAPL", "NVDA", "AMD", "ORCL", "CRM"],
        "Financial Services": ["JPM", "BAC", "GS", "MS", "C", "WFC"],
        "Healthcare": ["LLY", "UNH", "JNJ", "PFE", "MRK", "ABBV"],
        "Energy": ["XOM", "CVX", "COP", "SLB", "EOG"],
        "Consumer Defensive": ["KO", "PEP", "WMT", "COST", "PG"],
        "Consumer Cyclical": ["AMZN", "TSLA", "HD", "MCD", "NKE"],
        "Industrials": ["CAT", "DE", "BA", "HON", "GE"],
    }
    sector_key = str(info.get("sector") or "")
    peer_candidates = [sym for sym in sector_peer_map.get(sector_key, ["AAPL", "MSFT", "NVDA", "AMZN", "GOOGL"]) if sym != ticker]
    peer_symbols = peer_candidates[:3]
    peer_rows: list[dict[str, Any]] = []
    for peer in peer_symbols:
        peer_info: dict[str, Any] = {}
        try:
            peer_info = yf.Ticker(peer).info or {}
            if not isinstance(peer_info, dict):
                peer_info = {}
        except Exception:
            peer_info = {}

        peer_pe = _as_float(peer_info.get("trailingPE"))
        peer_dte = _as_float(peer_info.get("debtToEquity"))
        if peer_dte is not None and peer_dte > 10:
            peer_dte = peer_dte / 100.0

        peer_sharpe = None
        try:
            hist = yf.download(peer, period="6mo", interval="1d", progress=False, auto_adjust=False)
            if hist is not None and getattr(hist, "empty", True) is False and "Close" in hist.columns:
                close = hist["Close"].astype(float).dropna()
                if len(close) > 15:
                    ret = close.pct_change().dropna()
                    vol = float(ret.std())
                    if vol > 0:
                        peer_sharpe = float((ret.mean() / vol) * math.sqrt(252))
        except Exception:
            peer_sharpe = None

        peer_rows.append(
            {
                "ticker": peer,
                "pe": None if peer_pe is None else round(peer_pe, 2),
                "debtToEquity": None if peer_dte is None else round(peer_dte, 2),
                "sharpeRatio": None if peer_sharpe is None else round(peer_sharpe, 2),
            }
        )

    calendar_raw: dict[str, Any] = {}
    try:
        cal = ticker_obj.calendar or {}
        if isinstance(cal, dict):
            calendar_raw = cal
    except Exception:
        calendar_raw = {}

    events: list[dict[str, Any]] = []
    for key, value in (calendar_raw or {}).items():
        if value is None:
            continue
        events.append({"label": str(key), "value": _serialize_for_firestore(value)})

    analyst = {
        "recommendationKey": info.get("recommendationKey"),
        "recommendationMean": info.get("recommendationMean"),
        "analystOpinions": info.get("numberOfAnalystOpinions"),
        "targetLowPrice": info.get("targetLowPrice"),
        "targetMeanPrice": info.get("targetMeanPrice"),
        "targetHighPrice": info.get("targetHighPrice"),
    }

    recommendation_trend: list[dict[str, Any]] = []
    try:
        rec = ticker_obj.recommendations
        if rec is not None and getattr(rec, "empty", True) is False:
            rec_rows = rec.reset_index(drop=True).head(6)
            for _, row in rec_rows.iterrows():
                recommendation_trend.append(
                    {
                        "period": str(row.get("period") or ""),
                        "strongBuy": int(row.get("strongBuy") or 0),
                        "buy": int(row.get("buy") or 0),
                        "hold": int(row.get("hold") or 0),
                        "sell": int(row.get("sell") or 0),
                        "strongSell": int(row.get("strongSell") or 0),
                    }
                )
    except Exception:
        recommendation_trend = []

    executive_summary = {
        "ticker": ticker,
        "exchange": profile.get("exchange"),
        "sector": profile.get("sector"),
        "marketCap": profile.get("marketCap"),
        "priceTarget12m": target_12m,
    }
    fundamental_deep_dive = {
        "revenueMechanics": {
            "totalRevenue": total_revenue,
            "grossProfit": gross_profit,
            "operatingIncome": operating_income,
            "segmentBreakdown": "Segment-level disclosure varies by issuer; this view highlights top-line and operating mechanics from reported statements.",
        },
        "profitability": {
            "netMargin": net_margin,
            "roi": roi,
            "returnOnAssets": _as_float(info.get("returnOnAssets")),
        },
        "capitalAllocation": {
            "dividendPolicy": dividend_policy,
            "shareBuybacks": buyback_summary,
            "freeCashFlow": free_cash_flow,
        },
    }
    risk_and_esg = {
        "riskMitigation": risk_mitigation_text,
        "liquidity": {
            "totalCash": liquidity_cash,
            "totalDebt": liquidity_debt,
            "currentRatio": current_ratio,
            "quickRatio": _as_float(info.get("quickRatio")),
        },
        "esg": {
            "environmental": env_score,
            "social": social_score,
            "governance": governance_score,
            "overall": total_esg,
        },
    }

    return {
        "ticker": ticker,
        "profile": _serialize_for_firestore(profile),
        "events": _serialize_for_firestore(events),
        "analyst": _serialize_for_firestore(analyst),
        "recommendationTrend": _serialize_for_firestore(recommendation_trend),
        "executiveSummary": _serialize_for_firestore(executive_summary),
        "fundamentalDeepDive": _serialize_for_firestore(fundamental_deep_dive),
        "riskAndEsg": _serialize_for_firestore(risk_and_esg),
        "balanceSheetHeatmap": _serialize_for_firestore(heatmap),
        "peerComparison": _serialize_for_firestore(peer_rows),
    }


@https_fn.on_call()
def get_ticker_news(req: https_fn.CallableRequest) -> dict[str, Any]:
    data = req.data or {}
    ticker = str(data.get("ticker") or "").upper()
    if not ticker:
        raise https_fn.HttpsError(https_fn.FunctionsErrorCode.INVALID_ARGUMENT, "Ticker is required.")

    seen: set[str] = set()

    def _extract_thumbnail(item: dict[str, Any]) -> str:
        thumb = item.get("thumbnail")
        if isinstance(thumb, dict):
            resolutions = thumb.get("resolutions")
            if isinstance(resolutions, list) and resolutions:
                # Prefer the largest image available.
                best = None
                for res in resolutions:
                    if not isinstance(res, dict):
                        continue
                    url = res.get("url")
                    width = res.get("width") or 0
                    if not url:
                        continue
                    if best is None or int(width or 0) > int(best.get("width") or 0):
                        best = {"url": str(url), "width": int(width or 0)}
                if best and best.get("url"):
                    return str(best["url"])
        return ""

    def _append_item(item: dict[str, Any]) -> None:
        title = str(item.get("title") or "").strip()
        link = str(item.get("link") or item.get("url") or "").strip()
        if not title:
            return
        key = (link or title).strip().lower()
        if key in seen:
            return
        seen.add(key)
        publisher = str(item.get("publisher") or item.get("source") or "").strip()
        published = item.get("publishedAt") or item.get("providerPublishTime")
        summary = str(item.get("summary") or item.get("description") or "").strip()
        news_items.append(
            {
                "title": title,
                "publisher": publisher,
                "link": link,
                "publishedAt": published,
                "summary": summary,
                "thumbnailUrl": _extract_thumbnail(item),
            }
        )

    news_items: list[dict[str, Any]] = []
    # 1) Prefer Yahoo search news (more reliable than yfinance .news)
    try:
        resp = requests.get(
            YAHOO_SEARCH_URL,
            headers=_yahoo_headers(),
            params={"q": ticker, "newsCount": 12, "quotesCount": 0, "listsCount": 0},
            timeout=10,
        )
        if resp.status_code < 400:
            payload = resp.json() if resp.text else {}
            for item in (payload.get("news") or [])[:12]:
                if isinstance(item, dict):
                    _append_item(item)
    except Exception:
        pass

    # 2) Fallback: yfinance news
    if not news_items:
        try:
            import yfinance as yf  # type: ignore

            ticker_obj = yf.Ticker(ticker)
            for item in (ticker_obj.news or [])[:10]:
                if not isinstance(item, dict):
                    continue
                _append_item(item)
        except Exception:
            pass

    # 3) Optional: RSS fallback (lightweight parse)
    if not news_items:
        try:
            import xml.etree.ElementTree as ET

            rss_url = f"https://feeds.finance.yahoo.com/rss/2.0/headline?s={ticker}&region=US&lang=en-US"
            rss = requests.get(rss_url, headers=_yahoo_headers(), timeout=10)
            if rss.status_code < 400 and rss.text:
                root = ET.fromstring(rss.text)
                for item in root.findall(".//item")[:10]:
                    title = (item.findtext("title") or "").strip()
                    link = (item.findtext("link") or "").strip()
                    pub_date = (item.findtext("pubDate") or "").strip()
                    if title:
                        key = (link or title).strip().lower()
                        if key in seen:
                            continue
                        seen.add(key)
                        news_items.append(
                            {
                                "title": title,
                                "publisher": "Yahoo Finance",
                                "link": link,
                                "publishedAt": pub_date,
                                "summary": "",
                                "thumbnailUrl": "",
                            }
                        )
        except Exception:
            pass

    return {"news": news_items}


@https_fn.on_call()
def get_ticker_x_trends(req: https_fn.CallableRequest) -> dict[str, Any]:
    data = req.data or {}
    ticker = str(data.get("ticker") or "").upper().strip()
    if not ticker:
        raise https_fn.HttpsError(https_fn.FunctionsErrorCode.INVALID_ARGUMENT, "Ticker is required.")
    query = f"${ticker} {ticker} stock"
    posts, warning = _fetch_x_social_posts(query, limit=8)
    stories, stories_warning = _fetch_x_news_stories(query, limit=6)
    combined_warning = " ".join([part for part in [warning, stories_warning] if str(part or "").strip()]).strip()
    return {
        "ticker": ticker,
        "posts": _serialize_for_firestore(posts),
        "stories": _serialize_for_firestore(stories),
        "warning": combined_warning,
    }


@https_fn.on_call()
def get_corporate_events_calendar(req: https_fn.CallableRequest) -> dict[str, Any]:
    data = req.data or {}
    ticker = _normalize_symbol_token(data.get("ticker"))
    tickers_raw = data.get("tickers") if isinstance(data.get("tickers"), list) else []
    tickers = [_normalize_symbol_token(item) for item in tickers_raw]
    tickers = [item for item in tickers if item]
    if ticker:
        tickers = [ticker] + tickers
    tickers = list(dict.fromkeys(tickers))

    if not tickers:
        tickers = ["AAPL", "MSFT", "NVDA", "AMZN", "GOOGL"]

    today = date.today()
    try:
        start_date = datetime.fromisoformat(str(data.get("startDate") or data.get("start") or today.isoformat())[:10]).date()
    except Exception:
        start_date = today
    try:
        default_end = today + timedelta(days=90)
        end_date = datetime.fromisoformat(str(data.get("endDate") or data.get("end") or default_end.isoformat())[:10]).date()
    except Exception:
        end_date = today + timedelta(days=90)

    if end_date < start_date:
        end_date = start_date + timedelta(days=30)
    if (end_date - start_date).days > 370:
        end_date = start_date + timedelta(days=370)

    country_code = _normalize_country_code(data.get("country"))
    use_massive = country_code == "US"
    if use_massive and tickers:
        # If all provided tickers look non-US, force Yahoo fallback.
        us_check = [_is_us_equity_symbol(item) for item in tickers[:10]]
        if us_check and not any(us_check):
            use_massive = False

    warnings: list[str] = []
    events: list[dict[str, Any]] = []
    source = "yahoo_finance"

    if use_massive:
        massive_events, massive_warning = _fetch_massive_corporate_events(
            tickers=tickers,
            start_date=start_date,
            end_date=end_date,
            limit=int(data.get("limit") or 250),
            event_types=data.get("eventTypes") if isinstance(data.get("eventTypes"), list) else None,
            statuses=data.get("statuses") if isinstance(data.get("statuses"), list) else None,
        )
        if massive_events:
            events = massive_events
            source = "massive_tmx"
        elif massive_warning:
            warnings.append(massive_warning)

    if not events:
        source = "yahoo_finance"
        for symbol in tickers[:30]:
            try:
                events.extend(
                    _fetch_yahoo_corporate_events_for_ticker(
                        symbol,
                        start_date=start_date,
                        end_date=end_date,
                    )
                )
            except Exception:
                continue

    dedup: dict[str, dict[str, Any]] = {}
    for row in events:
        if not isinstance(row, dict):
            continue
        event_id = str(row.get("id") or "").strip()
        if not event_id:
            event_id = f"{row.get('ticker')}_{row.get('date')}_{row.get('type')}".lower()
            row["id"] = event_id
        dedup[event_id] = row

    sorted_events = sorted(
        list(dedup.values()),
        key=lambda item: (str(item.get("date") or ""), str(item.get("ticker") or ""), str(item.get("name") or "")),
    )

    return {
        "tickers": tickers,
        "country": country_code,
        "startDate": start_date.isoformat(),
        "endDate": end_date.isoformat(),
        "source": source,
        "fallbackUsed": source != "massive_tmx",
        "warnings": warnings,
        "events": _serialize_for_firestore(sorted_events[: max(1, min(int(data.get("limit") or 250), 500))]),
    }


@https_fn.on_call()
def query_ticker_insight(req: https_fn.CallableRequest) -> dict[str, Any]:
    import yfinance as yf  # type: ignore

    data = req.data or {}
    meta = data.get("meta") if isinstance(data.get("meta"), dict) else {}
    context = _remote_config_context(req, None, meta if isinstance(meta, dict) else None)
    ticker = _normalize_symbol_token(data.get("ticker"))
    question = str(data.get("question") or data.get("query") or "").strip()
    language = str(data.get("language") or "en").strip().lower()[:8]
    if language in {"", "auto"}:
        language = "en"
    if not re.match(r"^[a-z]{2}(?:-[a-z]{2})?$", language):
        language = "en"
    language_label = {
        "en": "English",
        "es": "Spanish",
        "fr": "French",
        "de": "German",
        "ar": "Arabic",
        "bn": "Bengali",
    }.get(language, "English")
    technical_raw = data.get("technicalContext") if isinstance(data.get("technicalContext"), dict) else {}
    technical_context: dict[str, Any] = {}
    if technical_raw:
        lookback_days = technical_raw.get("lookbackDays")
        interval = str(technical_raw.get("interval") or "").strip().lower()
        if interval not in {"1d", "1h"}:
            interval = "1d"
        latest_rows = []
        for row in (technical_raw.get("latest") if isinstance(technical_raw.get("latest"), list) else [])[:24]:
            if not isinstance(row, dict):
                continue
            name = str(row.get("name") or "").strip()[:40]
            if not name:
                continue
            value = _safe_float(row.get("value"))
            if value is None:
                continue
            latest_rows.append({"name": name, "value": round(float(value), 6)})

        trend_rows = []
        for row in (technical_raw.get("trend") if isinstance(technical_raw.get("trend"), list) else [])[:16]:
            if not isinstance(row, dict):
                continue
            name = str(row.get("name") or "").strip()[:40]
            if not name:
                continue
            trend_rows.append(
                {
                    "name": name,
                    "direction": str(row.get("direction") or "").strip()[:16],
                    "pctChange": _safe_float(row.get("pctChange")),
                    "delta": _safe_float(row.get("delta")),
                }
            )

        heuristics = []
        for line in (technical_raw.get("heuristics") if isinstance(technical_raw.get("heuristics"), list) else [])[:12]:
            text = str(line or "").strip()
            if text:
                heuristics.append(text[:220])

        lookback_normalized = None
        try:
            lookback_candidate = int(float(lookback_days))
            if lookback_candidate > 0:
                lookback_normalized = lookback_candidate
        except Exception:
            lookback_normalized = None

        technical_context = {
            "lookbackDays": lookback_normalized,
            "interval": interval,
            "latest": latest_rows,
            "trend": trend_rows,
            "heuristics": heuristics,
        }
    if not ticker:
        raise https_fn.HttpsError(https_fn.FunctionsErrorCode.INVALID_ARGUMENT, "Ticker is required.")
    if not question:
        raise https_fn.HttpsError(https_fn.FunctionsErrorCode.INVALID_ARGUMENT, "Question is required.")

    tk = yf.Ticker(ticker)
    info: dict[str, Any] = {}
    try:
        raw_info = tk.info or {}
        info = raw_info if isinstance(raw_info, dict) else {}
    except Exception:
        info = {}

    history_summary: dict[str, Any] = {}
    try:
        hist = tk.history(period="6mo", interval="1d")
        if hist is not None and getattr(hist, "empty", True) is False and "Close" in hist.columns:
            close = hist["Close"].astype(float).dropna()
            if len(close) >= 2:
                last = float(close.iloc[-1])
                start = float(close.iloc[0])
                one_month = float(close.iloc[-22]) if len(close) >= 22 else None
                history_summary = {
                    "lastClose": round(last, 4),
                    "periodStartClose": round(start, 4),
                    "change6mPct": round(((last / start) - 1.0) * 100.0, 3) if start else None,
                    "change1mPct": round(((last / one_month) - 1.0) * 100.0, 3) if one_month else None,
                }
    except Exception:
        history_summary = {}

    headlines = _fetch_yahoo_news_query(ticker, limit=6)
    headline_compact = [
        {
            "title": str(item.get("title") or "").strip(),
            "publisher": str(item.get("publisher") or "").strip(),
            "publishedAt": item.get("publishedAt"),
            "link": str(item.get("link") or "").strip(),
        }
        for item in headlines[:6]
    ]

    context_payload = {
        "ticker": ticker,
        "country": str(info.get("country") or "").strip(),
        "exchange": str(info.get("fullExchangeName") or info.get("exchange") or "").strip(),
        "sector": str(info.get("sector") or "").strip(),
        "industry": str(info.get("industry") or "").strip(),
        "marketCap": info.get("marketCap"),
        "trailingPE": info.get("trailingPE"),
        "forwardPE": info.get("forwardPE"),
        "dividendYield": info.get("dividendYield"),
        "beta": info.get("beta"),
        "historySummary": history_summary,
        "headlines": headline_compact,
        "technicalContext": technical_context,
    }

    answer = ""
    model_used = "heuristic"
    provider = "local"
    allowed_models = _get_llm_allowed_models(context=context)
    allowed_set = set(allowed_models)
    requested_model = _normalize_ai_model_id(data.get("model") or "gpt-5")
    if not _is_supported_llm_model(requested_model):
        requested_model = "gpt-5"
    if allowed_set and requested_model not in allowed_set:
        requested_model = allowed_models[0] if allowed_models else "gpt-5"
    requested_provider = _model_provider_from_id(requested_model)

    if requested_provider == "openai" and OPENAI_API_KEY:
        model_used = requested_model
        provider = "openai"
        system_prompt = (
            "You are Quantura's market assistant. Answer using the supplied ticker context only. "
            "Be explicit when data is missing. Keep the answer concise, practical, and professional. "
            f"Write the final response in {language_label}."
        )
        user_prompt = (
            f"Ticker: {ticker}\n"
            f"Question: {question}\n"
            f"Preferred language: {language_label}\n"
            f"Context JSON:\n{json.dumps(context_payload, ensure_ascii=False)}\n\n"
            "Respond in plain text with: (1) direct answer, (2) key evidence bullets, (3) one risk caveat."
        )
        payload = {
            "model": model_used,
            "input": [
                {"role": "system", "content": [{"type": "input_text", "text": system_prompt}]},
                {"role": "user", "content": [{"type": "input_text", "text": user_prompt}]},
            ],
            "max_output_tokens": 900,
            "tools": [{"type": "web_search_preview"}],
        }
        try:
            response = requests.post(
                "https://api.openai.com/v1/responses",
                headers={"Authorization": f"Bearer {OPENAI_API_KEY}", "Content-Type": "application/json"},
                json=payload,
                timeout=35,
            )
            response.raise_for_status()
            answer = _extract_responses_output_text(response.json())
        except Exception as exc:
            provider = "fallback"
            model_used = "heuristic"
            answer = f"Unable to complete GPT-5 query right now. Fallback summary: {str(exc)[:160]}"
    elif requested_provider == "amazon_nova":
        system_prompt = (
            "You are Quantura's market assistant. Answer using the supplied ticker context only. "
            "Be explicit when data is missing. Keep the answer concise, practical, and professional. "
            f"Write the final response in {language_label}."
        )
        user_prompt = (
            f"Ticker: {ticker}\n"
            f"Question: {question}\n"
            f"Preferred language: {language_label}\n"
            f"Context JSON:\n{json.dumps(context_payload, ensure_ascii=False)}\n\n"
            "Respond in plain text with: (1) direct answer, (2) key evidence bullets, (3) one risk caveat."
        )
        answer = _invoke_amazon_nova_text(
            model_id=requested_model,
            system_prompt=system_prompt,
            user_prompt=user_prompt,
            max_tokens=900,
            temperature=0.2,
        )
        if answer:
            provider = "amazon_nova"
            model_used = requested_model

    if not answer:
        last_close = history_summary.get("lastClose")
        change_1m = history_summary.get("change1mPct")
        sector = str(context_payload.get("sector") or "Unknown sector")
        pe = context_payload.get("trailingPE")
        technical_brief = ""
        if technical_context:
            lookback_label = technical_context.get("lookbackDays") or "n/a"
            interval_label = technical_context.get("interval") or "1d"
            heuristics = technical_context.get("heuristics") if isinstance(technical_context.get("heuristics"), list) else []
            latest = technical_context.get("latest") if isinstance(technical_context.get("latest"), list) else []
            latest_text = ", ".join(
                [f"{str(item.get('name') or '')}:{item.get('value')}" for item in latest[:6] if isinstance(item, dict)]
            )
            heuristic_text = str(heuristics[0]) if heuristics else ""
            technical_brief = f"Technical snapshot ({lookback_label}d/{interval_label}) {latest_text}. {heuristic_text}".strip()
        answer = (
            f"{ticker} currently screens in {sector}. "
            f"Last close: {last_close if last_close is not None else 'n/a'}, "
            f"1M move: {change_1m if change_1m is not None else 'n/a'}%. "
            f"Trailing P/E: {pe if pe is not None else 'n/a'}. "
            f"Recent headlines were attached for context.{(' ' + technical_brief) if technical_brief else ''}"
        )

    return {
        "ticker": ticker,
        "question": question,
        "answer": answer.strip(),
        "provider": provider,
        "model": model_used,
        "context": _serialize_for_firestore(context_payload),
        "language": language,
        "usedYahooFallback": str(context_payload.get("country") or "").strip().lower() not in {"united states", "usa", "us"},
    }


@https_fn.on_call()
def get_market_headlines_feed(req: https_fn.CallableRequest) -> dict[str, Any]:
    data = req.data or {}
    country_code = _normalize_country_code(data.get("country"))
    limit = max(5, min(int(data.get("limit") or 14), 40))

    country_query = {
        "US": "US stock market top headlines today",
        "CA": "Canada stock market top headlines today",
        "GB": "UK stock market top headlines today",
        "DE": "Germany stock market top headlines today",
        "FR": "France stock market top headlines today",
        "JP": "Japan stock market top headlines today",
        "CN": "China stock market top headlines today",
        "IN": "India stock market top headlines today",
        "AU": "Australia stock market top headlines today",
        "BR": "Brazil stock market top headlines today",
    }.get(country_code, f"{country_code} stock market top headlines today")

    headlines = _fetch_yahoo_news_query(country_query, limit=limit)
    if not headlines:
        headlines = _fetch_yahoo_news_query("stock market top headlines today", limit=limit)

    warnings: list[str] = []
    x_posts, x_warning = _fetch_x_social_posts(country_query, limit=8)
    if x_warning:
        warnings.append(x_warning)
    x_stories, x_story_warning = _fetch_x_news_stories(country_query, limit=6)
    if x_story_warning:
        warnings.append(x_story_warning)
    if not x_posts and x_stories:
        x_posts = _x_news_stories_as_pseudo_posts(x_stories, limit=8)
        warnings.append("Showing X News stories because live X posts were unavailable.")
    reddit_query = f"{country_code} stock market investing reddit"
    reddit_posts = _fetch_reddit_social_posts(reddit_query, limit=8)
    facebook_query = f"{country_code} stock market investing"
    facebook_posts, facebook_warning = _fetch_meta_social_posts(facebook_query, platform="facebook", limit=6)
    if facebook_warning:
        warnings.append(facebook_warning)
    instagram_query = f"{country_code} stock market investing"
    instagram_posts, instagram_warning = _fetch_meta_social_posts(instagram_query, platform="instagram", limit=6)
    if instagram_warning:
        warnings.append(instagram_warning)

    return {
        "country": country_code,
        "query": country_query,
        "headlines": _serialize_for_firestore(headlines[:limit]),
        "social": {
            "x": _serialize_for_firestore(x_posts),
            "xStories": _serialize_for_firestore(x_stories),
            "reddit": _serialize_for_firestore(reddit_posts),
            "facebook": _serialize_for_firestore(facebook_posts),
            "instagram": _serialize_for_firestore(instagram_posts),
        },
        "warnings": warnings,
    }


@https_fn.on_call()
def get_options_chain(req: https_fn.CallableRequest) -> dict[str, Any]:
    import pandas as pd  # type: ignore
    import yfinance as yf  # type: ignore

    _require_auth(req)
    data = req.data or {}
    ticker = str(data.get("ticker") or "").upper().strip()
    expiration = str(data.get("expiration") or "").strip()
    limit = int(data.get("limit") or 24)
    limit = max(10, min(limit, 200))
    if not ticker:
        raise https_fn.HttpsError(https_fn.FunctionsErrorCode.INVALID_ARGUMENT, "Ticker is required.")

    ticker_obj = yf.Ticker(ticker)
    expirations = []
    try:
        expirations = list(ticker_obj.options or [])
    except Exception:
        expirations = []

    if not expirations:
        return {"ticker": ticker, "underlyingPrice": None, "expirations": [], "selectedExpiration": "", "calls": [], "puts": []}

    selected = expiration if expiration in expirations else expirations[0]

    underlying_price = None
    try:
        fast_info = getattr(ticker_obj, "fast_info", None) or {}
        underlying_price = _safe_float(fast_info.get("last_price") or fast_info.get("lastPrice"))
    except Exception:
        underlying_price = None
    if underlying_price is None:
        try:
            hist = ticker_obj.history(period="5d", interval="1d")
            if not hist.empty and "Close" in hist.columns:
                underlying_price = float(hist["Close"].dropna().iloc[-1])
        except Exception:
            underlying_price = None

    # Time to expiry in years (approx). If we cannot parse, probabilities will be null.
    T = None
    try:
        exp_dt = datetime.strptime(selected, "%Y-%m-%d").date()
        today = datetime.now(tz=timezone.utc).date()
        days = max((exp_dt - today).days, 0)
        T = days / 365.0
    except Exception:
        T = None

    normal = NormalDist()

    def _option_metrics(
        strike: float | None, iv: float | None, is_call: bool
    ) -> tuple[float, float, float, float] | None:
        if underlying_price is None or T is None or T <= 0:
            return None
        if strike is None or strike <= 0:
            return None
        if iv is None or iv <= 0:
            return None
        S = float(underlying_price)
        K = float(strike)
        sigma = float(iv)
        # Protect the probability math from extreme / bad IV values.
        if sigma > 5:
            return None
        denom = sigma * math.sqrt(T)
        if denom <= 0:
            return None
        d1 = (math.log(S / K) + (RISK_FREE_RATE + 0.5 * sigma * sigma) * T) / denom
        d2 = d1 - sigma * math.sqrt(T)
        p_call = normal.cdf(d2)
        prob_itm = p_call if is_call else normal.cdf(-d2)
        delta = normal.cdf(d1) if is_call else normal.cdf(d1) - 1.0
        return (prob_itm, delta, d1, d2)

    def _format_chain(df: pd.DataFrame, opt_type: str) -> list[dict[str, Any]]:
        if df is None or df.empty:
            return []
        frame = df.copy()
        if "strike" in frame.columns:
            frame["strike"] = pd.to_numeric(frame["strike"], errors="coerce")
        for col in ["openInterest", "volume"]:
            if col in frame.columns:
                frame[col] = pd.to_numeric(frame[col], errors="coerce").fillna(0)
        if "strike" in frame.columns:
            frame = frame.dropna(subset=["strike"])

        # Select a compact window around the underlying, then return results ordered by strike.
        if underlying_price is not None and "strike" in frame.columns:
            frame["distance"] = (frame["strike"] - float(underlying_price)).abs()
            sort_cols: list[str] = ["distance"]
            ascending: list[bool] = [True]
            if "openInterest" in frame.columns:
                sort_cols.append("openInterest")
                ascending.append(False)
            elif "volume" in frame.columns:
                sort_cols.append("volume")
                ascending.append(False)
            frame = frame.sort_values(by=sort_cols, ascending=ascending)
        elif "strike" in frame.columns:
            frame = frame.sort_values(by="strike", ascending=True)
        else:
            sort_col = "openInterest" if "openInterest" in frame.columns else ("volume" if "volume" in frame.columns else None)
            if sort_col:
                frame = frame.sort_values(by=sort_col, ascending=False)

        frame = frame.head(limit)
        if "strike" in frame.columns:
            frame = frame.sort_values(by="strike", ascending=True)

        items: list[dict[str, Any]] = []
        for _, row in frame.iterrows():
            strike = _safe_float(row.get("strike"))
            iv = _safe_float(row.get("impliedVolatility"))
            metrics = _option_metrics(strike, iv, is_call=(opt_type == "call"))
            p_itm = metrics[0] if metrics else None
            delta = metrics[1] if metrics else None
            bid = _safe_float(row.get("bid"))
            ask = _safe_float(row.get("ask"))
            mid = None
            if bid is not None and ask is not None and bid >= 0 and ask >= 0:
                mid = (bid + ask) / 2.0
            items.append(
                {
                    "contractSymbol": row.get("contractSymbol"),
                    "type": opt_type,
                    "strike": strike,
                    "lastPrice": _safe_float(row.get("lastPrice")),
                    "bid": bid,
                    "ask": ask,
                    "mid": None if mid is None else round(mid, 4),
                    "impliedVolatility": None if iv is None else round(iv, 4),
                    "delta": None if delta is None else round(delta, 4),
                    "volume": int(row.get("volume") or 0),
                    "openInterest": int(row.get("openInterest") or 0),
                    "inTheMoney": bool(row.get("inTheMoney")) if "inTheMoney" in row else None,
                    "probabilityITM": None if p_itm is None else round(p_itm * 100.0, 2),
                }
            )
        return items

    try:
        chain = ticker_obj.option_chain(selected)
        calls = _format_chain(chain.calls, "call")
        puts = _format_chain(chain.puts, "put")
    except Exception:
        calls = []
        puts = []

    return {
        "ticker": ticker,
        "underlyingPrice": underlying_price,
        "timeToExpiryYears": None if T is None else round(float(T), 6),
        "riskFreeRate": RISK_FREE_RATE,
        "expirations": expirations,
        "selectedExpiration": selected,
        "calls": _serialize_for_firestore(calls),
        "puts": _serialize_for_firestore(puts),
    }


@https_fn.on_call()
def get_technicals(req: https_fn.CallableRequest) -> dict[str, Any]:
    import pandas as pd  # type: ignore
    import yfinance as yf  # type: ignore
    from finta import TA  # type: ignore

    data = req.data or {}
    ticker = str(data.get("ticker") or "").upper()
    if not ticker:
        raise https_fn.HttpsError(https_fn.FunctionsErrorCode.INVALID_ARGUMENT, "Ticker is required.")

    interval = str(data.get("interval") or "1d")
    lookback = int(data.get("lookback") or 120)
    indicators = data.get("indicators") or ["RSI", "MACD"]
    include_series = bool(data.get("includeSeries"))
    max_points = int(data.get("maxPoints") or (240 if interval == "1h" else 260))
    max_points = max(30, min(max_points, 900))

    history = yf.download(ticker, period=f"{lookback}d", interval=interval, progress=False)
    if isinstance(history.columns, pd.MultiIndex):
        history.columns = history.columns.get_level_values(0)
    history = history.dropna()
    if history.empty:
        return {"latest": [], "series": {"dates": [], "items": []} if include_series else None}

    history = history.rename(
        columns={
            "Open": "open",
            "High": "high",
            "Low": "low",
            "Close": "close",
            "Volume": "volume",
        }
    )

    def _ta(name: str, frame: pd.DataFrame):
        func = getattr(TA, name, None)
        if not callable(func):
            return None
        return func(frame)

    latest_values = []
    series_items: list[dict[str, Any]] = []
    dates_index = history.index[-max_points:]
    dates_out = [pd.Timestamp(ts).isoformat() for ts in dates_index]
    indicator_map = {
        "RSI": lambda frame: _ta("RSI", frame),
        "MACD": lambda frame: TA.MACD(frame).get("MACD"),
        "SMA": lambda frame: TA.SMA(frame, 20),
        "EMA": lambda frame: TA.EMA(frame, 20),
        "BBANDS": lambda frame: _ta("BBANDS", frame),
        "ATR": lambda frame: TA.ATR(frame, 14),
        "ADX": lambda frame: _ta("ADX", frame),
        "CCI": lambda frame: _ta("CCI", frame),
        "MFI": lambda frame: _ta("MFI", frame),
        "OBV": lambda frame: _ta("OBV", frame),
        "ROC": lambda frame: _ta("ROC", frame),
        "STOCH": lambda frame: _ta("STOCH", frame),
        "WILLR": lambda frame: _ta("WILLIAMS", frame),
    }

    for indicator in indicators:
        func = indicator_map.get(indicator)
        if not func:
            continue
        try:
            series = func(history)
            if series is None or len(series) == 0:
                continue
            if indicator == "BBANDS" and isinstance(series, pd.DataFrame):
                for key, label in [("BB_UPPER", "BBANDS_UPPER"), ("BB_MIDDLE", "BBANDS_MIDDLE"), ("BB_LOWER", "BBANDS_LOWER")]:
                    band = series.get(key)
                    if band is None or len(band) == 0:
                        continue
                    value = band.dropna().iloc[-1]
                    latest_values.append({"name": label, "value": round(float(value), 4)})
                    if include_series:
                        aligned = band.reindex(dates_index)
                        series_items.append(
                            {
                                "name": label,
                                "values": [None if pd.isna(v) else round(float(v), 6) for v in aligned.tolist()],
                            }
                        )
                continue

            if isinstance(series, pd.DataFrame):
                series = series.iloc[:, 0]
            value = series.dropna().iloc[-1]
            latest_values.append({"name": indicator, "value": round(float(value), 4)})
            if include_series:
                aligned = series.reindex(dates_index)
                series_items.append(
                    {
                        "name": indicator,
                        "values": [None if pd.isna(v) else round(float(v), 6) for v in aligned.tolist()],
                    }
                )
        except Exception:
            continue

    out: dict[str, Any] = {"latest": latest_values}
    if include_series:
        out["series"] = {"dates": dates_out, "items": series_items}
    return out


@https_fn.on_call()
def queue_screener_run(req: https_fn.CallableRequest) -> dict[str, Any]:
    token = _require_auth(req)
    data = req.data or {}
    payload = {
        "userId": req.auth.uid,
        "userEmail": token.get("email"),
        "universe": data.get("universe"),
        "market": data.get("market"),
        "minCap": data.get("minCap"),
        "maxNames": data.get("maxNames"),
        "notes": data.get("notes"),
        "status": "queued",
        "createdAt": firestore.SERVER_TIMESTAMP,
        "updatedAt": firestore.SERVER_TIMESTAMP,
        "meta": data.get("meta") or {},
    }

    doc_ref = db.collection("screener_runs").document()
    doc_ref.set(payload)
    _audit_event(req.auth.uid, token.get("email"), "screener_queued", {"runId": doc_ref.id})
    return {"runId": doc_ref.id}


@https_fn.on_call(memory=MemoryOption.GB_1, timeout_sec=180)
def run_quick_screener(req: https_fn.CallableRequest) -> dict[str, Any]:
    try:
        import numpy as np  # type: ignore
        import pandas as pd  # type: ignore
        import yfinance as yf  # type: ignore
    except Exception as exc:
        _raise_structured_error(
            https_fn.FunctionsErrorCode.FAILED_PRECONDITION,
            "screener_dependency_error",
            "Screener dependencies are unavailable.",
            {"raw": str(exc)[:400]},
        )

    token = _require_auth(req)
    data = req.data or {}

    workspace_id = str(data.get("workspaceId") or req.auth.uid or "").strip()
    if not workspace_id:
        workspace_id = req.auth.uid

    # Allow writing into your own workspace, or an editor-level shared workspace.
    if workspace_id != req.auth.uid and token.get("email") != ADMIN_EMAIL:
        shared = (
            db.collection("users")
            .document(workspace_id)
            .collection("collaborators")
            .document(req.auth.uid)
            .get()
        )
        if not shared.exists:
            raise https_fn.HttpsError(
                https_fn.FunctionsErrorCode.PERMISSION_DENIED,
                "Workspace access required to run a screener for this workspace.",
            )
        role = str((shared.to_dict() or {}).get("role") or "viewer").strip().lower()
        if role != "editor":
            raise https_fn.HttpsError(
                https_fn.FunctionsErrorCode.PERMISSION_DENIED,
                "Editor access is required to run a screener for this workspace.",
            )

    context = _remote_config_context(req, token, data.get("meta") if isinstance(data.get("meta"), dict) else None)
    tier_key, tier = _resolve_ai_tier(req.auth.uid, token, context=context)
    allowed_models_raw = tier.get("allowed_models") if isinstance(tier, dict) else []
    allowed_models = [_normalize_ai_model_id(item) for item in (allowed_models_raw or []) if str(item).strip()]
    allowed_models = list(dict.fromkeys([item for item in allowed_models if item]))
    default_weekly = 3 if tier_key == "free" else 25 if tier_key == "pro" else 75
    weekly_limit = int(tier.get("weekly_limit") or tier.get("daily_limit") or default_weekly)
    weekly_limit = max(1, weekly_limit)

    personality = str(data.get("personality") or "").strip().lower()
    personality_model_map = {
        "balanced": "gpt-5-mini",
        "deep_research": "gpt-5",
        "deep-research": "gpt-5",
        "momentum": "gpt-5.1",
        "contrarian": "gpt-5.2",
        "efficient": "gpt-5-nano",
    }

    selected_model = _normalize_ai_model_id(
        data.get("model")
        or data.get("selectedModel")
        or personality_model_map.get(personality)
        or "gpt-5-mini"
    )
    if not selected_model:
        selected_model = allowed_models[0] if allowed_models else "gpt-5-mini"
    if allowed_models and selected_model not in allowed_models:
        selected_model = allowed_models[0]

    _enforce_weekly_usage(
        req.auth.uid,
        "aiScreenerRuns",
        max(1, weekly_limit),
        "Weekly AI screener limit reached. Upgrade to Pro for higher throughput.",
    )

    market = str(data.get("market") or "us").lower()
    try:
        max_names = int(data.get("maxNames") or 25)
    except Exception:
        max_names = 25
    max_names = max(5, min(max_names, 50))

    market_cap_filter_raw = data.get("marketCapFilter") if isinstance(data.get("marketCapFilter"), dict) else {}
    filter_mode = str(market_cap_filter_raw.get("type") or "greater_than").strip().lower() or "greater_than"
    filter_value_raw = market_cap_filter_raw.get("value")
    if filter_value_raw in (None, ""):
        filter_value_raw = data.get("minCap")

    market_cap_target_abs: float | None = None
    try:
        if filter_value_raw not in (None, ""):
            parsed = float(filter_value_raw)
            # Backward compatibility: historical API used billions.
            market_cap_target_abs = parsed * 1_000_000_000 if parsed < 1_000_000_000 else parsed
    except Exception:
        _raise_structured_error(
            https_fn.FunctionsErrorCode.INVALID_ARGUMENT,
            "invalid_market_cap_filter",
            "Market cap filter value is invalid.",
        )

    if filter_mode not in {"greater_than", "less_than", "any"}:
        filter_mode = "greater_than"

    raw_filters = data.get("filters") if isinstance(data.get("filters"), dict) else {}
    screener_filters: dict[str, str] = {}
    for key, value in (raw_filters or {}).items():
        name = str(key or "").strip()
        if not name.startswith("filter"):
            continue
        filter_token = str(value or "").strip()
        if filter_token:
            screener_filters[name] = filter_token

    notes_text = str(data.get("notes") or "")
    note_signals = _resolve_screener_note_signals(notes_text, selected_model, context=context)

    universe_key = str(data.get("universe") or "trending").strip().lower()
    universe_map: dict[str, list[str]] = {
        "large-cap": [
            "AAPL", "MSFT", "NVDA", "GOOGL", "AMZN", "META", "TSLA", "AVGO", "JPM", "V",
            "MA", "WMT", "COST", "XOM", "LLY", "UNH", "JNJ", "PG", "HD", "BAC",
            "KO", "PEP", "ORCL", "CRM", "NFLX", "ADBE", "AMD", "CSCO", "CVX", "MRK",
        ],
        "mid-cap": [
            "HUBS", "ROST", "VRSK", "ANSS", "TTWO", "PAYC", "NVR", "RJF", "PODD", "WAB",
            "FDS", "MGM", "DKNG", "GL", "MTCH", "OKTA", "ZS", "ULTA", "ETSY", "ON",
            "COIN", "RBLX", "MELI", "CFLT", "NET", "SQ", "DDOG", "SNAP", "SHOP", "U",
        ],
        "small-cap": [
            "IWM", "CROX", "RXRX", "CRSR", "SMR", "SOFI", "OPEN", "RUN", "TOST", "AI",
            "SOUN", "UPST", "LMND", "RKLB", "ASTS", "IONQ", "DNA", "BBAI", "ENVX", "ARRY",
            "CHPT", "QS", "BMBL", "RIVN", "NOVA", "S", "PATH", "JOBY", "ACHR", "PL",
        ],
    }

    trending: list[str] = []
    if universe_key == "trending" or universe_key not in universe_map:
        try:
            response = requests.get(TRENDING_URL, headers=_yahoo_headers(), timeout=10)
            response.raise_for_status()
            payload = response.json()
            quotes = payload.get("finance", {}).get("result", [{}])[0].get("quotes", [])
            trending = [str(item.get("symbol") or "").upper() for item in quotes if item.get("symbol")]
        except Exception:
            trending = []

    fallback_diverse = universe_map["large-cap"]
    base_tickers = universe_map.get(universe_key) if universe_key in universe_map else trending
    if not base_tickers:
        base_tickers = trending or fallback_diverse

    note_tickers = [
        _normalize_symbol_token(item)
        for item in (note_signals.get("tickers") if isinstance(note_signals, dict) else [])
        if str(item or "").strip()
    ]
    note_tickers = [item for item in note_tickers if item]

    candidate_pool = max(30, min(max_names * 4, 120))
    seed_tickers = note_tickers + list(base_tickers or []) + list(trending or []) + list(fallback_diverse)
    tickers = list(dict.fromkeys([str(t).upper().strip() for t in seed_tickers if str(t).strip()]))[:candidate_pool]

    results: list[dict[str, Any]] = []
    market_cap_cache: dict[str, float | None] = {}
    info_cache: dict[str, dict[str, Any]] = {}
    dividend_growth_cache: dict[str, dict[str, Any]] = {}

    index_members: dict[str, set[str]] = {
        "sp500": set(universe_map["large-cap"]),
        "nasdaq100": {
            "AAPL", "MSFT", "NVDA", "GOOGL", "AMZN", "META", "AVGO", "TSLA", "AMD", "NFLX",
            "ADBE", "ORCL", "CSCO", "INTC", "QCOM", "TXN", "AMAT", "MU", "LRCX", "KLAC",
        },
        "djia": {
            "AAPL", "AMGN", "AXP", "BA", "CAT", "CRM", "CSCO", "CVX", "DIS", "GS",
            "HD", "HON", "IBM", "JNJ", "JPM", "KO", "MCD", "MMM", "MRK", "MSFT",
            "NKE", "PG", "TRV", "UNH", "V", "VZ", "WBA", "WMT", "XOM",
        },
        "russell2000": set(universe_map["small-cap"]),
    }
    theme_keywords: dict[str, list[str]] = {
        "artificialintelligence": ["artificial intelligence", "machine learning", "ai", "gpu", "accelerator"],
        "cloudcomputing": ["cloud", "saas", "data center", "infrastructure"],
        "cybersecurity": ["cybersecurity", "security", "identity", "threat"],
        "fintech": ["fintech", "payments", "digital banking", "credit"],
        "semiconductors": ["semiconductor", "chip", "foundry", "fab"],
        "robotics": ["robot", "automation", "autonomous"],
        "energyrenewable": ["renewable", "solar", "wind", "battery", "clean energy"],
        "defenseaerospace": ["defense", "aerospace", "drone", "missile"],
    }
    subtheme_keywords: dict[str, list[str]] = {
        "aicompute": ["gpu", "accelerator", "ai compute", "inference"],
        "aienterprise": ["enterprise software", "productivity", "workflow"],
        "cloudsecurity": ["cloud security", "zero trust", "identity"],
        "fintechpayments": ["payments", "merchant", "card network", "wallet"],
        "semiscompute": ["cpu", "gpu", "logic", "compute"],
        "roboticsautomation": ["industrial automation", "robotics", "factory"],
        "energycleansolar": ["solar", "photovoltaic", "inverter"],
        "defensedrones": ["drone", "uav", "autonomous defense"],
    }

    def _info_for_symbol(symbol: str) -> dict[str, Any]:
        cached = info_cache.get(symbol)
        if cached is not None:
            return cached
        info: dict[str, Any] = {}
        try:
            tk = yf.Ticker(symbol)
            raw = tk.info or {}
            if isinstance(raw, dict):
                info = dict(raw)

            fast_info = getattr(tk, "fast_info", None)
            if fast_info is not None:
                fast_market_cap = None
                fast_last_price = None
                try:
                    fast_market_cap = fast_info.get("market_cap")
                    fast_last_price = fast_info.get("last_price")
                except Exception:
                    fast_market_cap = getattr(fast_info, "market_cap", None)
                    fast_last_price = getattr(fast_info, "last_price", None)

                if fast_market_cap not in (None, "") and info.get("marketCap") in (None, ""):
                    info["marketCap"] = fast_market_cap
                if fast_market_cap not in (None, "") and info.get("market_cap") in (None, ""):
                    info["market_cap"] = fast_market_cap
                if fast_last_price not in (None, "") and info.get("currentPrice") in (None, ""):
                    info["currentPrice"] = fast_last_price
        except Exception:
            info = {}
        info_cache[symbol] = info
        return info

    def _to_float(value: Any) -> float | None:
        try:
            num = float(value)
            if math.isfinite(num):
                return num
        except Exception:
            return None
        return None

    def _to_percent(value: Any) -> float | None:
        num = _to_float(value)
        if num is None:
            return None
        if abs(num) <= 1:
            return num * 100.0
        return num

    def _parse_threshold(token: str) -> tuple[str, float] | None:
        text = str(token or "").strip().lower()
        if len(text) < 2 or text[0] not in {"u", "o"}:
            return None
        try:
            return ("under" if text[0] == "u" else "over", float(text[1:]))
        except Exception:
            return None

    def _match_threshold(value: float | None, token: str) -> bool:
        if value is None:
            return False
        text = str(token or "").strip().lower()
        if not text:
            return True
        if "to" in text and not text.startswith(("u", "o")):
            try:
                low_text, high_text = text.split("to", 1)
                low = float(low_text)
                high = float(high_text)
                return low <= value <= high
            except Exception:
                return False
        parsed = _parse_threshold(text)
        if not parsed:
            return False
        mode, threshold = parsed
        if mode == "under":
            return value < threshold
        return value > threshold

    def _match_numeric_filter_token(
        value: float | None,
        token: str,
        *,
        low: float | None = None,
        high: float | None = None,
        positive_label: str = "profitable",
        negative_label: str = "negative",
    ) -> bool:
        if value is None:
            return False
        key = str(token or "").strip().lower()
        if not key:
            return True
        if key == "low" and low is not None:
            return value < low
        if key == "high" and high is not None:
            return value > high
        if key in {positive_label, "positive"}:
            return value > 0
        if key == negative_label:
            return value < 0
        return _match_threshold(value, key)

    def _match_percent_filter_token(
        value: float | None,
        token: str,
        *,
        low: float | None = None,
        high: float | None = None,
        very_high: float | None = None,
        very_neg: float | None = None,
    ) -> bool:
        if value is None:
            return False
        key = str(token or "").strip().lower()
        if not key:
            return True
        if key in {"pos", "positive"}:
            return value > 0
        if key in {"neg", "negative"}:
            return value < 0
        if key == "none":
            return abs(value) <= 1e-9
        if key == "poslow":
            return 0 < value < 10
        if key == "low" and low is not None:
            return value < low
        if key == "high" and high is not None:
            return value > high
        if key == "veryhigh":
            threshold = very_high if very_high is not None else high
            return threshold is not None and value > threshold
        if key == "verypos":
            threshold = very_high if very_high is not None else high
            return threshold is not None and value > threshold
        if key == "veryneg":
            threshold = very_neg if very_neg is not None else -20.0
            return value < threshold
        return _match_threshold(value, key)

    def _match_market_cap_profile(cap: float | None, token: str) -> bool:
        if cap is None:
            return False
        key = str(token or "").strip().lower()
        if not key:
            return True
        ranges: dict[str, tuple[float, float]] = {
            "mega": (200_000_000_000, float("inf")),
            "large": (10_000_000_000, 200_000_000_000),
            "mid": (2_000_000_000, 10_000_000_000),
            "small": (300_000_000, 2_000_000_000),
            "micro": (50_000_000, 300_000_000),
            "nano": (0, 50_000_000),
        }
        if key in ranges:
            low, high = ranges[key]
            return low <= cap < high
        if key == "largeover":
            return cap >= 10_000_000_000
        if key == "midover":
            return cap >= 2_000_000_000
        if key == "smallover":
            return cap >= 300_000_000
        if key == "microover":
            return cap >= 50_000_000
        if key == "largeunder":
            return cap < 200_000_000_000
        if key == "midunder":
            return cap < 10_000_000_000
        if key == "smallunder":
            return cap < 2_000_000_000
        if key == "microunder":
            return cap < 300_000_000
        return True

    def _match_theme_token(text_blob: str, token: str, keyword_map: dict[str, list[str]]) -> bool:
        key = str(token or "").strip().lower()
        if not key:
            return True
        keywords = keyword_map.get(key) or []
        if not keywords:
            return True
        return any(kw in text_blob for kw in keywords)

    def _dividend_growth_stats(symbol: str) -> dict[str, Any]:
        cached = dividend_growth_cache.get(symbol)
        if cached is not None:
            return cached

        stats = {
            "growth1y": None,
            "growth3y": None,
            "growth5y": None,
            "streak": 0,
        }
        try:
            dividends = yf.Ticker(symbol).dividends
            if dividends is not None and getattr(dividends, "empty", True) is False:
                annual = dividends.resample("Y").sum().dropna()
                annual = annual[annual > 0]
                if len(annual) >= 2:
                    last = float(annual.iloc[-1])
                    prev1 = float(annual.iloc[-2])
                    if prev1 > 0:
                        stats["growth1y"] = ((last / prev1) - 1.0) * 100.0
                if len(annual) >= 4:
                    prev3 = float(annual.iloc[-4])
                    last = float(annual.iloc[-1])
                    if prev3 > 0:
                        stats["growth3y"] = (((last / prev3) ** (1.0 / 3.0)) - 1.0) * 100.0
                if len(annual) >= 6:
                    prev5 = float(annual.iloc[-6])
                    last = float(annual.iloc[-1])
                    if prev5 > 0:
                        stats["growth5y"] = (((last / prev5) ** (1.0 / 5.0)) - 1.0) * 100.0

                streak = 0
                values = [float(x) for x in annual.tail(12).tolist()]
                for idx in range(len(values) - 1, 0, -1):
                    if values[idx] > values[idx - 1]:
                        streak += 1
                    else:
                        break
                stats["streak"] = streak
        except Exception:
            pass

        dividend_growth_cache[symbol] = stats
        return stats

    def _fmt_market_cap(value: float | None) -> str:
        if value is None:
            return "—"
        abs_value = abs(float(value))
        if abs_value >= 1_000_000_000_000:
            return f"${value / 1_000_000_000_000:.2f}T"
        if abs_value >= 1_000_000_000:
            return f"${value / 1_000_000_000:.2f}B"
        if abs_value >= 1_000_000:
            return f"${value / 1_000_000:.2f}M"
        return f"${value:,.0f}"

    def _market_cap_for_symbol(symbol: str) -> float | None:
        if symbol in market_cap_cache:
            return market_cap_cache[symbol]
        cap: float | None = None
        try:
            info = _info_for_symbol(symbol)
            cap_raw = info.get("marketCap")
            if cap_raw is not None:
                cap_num = float(cap_raw)
                if math.isfinite(cap_num) and cap_num > 0:
                    cap = cap_num
            if cap is None:
                fast_cap = info.get("market_cap")
                if fast_cap is not None:
                    fast_cap_num = float(fast_cap)
                    if math.isfinite(fast_cap_num) and fast_cap_num > 0:
                        cap = fast_cap_num
        except Exception:
            cap = None
        market_cap_cache[symbol] = cap
        return cap

    def _compute_rsi(close_series: Any, period: int = 14) -> float | None:
        try:
            series = pd.Series(close_series, dtype="float64").dropna()
            if len(series) < period + 1:
                return None
            delta = series.diff()
            gain = delta.clip(lower=0)
            loss = -delta.clip(upper=0)
            avg_gain = gain.ewm(alpha=1 / period, min_periods=period, adjust=False).mean()
            avg_loss = loss.ewm(alpha=1 / period, min_periods=period, adjust=False).mean()
            rs = avg_gain / avg_loss.replace(0, np.nan)
            rsi = 100 - (100 / (1 + rs))
            clean = rsi.dropna()
            if clean.empty:
                return None
            return float(clean.iloc[-1])
        except Exception:
            return None

    def _safe_pct_change(current: float, base: float) -> float | None:
        try:
            curr = float(current)
            prev = float(base)
            if not math.isfinite(curr) or not math.isfinite(prev) or abs(prev) < 1e-12:
                return None
            value = (curr / prev) - 1.0
            return value if math.isfinite(value) else None
        except Exception:
            return None

    def _passes_screener_filters(symbol: str, frame: Any, last_close: float, market_cap: float | None) -> bool:
        if not screener_filters:
            return True

        info = _info_for_symbol(symbol)
        sector = str(info.get("sector") or "").strip()
        industry = str(info.get("industry") or "").strip()
        country = str(info.get("country") or "").strip()
        exchange_raw = str(info.get("exchange") or info.get("fullExchangeName") or "").strip().lower()
        summary_blob = " ".join(
            [
                str(info.get("longBusinessSummary") or ""),
                str(info.get("shortName") or ""),
                sector,
                industry,
            ]
        ).strip().lower()
        cap_value = market_cap if market_cap is not None else _market_cap_for_symbol(symbol)

        def _match_exchange(token: str) -> bool:
            key = str(token or "").strip().lower()
            if not key:
                return True
            aliases = {
                "nasdaq": ["nasdaq", "nms", "ngm", "ncm", "nas"],
                "nyse": ["nyse", "nyq"],
                "amex": ["amex", "ase", "american"],
                "cboe": ["cboe", "bats", "cbo"],
            }
            keys = aliases.get(key) or [key]
            return any(alias in exchange_raw for alias in keys)

        def _match_volume_k(value: float | None, token: str) -> bool:
            if value is None:
                return False
            parsed = _parse_threshold(token)
            if not parsed:
                return False
            mode, amount_k = parsed
            threshold = amount_k * 1_000.0
            return value < threshold if mode == "under" else value > threshold

        def _match_shares_millions(value: float | None, token: str) -> bool:
            if value is None:
                return False
            parsed = _parse_threshold(token)
            if not parsed:
                return False
            mode, amount_m = parsed
            threshold = amount_m * 1_000_000.0
            return value < threshold if mode == "under" else value > threshold

        exchange_filter = screener_filters.get("filterExchange")
        if exchange_filter and not _match_exchange(exchange_filter):
            return False

        index_filter = screener_filters.get("filterIndex")
        if index_filter:
            index_key = str(index_filter).strip().lower()
            members = index_members.get(index_key)
            if members is not None and symbol not in members:
                return False

        sector_filter = screener_filters.get("filterSector")
        if sector_filter and str(sector_filter).strip().lower() != sector.lower():
            return False

        industry_filter = screener_filters.get("filterIndustry")
        if industry_filter and str(industry_filter).strip().lower() not in industry.lower():
            return False

        country_filter = screener_filters.get("filterCountry")
        if country_filter:
            country_key = str(country_filter).strip().lower()
            country_text = country.lower()
            if country_key == "united states":
                if country_text not in {"united states", "usa", "us"}:
                    return False
            elif country_key != country_text:
                return False

        cap_profile_filter = screener_filters.get("filterCapProfile")
        if cap_profile_filter and not _match_market_cap_profile(cap_value, cap_profile_filter):
            return False

        theme_filter = screener_filters.get("filterTheme")
        if theme_filter and not _match_theme_token(summary_blob, theme_filter, theme_keywords):
            return False

        subtheme_filter = screener_filters.get("filterSubtheme")
        if subtheme_filter and not _match_theme_token(summary_blob, subtheme_filter, subtheme_keywords):
            return False

        trailing_pe = _to_float(info.get("trailingPE"))
        forward_pe = _to_float(info.get("forwardPE"))
        peg = _to_float(info.get("pegRatio"))
        price_to_sales = _to_float(info.get("priceToSalesTrailing12Months"))
        price_to_book = _to_float(info.get("priceToBook"))
        total_cash_per_share = _to_float(info.get("totalCashPerShare"))
        free_cash_flow = _to_float(info.get("freeCashflow"))
        enterprise_to_ebitda = _to_float(info.get("enterpriseToEbitda"))
        enterprise_to_sales = _to_float(info.get("enterpriseToRevenue"))
        dividend_yield_pct = _to_percent(info.get("dividendYield"))
        roe_pct = _to_percent(info.get("returnOnEquity"))
        roa_pct = _to_percent(info.get("returnOnAssets"))
        roi_pct = _to_percent(info.get("returnOnCapital")) or roa_pct
        debt_to_equity = _to_float(info.get("debtToEquity"))
        if debt_to_equity is not None and debt_to_equity > 10:
            debt_to_equity = debt_to_equity / 100.0
        lt_debt_to_equity = debt_to_equity
        net_margin_pct = _to_percent(info.get("profitMargins"))
        gross_margin_pct = _to_percent(info.get("grossMargins"))
        operating_margin_pct = _to_percent(info.get("operatingMargins"))
        payout_ratio_pct = _to_percent(info.get("payoutRatio"))
        current_ratio = _to_float(info.get("currentRatio"))
        quick_ratio = _to_float(info.get("quickRatio"))
        insider_ownership_pct = _to_percent(info.get("heldPercentInsiders"))
        institutional_ownership_pct = _to_percent(info.get("heldPercentInstitutions"))
        insider_transactions_pct = _to_percent(info.get("netInsiderBuying"))
        institutional_transactions_pct = _to_percent(info.get("netPercentInstitutionalShares"))
        analyst_recommendation = _to_float(info.get("recommendationMean"))
        earnings_growth_pct = _to_percent(info.get("earningsGrowth"))
        earnings_qoq_pct = _to_percent(info.get("earningsQuarterlyGrowth"))
        sales_growth_pct = _to_percent(info.get("revenueGrowth"))

        price_to_cash = (last_close / total_cash_per_share) if total_cash_per_share and total_cash_per_share > 0 else None
        price_to_fcf = (cap_value / free_cash_flow) if cap_value and free_cash_flow and free_cash_flow > 0 else None

        pe_filter = screener_filters.get("filterPe")
        if pe_filter and not _match_numeric_filter_token(trailing_pe, pe_filter, low=15, high=50):
            return False
        forward_pe_filter = screener_filters.get("filterForwardPe")
        if forward_pe_filter and not _match_numeric_filter_token(forward_pe, forward_pe_filter, low=15, high=50):
            return False
        peg_filter = screener_filters.get("filterPeg")
        if peg_filter and not _match_numeric_filter_token(peg, peg_filter, low=1, high=2):
            return False
        ps_filter = screener_filters.get("filterPs")
        if ps_filter and not _match_numeric_filter_token(price_to_sales, ps_filter, low=1, high=10):
            return False
        pb_filter = screener_filters.get("filterPb")
        if pb_filter and not _match_numeric_filter_token(price_to_book, pb_filter, low=1, high=5):
            return False

        div_filter = screener_filters.get("filterDividendYield")
        if div_filter and not _match_percent_filter_token(dividend_yield_pct, div_filter, high=5, very_high=10):
            return False

        roe_filter = screener_filters.get("filterRoe")
        if roe_filter and not _match_percent_filter_token(roe_pct, roe_filter, high=20, very_high=30, very_neg=-15):
            return False

        roa_filter = screener_filters.get("filterRoa")
        if roa_filter and not _match_percent_filter_token(roa_pct, roa_filter, high=15, very_high=15, very_neg=-15):
            return False

        debt_filter = screener_filters.get("filterDebtEq")
        if debt_filter and not _match_numeric_filter_token(debt_to_equity, debt_filter, low=0.1, high=0.5):
            return False

        net_margin_filter = screener_filters.get("filterNetMargin")
        if net_margin_filter and not _match_percent_filter_token(net_margin_pct, net_margin_filter, high=20, very_neg=-20):
            return False

        analyst_filter = screener_filters.get("filterAnalystRecom")
        if analyst_filter:
            if analyst_recommendation is None:
                return False
            key = str(analyst_filter).strip().lower()
            if key == "strongbuy" and not (analyst_recommendation <= 1.5):
                return False
            if key == "buybetter" and not (analyst_recommendation <= 2.0):
                return False
            if key == "buy" and not (analyst_recommendation <= 2.5):
                return False
            if key == "holdbetter" and not (analyst_recommendation <= 3.0):
                return False
            if key == "holdworse" and not (analyst_recommendation >= 3.0):
                return False
            if key == "strongsell" and not (analyst_recommendation >= 4.5):
                return False

        price_cash_filter = screener_filters.get("filterPriceCash")
        if price_cash_filter and not _match_numeric_filter_token(price_to_cash, price_cash_filter, low=3, high=50):
            return False

        price_fcf_filter = screener_filters.get("filterPriceFcf")
        if price_fcf_filter and not _match_numeric_filter_token(price_to_fcf, price_fcf_filter, low=15, high=50):
            return False

        ev_ebitda_filter = screener_filters.get("filterEvEbitda")
        if ev_ebitda_filter and not _match_numeric_filter_token(
            enterprise_to_ebitda,
            ev_ebitda_filter,
            low=15,
            high=50,
            positive_label="profitable",
            negative_label="negative",
        ):
            return False

        ev_sales_filter = screener_filters.get("filterEvSales")
        if ev_sales_filter and not _match_numeric_filter_token(
            enterprise_to_sales,
            ev_sales_filter,
            low=1,
            high=10,
            positive_label="positive",
            negative_label="negative",
        ):
            return False

        dividend_growth_filter = screener_filters.get("filterDividendGrowth")
        if dividend_growth_filter:
            div_stats = _dividend_growth_stats(symbol)
            key = str(dividend_growth_filter).strip().lower()
            growth1 = _to_float(div_stats.get("growth1y"))
            growth3 = _to_float(div_stats.get("growth3y"))
            growth5 = _to_float(div_stats.get("growth5y"))
            streak = int(div_stats.get("streak") or 0)
            if key == "1ypos" and not (growth1 is not None and growth1 > 0):
                return False
            if key.startswith("1yo"):
                threshold = _to_float(key[3:])
                if threshold is None or growth1 is None or growth1 <= threshold:
                    return False
            if key == "3ypos" and not (growth3 is not None and growth3 > 0):
                return False
            if key.startswith("3yo"):
                threshold = _to_float(key[3:])
                if threshold is None or growth3 is None or growth3 <= threshold:
                    return False
            if key == "5ypos" and not (growth5 is not None and growth5 > 0):
                return False
            if key.startswith("5yo"):
                threshold = _to_float(key[3:])
                if threshold is None or growth5 is None or growth5 <= threshold:
                    return False
            if key.startswith("cy"):
                years = int(_to_float(key[2:]) or 0)
                if years <= 0 or streak < years:
                    return False

        def _match_growth_filter(value: float | None, token: str) -> bool:
            return _match_percent_filter_token(value, token, high=25)

        eps_this_filter = screener_filters.get("filterEpsGrowthThisYear")
        if eps_this_filter and not _match_growth_filter(earnings_growth_pct, eps_this_filter):
            return False
        eps_next_filter = screener_filters.get("filterEpsGrowthNextYear")
        if eps_next_filter and not _match_growth_filter(earnings_growth_pct, eps_next_filter):
            return False
        eps_qoq_filter = screener_filters.get("filterEpsGrowthQoq")
        if eps_qoq_filter and not _match_growth_filter(earnings_qoq_pct, eps_qoq_filter):
            return False
        eps_ttm_filter = screener_filters.get("filterEpsGrowthTtm")
        if eps_ttm_filter and not _match_growth_filter(earnings_qoq_pct, eps_ttm_filter):
            return False
        eps_3y_filter = screener_filters.get("filterEpsGrowth3Years")
        if eps_3y_filter and not _match_growth_filter(earnings_growth_pct, eps_3y_filter):
            return False
        eps_5y_filter = screener_filters.get("filterEpsGrowth5Years")
        if eps_5y_filter and not _match_growth_filter(earnings_growth_pct, eps_5y_filter):
            return False
        eps_n5y_filter = screener_filters.get("filterEpsGrowthNext5Years")
        if eps_n5y_filter and not _match_growth_filter(earnings_growth_pct, eps_n5y_filter):
            return False

        sales_qoq_filter = screener_filters.get("filterSalesGrowthQoq")
        if sales_qoq_filter and not _match_growth_filter(sales_growth_pct, sales_qoq_filter):
            return False
        sales_ttm_filter = screener_filters.get("filterSalesGrowthTtm")
        if sales_ttm_filter and not _match_growth_filter(sales_growth_pct, sales_ttm_filter):
            return False
        sales_3y_filter = screener_filters.get("filterSalesGrowth3Years")
        if sales_3y_filter and not _match_growth_filter(sales_growth_pct, sales_3y_filter):
            return False
        sales_5y_filter = screener_filters.get("filterSalesGrowth5Years")
        if sales_5y_filter and not _match_growth_filter(sales_growth_pct, sales_5y_filter):
            return False

        eps_rev_surprise_filter = screener_filters.get("filterEarningsRevenueSurprise")
        if eps_rev_surprise_filter:
            eps_proxy = earnings_qoq_pct
            rev_proxy = sales_growth_pct
            key = str(eps_rev_surprise_filter).strip().lower()
            if key == "bp" and not (eps_proxy is not None and rev_proxy is not None and eps_proxy > 0 and rev_proxy > 0):
                return False
            if key == "bm" and not (eps_proxy is not None and rev_proxy is not None and abs(eps_proxy) < 0.5 and abs(rev_proxy) < 0.5):
                return False
            if key == "bn" and not (eps_proxy is not None and rev_proxy is not None and eps_proxy < 0 and rev_proxy < 0):
                return False
            if key == "ep" and not (eps_proxy is not None and eps_proxy > 0):
                return False
            if key == "em" and not (eps_proxy is not None and abs(eps_proxy) < 0.5):
                return False
            if key == "en" and not (eps_proxy is not None and eps_proxy < 0):
                return False
            if key == "rp" and not (rev_proxy is not None and rev_proxy > 0):
                return False
            if key == "rm" and not (rev_proxy is not None and abs(rev_proxy) < 0.5):
                return False
            if key == "rn" and not (rev_proxy is not None and rev_proxy < 0):
                return False

        roi_filter = screener_filters.get("filterRoi")
        if roi_filter and not _match_percent_filter_token(roi_pct, roi_filter, high=25, very_high=25, very_neg=-10):
            return False

        current_ratio_filter = screener_filters.get("filterCurrentRatio")
        if current_ratio_filter and not _match_numeric_filter_token(current_ratio, current_ratio_filter, low=1, high=3):
            return False

        quick_ratio_filter = screener_filters.get("filterQuickRatio")
        if quick_ratio_filter and not _match_numeric_filter_token(quick_ratio, quick_ratio_filter, low=0.5, high=3):
            return False

        lt_debt_filter = screener_filters.get("filterLtDebtEquity")
        if lt_debt_filter and not _match_numeric_filter_token(lt_debt_to_equity, lt_debt_filter, low=0.1, high=0.5):
            return False

        gross_margin_filter = screener_filters.get("filterGrossMargin")
        if gross_margin_filter and not _match_percent_filter_token(gross_margin_pct, gross_margin_filter, high=50):
            return False

        operating_margin_filter = screener_filters.get("filterOperatingMargin")
        if operating_margin_filter and not _match_percent_filter_token(
            operating_margin_pct,
            operating_margin_filter,
            high=25,
            very_neg=-20,
        ):
            return False

        payout_ratio_filter = screener_filters.get("filterPayoutRatio")
        if payout_ratio_filter and not _match_percent_filter_token(payout_ratio_pct, payout_ratio_filter, low=20, high=50):
            return False

        insider_own_filter = screener_filters.get("filterInsiderOwnership")
        if insider_own_filter and not _match_percent_filter_token(
            insider_ownership_pct,
            insider_own_filter,
            low=5,
            high=30,
            very_high=50,
        ):
            return False

        insider_tx_filter = screener_filters.get("filterInsiderTransactions")
        if insider_tx_filter and not _match_percent_filter_token(
            insider_transactions_pct,
            insider_tx_filter,
            very_high=20,
            very_neg=-20,
        ):
            return False

        inst_own_filter = screener_filters.get("filterInstitutionalOwnership")
        if inst_own_filter and not _match_percent_filter_token(institutional_ownership_pct, inst_own_filter, low=5, high=90):
            return False

        inst_tx_filter = screener_filters.get("filterInstitutionalTransactions")
        if inst_tx_filter and not _match_percent_filter_token(
            institutional_transactions_pct,
            inst_tx_filter,
            very_high=20,
            very_neg=-20,
        ):
            return False

        current_volume = _to_float(frame["Volume"].iloc[-1]) if "Volume" in frame.columns and len(frame) else None
        avg_volume = _to_float(info.get("averageVolume")) or _to_float(info.get("averageVolume10days"))
        relative_volume = (current_volume / avg_volume) if current_volume and avg_volume and avg_volume > 0 else None
        short_float_pct = _to_percent(info.get("shortPercentOfFloat"))
        shares_outstanding = _to_float(info.get("sharesOutstanding"))
        float_shares = _to_float(info.get("floatShares"))
        float_pct = (float_shares / shares_outstanding * 100.0) if float_shares and shares_outstanding and shares_outstanding > 0 else None

        price_filter = screener_filters.get("filterPrice")
        if price_filter and not _match_threshold(last_close, price_filter):
            return False

        short_float_filter = screener_filters.get("filterShortFloat")
        if short_float_filter and not _match_threshold(short_float_pct, short_float_filter):
            return False

        avg_vol_filter = screener_filters.get("filterAverageVolume")
        if avg_vol_filter and not _match_volume_k(avg_volume, avg_vol_filter):
            return False

        rel_vol_filter = screener_filters.get("filterRelativeVolume")
        if rel_vol_filter and not _match_threshold(relative_volume, rel_vol_filter):
            return False

        cur_vol_filter = screener_filters.get("filterCurrentVolume")
        if cur_vol_filter and not _match_volume_k(current_volume, cur_vol_filter):
            return False

        option_short_filter = screener_filters.get("filterOptionShort")
        if option_short_filter:
            quote_type = str(info.get("quoteType") or "").strip().upper()
            optionable = quote_type == "EQUITY" and (
                _match_exchange("nasdaq") or _match_exchange("nyse") or _match_exchange("amex")
            )
            shortable = bool(optionable and (avg_volume or 0) >= 250_000 and (cap_value or 0) >= 300_000_000 and last_close >= 2.0)
            key = str(option_short_filter).strip().lower()
            if key == "option" and not optionable:
                return False
            if key == "short" and not shortable:
                return False
            if key == "optionshort" and not (optionable and shortable):
                return False

        earnings_filter = screener_filters.get("filterEarningsDate")
        if earnings_filter:
            raw_earnings = info.get("earningsTimestamp") or info.get("earningsDate")
            earnings_date: date | None = None
            if isinstance(raw_earnings, (list, tuple)) and raw_earnings:
                raw_earnings = raw_earnings[0]
            try:
                if isinstance(raw_earnings, datetime):
                    earnings_date = raw_earnings.date()
                elif isinstance(raw_earnings, (int, float)):
                    earnings_date = datetime.fromtimestamp(float(raw_earnings), tz=timezone.utc).date()
                elif raw_earnings:
                    earnings_date = datetime.fromisoformat(str(raw_earnings).replace("Z", "+00:00")).date()
            except Exception:
                earnings_date = None
            if earnings_date is None:
                return False
            today = datetime.now(tz=timezone.utc).date()
            key = str(earnings_filter).strip().lower()
            if key == "today" and earnings_date != today:
                return False
            if key == "tomorrow" and earnings_date != (today + timedelta(days=1)):
                return False
            if key == "nextdays5" and not (today <= earnings_date <= today + timedelta(days=5)):
                return False
            if key == "thisweek":
                now_week = today.isocalendar()[:2]
                target_week = earnings_date.isocalendar()[:2]
                if target_week != now_week:
                    return False

        target_filter = screener_filters.get("filterTargetPrice")
        if target_filter:
            target_price = _to_float(info.get("targetMeanPrice"))
            if target_price is None or last_close <= 0:
                return False
            key = str(target_filter).strip().lower()
            if key == "above" and not (target_price > last_close):
                return False
            if key == "below" and not (target_price < last_close):
                return False
            if key.startswith("a"):
                pct = _to_float(key[1:])
                if pct is None or not (target_price >= last_close * (1 + (pct / 100.0))):
                    return False
            if key.startswith("b"):
                pct = _to_float(key[1:])
                if pct is None or not (target_price <= last_close * (1 - (pct / 100.0))):
                    return False

        ipo_filter = screener_filters.get("filterIpoDate")
        if ipo_filter:
            ipo_epoch = _to_float(info.get("firstTradeDateEpochUtc"))
            if ipo_epoch is None:
                return False
            ipo_date = datetime.fromtimestamp(ipo_epoch, tz=timezone.utc).date()
            age_days = (datetime.now(tz=timezone.utc).date() - ipo_date).days
            key = str(ipo_filter).strip().lower()
            if key == "prevyear" and not (0 <= age_days <= 365):
                return False
            if key == "prev5yrs" and not (0 <= age_days <= 365 * 5):
                return False
            if key == "more5" and not (age_days > 365 * 5):
                return False

        shares_filter = screener_filters.get("filterSharesOutstanding")
        if shares_filter and not _match_shares_millions(shares_outstanding, shares_filter):
            return False

        float_filter = screener_filters.get("filterFloat")
        if float_filter:
            key = str(float_filter).strip().lower()
            if key.endswith("p"):
                percent_token = key[:-1]
                if not _match_threshold(float_pct, percent_token):
                    return False
            elif not _match_shares_millions(float_shares, key):
                return False

        return True

    if tickers:
        try:
            history = yf.download(
                " ".join(tickers),
                period="6mo",
                interval="1d",
                group_by="ticker",
                progress=False,
                threads=True,
                auto_adjust=False,
            )
        except Exception:
            history = pd.DataFrame()

        per_symbol_history: dict[str, Any] = {}

        def _extract(history_frame: Any, symbol: str) -> Any:
            if isinstance(history_frame, pd.DataFrame) and not history_frame.empty:
                if not isinstance(history_frame.columns, pd.MultiIndex):
                    return history_frame.copy()
                if symbol in history_frame.columns.get_level_values(0):
                    frame = history_frame[symbol].copy()
                    if isinstance(frame.columns, pd.MultiIndex):
                        frame.columns = frame.columns.get_level_values(-1)
                    return frame
                if symbol in history_frame.columns.get_level_values(1):
                    frame = history_frame.xs(symbol, level=1, axis=1).copy()
                    if isinstance(frame.columns, pd.MultiIndex):
                        frame.columns = frame.columns.get_level_values(-1)
                    return frame
            cached = per_symbol_history.get(symbol)
            if cached is not None:
                return cached
            try:
                frame = yf.download(
                    symbol,
                    period="6mo",
                    interval="1d",
                    progress=False,
                    auto_adjust=False,
                    threads=False,
                )
            except Exception:
                frame = pd.DataFrame()
            per_symbol_history[symbol] = frame
            return frame

        for symbol in tickers:
            market_cap: float | None = None
            if market_cap_target_abs is not None and filter_mode != "any":
                market_cap = _market_cap_for_symbol(symbol)
                if market_cap is None:
                    continue
                if filter_mode == "greater_than" and not (market_cap >= market_cap_target_abs):
                    continue
                if filter_mode == "less_than" and not (market_cap <= market_cap_target_abs):
                    continue

            frame = _extract(history, symbol)
            if frame is None or getattr(frame, "empty", True) or "Close" not in frame.columns:
                continue
            frame = frame.dropna()
            if frame.empty:
                continue

            close = frame["Close"].astype(float).dropna()
            if len(close) < 30:
                continue

            last_close = float(close.iloc[-1])
            if screener_filters and not _passes_screener_filters(symbol, frame, last_close, market_cap):
                continue
            ret_1m = _safe_pct_change(last_close, float(close.iloc[-21])) if len(close) >= 22 else None
            ret_3m = _safe_pct_change(last_close, float(close.iloc[-63])) if len(close) >= 64 else None

            returns = close.pct_change().dropna()
            vol_raw = float(returns.std() * np.sqrt(252)) if len(returns) else None
            vol_ann = vol_raw if vol_raw is not None and math.isfinite(vol_raw) else None
            rsi_val = _compute_rsi(close, period=14)

            score = 0.0
            if isinstance(ret_3m, float):
                score += 0.65 * ret_3m
            if isinstance(ret_1m, float):
                score += 0.35 * ret_1m
            if isinstance(rsi_val, float):
                score += 0.05 * ((rsi_val - 50.0) / 50.0)

            results.append(
                {
                    "symbol": symbol,
                    "lastClose": round(last_close, 4),
                    "return1m": None if ret_1m is None else round(ret_1m * 100.0, 2),
                    "return3m": None if ret_3m is None else round(ret_3m * 100.0, 2),
                    "rsi14": None if rsi_val is None else round(rsi_val, 2),
                    "volatility": None if vol_ann is None else round(vol_ann, 4),
                    "score": round(score, 6),
                    "marketCap": None if market_cap is None else int(round(market_cap)),
                    "marketCapLabel": _fmt_market_cap(market_cap),
                }
            )

    results.sort(key=lambda item: float(item.get("score") or 0.0), reverse=True)
    results = results[:max_names]

    # Safety fallback: always return a portfolio list even if data providers are partially unavailable.
    fallback_used = False
    if not results:
        fallback_used = True
        backup_universe = tickers[:max_names] or ["AAPL", "MSFT", "NVDA", "AMZN", "GOOGL"][:max_names]
        for idx, symbol in enumerate(backup_universe):
            info = _info_for_symbol(symbol)
            last_close = _to_float(info.get("currentPrice")) or _to_float(info.get("regularMarketPrice"))
            if last_close is None:
                try:
                    frame = yf.download(symbol, period="1mo", interval="1d", progress=False, auto_adjust=False, threads=False)
                    if frame is not None and not frame.empty and "Close" in frame.columns:
                        close_series = frame["Close"].astype(float).dropna()
                        if not close_series.empty:
                            last_close = float(close_series.iloc[-1])
                except Exception:
                    last_close = None
            cap = _market_cap_for_symbol(symbol)
            backup_score = max(0.0, 1.0 - (idx * 0.015))
            results.append(
                {
                    "symbol": symbol,
                    "lastClose": None if last_close is None else round(last_close, 4),
                    "return1m": None,
                    "return3m": None,
                    "rsi14": None,
                    "volatility": None,
                    "score": round(backup_score, 6),
                    "marketCap": None if cap is None else int(round(cap)),
                    "marketCapLabel": _fmt_market_cap(cap),
                }
            )

    if market_cap_target_abs is None or filter_mode == "any":
        for item in results:
            symbol = str(item.get("symbol") or "").strip().upper()
            cap = _market_cap_for_symbol(symbol) if symbol else None
            item["marketCap"] = None if cap is None else int(round(cap))
            item["marketCapLabel"] = _fmt_market_cap(cap)

    run_title = str(data.get("title") or "").strip()
    if not run_title:
        run_title = f"{str(data.get('universe') or 'Trending').title()} screener"

    workspace_profile: dict[str, Any] = {}
    try:
        workspace_snap = db.collection("users").document(workspace_id).get()
        if workspace_snap.exists:
            workspace_doc = workspace_snap.to_dict() or {}
            profile_doc = workspace_doc.get("profile")
            if isinstance(profile_doc, dict):
                workspace_profile = profile_doc
    except Exception:
        workspace_profile = {}

    owner_username = str(workspace_profile.get("username") or "").strip()
    owner_bio = str(workspace_profile.get("bio") or "").strip()
    owner_avatar = str(workspace_profile.get("avatar") or "").strip()
    owner_public_profile = bool(workspace_profile.get("publicProfile"))
    owner_public_screener_pref = bool(workspace_profile.get("publicScreenerSharing"))
    requested_public = bool(data.get("isPublicScreener") or data.get("isPublic"))
    is_public = bool(requested_public or owner_public_screener_pref)

    doc_ref = db.collection("screener_runs").document()
    run_doc = {
        "userId": workspace_id,
        "userEmail": token.get("email"),
        "createdByUid": req.auth.uid,
        "createdByEmail": token.get("email"),
        "market": market,
        "universe": str(data.get("universe") or "trending"),
        "maxNames": max_names,
        "status": "completed",
        "title": run_title,
        "results": _serialize_for_firestore(results),
        "notes": str(data.get("notes") or ""),
        "noteSignals": _serialize_for_firestore(note_signals),
        "fallbackUsed": fallback_used,
        "personality": personality or "",
        "modelUsed": selected_model,
        "modelTier": tier_key,
        "allowedModels": allowed_models,
        "weeklyLimit": weekly_limit,
        "dailyLimit": weekly_limit,
        "isPublic": is_public,
        "ownerUsername": owner_username,
        "ownerBio": owner_bio,
        "ownerAvatar": owner_avatar,
        "ownerPublicProfile": owner_public_profile,
        "publishedAt": firestore.SERVER_TIMESTAMP if is_public else None,
        "marketCapFilter": {
            "type": filter_mode,
            "value": market_cap_target_abs,
        },
        "filters": screener_filters,
        "createdAt": firestore.SERVER_TIMESTAMP,
        "updatedAt": firestore.SERVER_TIMESTAMP,
        "meta": data.get("meta") or {},
    }
    try:
        doc_ref.set(run_doc)
    except Exception as exc:
        _raise_structured_error(
            https_fn.FunctionsErrorCode.INTERNAL,
            "screener_persist_error",
            "Screener completed, but saving the run failed.",
            {"raw": str(exc)[:400]},
        )
    _audit_event(req.auth.uid, token.get("email"), "screener_completed", {"runId": doc_ref.id, "count": len(results)})

    return {
        "runId": doc_ref.id,
        "status": "completed",
        "title": run_title,
        "results": results,
        "workspaceId": workspace_id,
        "modelUsed": selected_model,
        "modelTier": tier_key,
        "personality": personality or "",
        "weeklyLimit": weekly_limit,
        "dailyLimit": weekly_limit,
        "isPublic": is_public,
        "allowedModels": allowed_models,
        "resultsFound": len(results),
        "noteSignals": note_signals,
        "fallbackUsed": fallback_used,
        "marketCapFilter": {"type": filter_mode, "value": market_cap_target_abs},
        "filters": screener_filters,
    }


@https_fn.on_call()
def submit_feature_vote(req: https_fn.CallableRequest) -> dict[str, Any]:
    token = _require_auth(req)
    data = req.data or {}

    feature_key = str(data.get("featureKey") or data.get("feature") or "").strip().lower()
    if feature_key not in FEATURE_VOTE_KEYS:
        raise https_fn.HttpsError(
            https_fn.FunctionsErrorCode.INVALID_ARGUMENT,
            "Feature key must be one of: uploads, autopilot.",
        )

    vote = str(data.get("vote") or "").strip().lower()
    if vote not in FEATURE_VOTE_CHOICES:
        raise https_fn.HttpsError(
            https_fn.FunctionsErrorCode.INVALID_ARGUMENT,
            "Vote must be yes or no.",
        )

    feedback = str(data.get("feedback") or "").strip()[:2000]
    meta = data.get("meta") if isinstance(data.get("meta"), dict) else {}
    user_id = req.auth.uid
    user_email = str(token.get("email") or "").strip()

    vote_ref = db.collection("feature_votes").document(f"{feature_key}__{user_id}")
    totals_ref = db.collection("feature_vote_totals").document(feature_key)
    transaction = db.transaction()

    vote_snap = vote_ref.get(transaction=transaction)
    previous_vote = ""
    if vote_snap.exists:
        previous_vote = str((vote_snap.to_dict() or {}).get("vote") or "").strip().lower()

    totals_snap = totals_ref.get(transaction=transaction)
    totals_doc = totals_snap.to_dict() if totals_snap.exists else {}

    def _to_int(value: Any) -> int:
        try:
            return max(0, int(value or 0))
        except Exception:
            return 0

    yes_count = _to_int(totals_doc.get("yes"))
    no_count = _to_int(totals_doc.get("no"))
    total_count = _to_int(totals_doc.get("total"))

    had_previous = previous_vote in FEATURE_VOTE_CHOICES
    if not had_previous:
        total_count += 1
    elif previous_vote != vote:
        if previous_vote == "yes":
            yes_count = max(0, yes_count - 1)
        elif previous_vote == "no":
            no_count = max(0, no_count - 1)

    if not had_previous or previous_vote != vote:
        if vote == "yes":
            yes_count += 1
        elif vote == "no":
            no_count += 1

    vote_payload = {
        "featureKey": feature_key,
        "vote": vote,
        "feedback": feedback,
        "userId": user_id,
        "userEmail": user_email,
        "meta": meta,
        "updatedAt": firestore.SERVER_TIMESTAMP,
    }
    if not vote_snap.exists:
        vote_payload["createdAt"] = firestore.SERVER_TIMESTAMP
    transaction.set(vote_ref, vote_payload, merge=True)

    transaction.set(
        totals_ref,
        {
            "featureKey": feature_key,
            "yes": yes_count,
            "no": no_count,
            "total": total_count,
            "updatedAt": firestore.SERVER_TIMESTAMP,
        },
        merge=True,
    )
    transaction.commit()

    _audit_event(
        user_id,
        user_email,
        "feature_vote_submitted",
        {"featureKey": feature_key, "vote": vote},
    )
    return {
        "featureKey": feature_key,
        "vote": vote,
        "totals": {
            "yes": yes_count,
            "no": no_count,
            "total": total_count,
        },
    }


@https_fn.on_call()
def get_feature_vote_summary(req: https_fn.CallableRequest) -> dict[str, Any]:
    token = _require_auth(req)
    _require_admin(token)
    data = req.data or {}

    try:
        limit = int(data.get("limit") or 25)
    except Exception:
        limit = 25
    limit = max(5, min(limit, 100))

    def _to_int(value: Any) -> int:
        try:
            return max(0, int(value or 0))
        except Exception:
            return 0

    features: dict[str, dict[str, Any]] = {}
    for feature_key in sorted(FEATURE_VOTE_KEYS):
        snap = db.collection("feature_vote_totals").document(feature_key).get()
        doc = snap.to_dict() if snap.exists else {}
        yes_count = _to_int(doc.get("yes"))
        no_count = _to_int(doc.get("no"))
        total_count = _to_int(doc.get("total") or (yes_count + no_count))
        yes_percent = round((yes_count / total_count) * 100, 1) if total_count > 0 else 0.0
        features[feature_key] = {
            "yes": max(0, yes_count),
            "no": max(0, no_count),
            "total": max(0, total_count),
            "yesPercent": yes_percent,
            "updatedAt": doc.get("updatedAt"),
        }

    recent_docs = (
        db.collection("feature_votes")
        .order_by("updatedAt", direction=firestore.Query.DESCENDING)
        .limit(limit)
        .stream()
    )
    recent: list[dict[str, Any]] = []
    for item in recent_docs:
        payload = item.to_dict() or {}
        feature_key = str(payload.get("featureKey") or "").strip().lower()
        if feature_key not in FEATURE_VOTE_KEYS:
            continue
        recent.append(
            {
                "id": item.id,
                "featureKey": feature_key,
                "vote": str(payload.get("vote") or "").strip().lower(),
                "feedback": str(payload.get("feedback") or "").strip(),
                "userEmail": str(payload.get("userEmail") or "").strip(),
                "updatedAt": payload.get("updatedAt"),
            }
        )

    return {
        "features": _serialize_for_firestore(features),
        "recent": _serialize_for_firestore(recent),
    }


@https_fn.on_call()
def queue_autopilot_run(req: https_fn.CallableRequest) -> dict[str, Any]:
    token = _require_auth(req)
    _require_admin(token)
    data = req.data or {}
    payload = {
        "userId": req.auth.uid,
        "userEmail": token.get("email"),
        "ticker": data.get("ticker"),
        "horizon": data.get("horizon"),
        "quantiles": data.get("quantiles"),
        "interval": data.get("interval"),
        "notes": data.get("notes"),
        "status": "queued",
        "createdAt": firestore.SERVER_TIMESTAMP,
        "updatedAt": firestore.SERVER_TIMESTAMP,
        "meta": data.get("meta") or {},
    }

    doc_ref = db.collection("autopilot_requests").document()
    doc_ref.set(payload)
    _audit_event(req.auth.uid, token.get("email"), "autopilot_queued", {"requestId": doc_ref.id})
    return {"requestId": doc_ref.id}


@https_fn.on_call()
def delete_autopilot_request(req: https_fn.CallableRequest) -> dict[str, Any]:
    token = _require_auth(req)
    _require_admin(token)
    data = req.data or {}
    request_id = str(data.get("requestId") or data.get("id") or "").strip()
    if not request_id:
        raise https_fn.HttpsError(https_fn.FunctionsErrorCode.INVALID_ARGUMENT, "Request ID is required.")

    ref = db.collection("autopilot_requests").document(request_id)
    snap = ref.get()
    if not snap.exists:
        raise https_fn.HttpsError(https_fn.FunctionsErrorCode.NOT_FOUND, "Autopilot request not found.")

    ref.delete()
    _audit_event(req.auth.uid, token.get("email"), "autopilot_deleted", {"requestId": request_id})
    return {"deleted": True, "requestId": request_id}


@https_fn.on_call()
def submit_feedback(req: https_fn.CallableRequest) -> dict[str, Any]:
    data = req.data or {}
    message = str(data.get("message") or "").strip()
    if not message:
        raise https_fn.HttpsError(https_fn.FunctionsErrorCode.INVALID_ARGUMENT, "Feedback message is required.")

    rating = data.get("rating")
    page_path = str(data.get("pagePath") or "").strip()
    meta = data.get("meta") if isinstance(data.get("meta"), dict) else {}

    user_id = req.auth.uid if req.auth else ""
    user_email = ""
    if req.auth and req.auth.token:
        user_email = str(req.auth.token.get("email") or "")
    if not user_email:
        user_email = str(data.get("email") or "").strip()

    doc_ref = db.collection("feedback").document()
    doc_ref.set(
        {
            "userId": user_id,
            "userEmail": user_email,
            "rating": rating,
            "message": message,
            "pagePath": page_path or str(meta.get("pagePath") or ""),
            "meta": meta,
            "createdAt": firestore.SERVER_TIMESTAMP,
        }
    )

    if user_id:
        _audit_event(user_id, user_email, "feedback_submitted", {"feedbackId": doc_ref.id})

    return {"feedbackId": doc_ref.id}


@https_fn.on_call()
def generate_social_campaign_drafts(req: https_fn.CallableRequest) -> dict[str, Any]:
    token = _require_auth(req)
    data = req.data or {}

    topic = str(data.get("topic") or "").strip() or "Weekly market positioning update"
    objective = str(data.get("objective") or "").strip() or "Drive qualified leads to Quantura"
    audience = str(data.get("audience") or "").strip() or "active investors and growth-focused teams"
    tone = str(data.get("tone") or "").strip() or "confident, concise, practical"
    channels = _normalize_social_channels(data.get("channels"))
    posting_routes = {channel: _social_channel_posting_route(channel) for channel in channels}
    routable_channels = [channel for channel, route in posting_routes.items() if route != "none"]
    unroutable_channels = [channel for channel in channels if channel not in routable_channels]
    posts_per_channel = max(1, min(int(data.get("postsPerChannel") or 2), 5))
    cta_url = str(data.get("ctaUrl") or SOCIAL_DEFAULT_CTA_URL).strip() or SOCIAL_DEFAULT_CTA_URL
    title = str(data.get("title") or f"{topic} campaign").strip()
    save_draft = bool(data.get("saveDraft", True))

    drafts, model_used = _generate_social_copy_with_openai(
        topic=topic,
        objective=objective,
        audience=audience,
        tone=tone,
        channels=channels,
        posts_per_channel=posts_per_channel,
        cta_url=cta_url,
    )

    campaign_id = ""
    if save_draft:
        campaign_ref = db.collection("social_campaigns").document()
        campaign_ref.set(
            {
                "userId": req.auth.uid,
                "userEmail": token.get("email"),
                "title": title,
                "topic": topic,
                "objective": objective,
                "audience": audience,
                "tone": tone,
                "channels": channels,
                "postsPerChannel": posts_per_channel,
                "ctaUrl": cta_url,
                "drafts": _serialize_for_firestore(drafts),
                "model": model_used,
                "status": "draft",
                "createdAt": firestore.SERVER_TIMESTAMP,
                "updatedAt": firestore.SERVER_TIMESTAMP,
                "meta": data.get("meta") or {},
            }
        )
        campaign_id = campaign_ref.id

    _audit_event(
        req.auth.uid,
        token.get("email"),
        "social_drafts_generated",
        {
            "campaignId": campaign_id,
            "channels": channels,
            "routableChannels": routable_channels,
            "unroutableChannels": unroutable_channels,
            "postsPerChannel": posts_per_channel,
            "model": model_used,
        },
    )

    return {
        "campaignId": campaign_id,
        "title": title,
        "topic": topic,
        "channels": channels,
        "routableChannels": routable_channels,
        "unroutableChannels": unroutable_channels,
        "postingRoutes": posting_routes,
        "postsPerChannel": posts_per_channel,
        "model": model_used,
        "drafts": drafts,
    }


@https_fn.on_call()
def queue_social_campaign_posts(req: https_fn.CallableRequest) -> dict[str, Any]:
    token = _require_auth(req)
    data = req.data or {}
    campaign_id = str(data.get("campaignId") or "").strip()
    if not campaign_id:
        raise https_fn.HttpsError(https_fn.FunctionsErrorCode.INVALID_ARGUMENT, "campaignId is required.")

    campaign_ref = db.collection("social_campaigns").document(campaign_id)
    campaign_snap = campaign_ref.get()
    if not campaign_snap.exists:
        raise https_fn.HttpsError(https_fn.FunctionsErrorCode.NOT_FOUND, "Campaign not found.")

    campaign = campaign_snap.to_dict() or {}
    owner_id = str(campaign.get("userId") or "").strip()
    if token.get("email") != ADMIN_EMAIL and owner_id != req.auth.uid:
        raise https_fn.HttpsError(https_fn.FunctionsErrorCode.PERMISSION_DENIED, "Access denied.")

    drafts = campaign.get("drafts") if isinstance(campaign.get("drafts"), dict) else {}
    if not drafts:
        raise https_fn.HttpsError(https_fn.FunctionsErrorCode.FAILED_PRECONDITION, "Campaign has no drafts to queue.")

    requested_channels = _normalize_social_channels(data.get("channels"))
    channels = [channel for channel in requested_channels if channel in drafts] or list(drafts.keys())
    if not channels:
        raise https_fn.HttpsError(https_fn.FunctionsErrorCode.FAILED_PRECONDITION, "No channels available to queue.")
    routable_channels = _configured_social_posting_channels(channels)
    unroutable_channels = [channel for channel in channels if channel not in routable_channels]
    channels = routable_channels
    if not channels:
        raise https_fn.HttpsError(
            https_fn.FunctionsErrorCode.FAILED_PRECONDITION,
            "No posting channels are configured. Add direct API credentials or SOCIAL_WEBHOOK_* targets.",
        )

    scheduled_for = _as_utc_datetime(data.get("scheduledFor")) or datetime.now(timezone.utc)
    if scheduled_for < datetime.now(timezone.utc) - timedelta(minutes=2):
        scheduled_for = datetime.now(timezone.utc)

    enqueue = _enqueue_social_posts(
        campaign_id=campaign_id,
        user_id=owner_id or req.auth.uid,
        user_email=str(campaign.get("userEmail") or token.get("email") or ""),
        channels=channels,
        drafts=drafts,
        scheduled_for=scheduled_for,
        meta=data.get("meta") if isinstance(data.get("meta"), dict) else {},
    )

    campaign_ref.set(
        {
            "status": "queued",
            "queuedAt": firestore.SERVER_TIMESTAMP,
            "queuedCount": enqueue["count"],
            "queuedChannels": channels,
            "unroutableChannels": unroutable_channels,
            "scheduledFor": scheduled_for,
            "updatedAt": firestore.SERVER_TIMESTAMP,
        },
        merge=True,
    )

    _audit_event(
        req.auth.uid,
        token.get("email"),
        "social_posts_queued",
        {
            "campaignId": campaign_id,
            "queuedCount": enqueue["count"],
            "channels": channels,
            "unroutableChannels": unroutable_channels,
        },
    )
    return {
        "campaignId": campaign_id,
        "queuedCount": enqueue["count"],
        "queueIds": enqueue["queueIds"],
        "scheduledFor": scheduled_for.isoformat(),
        "channels": channels,
        "unroutableChannels": unroutable_channels,
    }


@https_fn.on_call()
def list_social_campaigns(req: https_fn.CallableRequest) -> dict[str, Any]:
    token = _require_auth(req)
    limit = max(1, min(int((req.data or {}).get("limit") or 40), 100))
    query = db.collection("social_campaigns")
    if token.get("email") != ADMIN_EMAIL:
        query = query.where("userId", "==", req.auth.uid)

    docs = query.order_by("createdAt", direction=firestore.Query.DESCENDING).limit(limit).stream()
    items = []
    for doc in docs:
        row = doc.to_dict() or {}
        row["id"] = doc.id
        items.append(row)
    return {"campaigns": _serialize_for_firestore(items)}


@https_fn.on_call()
def list_social_queue(req: https_fn.CallableRequest) -> dict[str, Any]:
    token = _require_auth(req)
    data = req.data or {}
    limit = max(1, min(int(data.get("limit") or 80), 200))
    status_filter = str(data.get("status") or "").strip().lower()
    query = db.collection("social_queue")
    if token.get("email") != ADMIN_EMAIL:
        query = query.where("userId", "==", req.auth.uid)
    if status_filter:
        query = query.where("status", "==", status_filter)

    docs = query.order_by("createdAt", direction=firestore.Query.DESCENDING).limit(limit).stream()
    items = []
    for doc in docs:
        row = doc.to_dict() or {}
        row["id"] = doc.id
        items.append(row)
    return {"queue": _serialize_for_firestore(items)}


@https_fn.on_call()
def publish_social_queue_now(req: https_fn.CallableRequest) -> dict[str, Any]:
    token = _require_auth(req)
    _require_admin(token)
    if not SOCIAL_AUTOMATION_ENABLED:
        raise https_fn.HttpsError(https_fn.FunctionsErrorCode.FAILED_PRECONDITION, "Social automation is disabled.")
    data = req.data or {}
    max_posts = max(1, min(int(data.get("maxPosts") or SOCIAL_DISPATCH_BATCH_SIZE), 100))
    summary = _dispatch_due_social_posts(max_posts=max_posts, trigger="manual", actor_uid=req.auth.uid)
    _audit_event(req.auth.uid, token.get("email"), "social_dispatch_manual", summary)
    return summary


@scheduler_fn.on_schedule(
    schedule="0 * * * *",
    timezone=scheduler_fn.Timezone(SOCIAL_AUTOMATION_TIMEZONE),
)
def social_dispatch_scheduler(event: scheduler_fn.ScheduledEvent) -> None:
    del event
    if not SOCIAL_AUTOMATION_ENABLED:
        return
    _dispatch_due_social_posts(max_posts=SOCIAL_DISPATCH_BATCH_SIZE, trigger="scheduler")


def _autopilot_social_channels() -> list[str]:
    channels = [channel for channel in SOCIAL_AUTOPILOT_CHANNELS if channel in SOCIAL_POPULAR_CHANNELS]
    if not channels:
        channels = ["x", "linkedin", "facebook", "instagram", "tiktok"]
    channels = list(dict.fromkeys(channels))
    return _configured_social_posting_channels(channels)


def _run_social_autopilot_plan(trigger: str, actor_uid: str = "") -> dict[str, Any]:
    if not SOCIAL_AUTOPILOT_ENABLED:
        return {"ok": False, "skipped": True, "reason": "autopilot_disabled"}

    channels = _autopilot_social_channels()
    if not channels:
        return {"ok": False, "skipped": True, "reason": "no_channels"}

    local_now = datetime.now(_posting_tzinfo())
    date_key = local_now.date().isoformat()
    campaign_id = f"autopilot_{date_key}"
    campaign_ref = db.collection("social_campaigns").document(campaign_id)
    snap = campaign_ref.get()
    if snap.exists:
        return {"ok": True, "skipped": True, "reason": "already_scheduled", "campaignId": campaign_id}

    drafts, model_used = _generate_social_copy_with_openai(
        topic=SOCIAL_AUTOPILOT_TOPIC,
        objective=SOCIAL_AUTOPILOT_OBJECTIVE,
        audience=SOCIAL_AUTOPILOT_AUDIENCE,
        tone=SOCIAL_AUTOPILOT_TONE,
        channels=channels,
        posts_per_channel=SOCIAL_AUTOPILOT_POSTS_PER_CHANNEL,
        cta_url=SOCIAL_DEFAULT_CTA_URL,
    )
    scheduled_at = datetime.now(timezone.utc)
    enqueue = _enqueue_social_posts(
        campaign_id=campaign_id,
        user_id=SOCIAL_AUTOPILOT_USER_ID,
        user_email=SOCIAL_AUTOPILOT_USER_EMAIL,
        channels=channels,
        drafts=drafts,
        scheduled_for=scheduled_at,
        meta={"trigger": trigger, "dateKey": date_key},
    )

    campaign_ref.set(
        {
            "id": campaign_id,
            "userId": SOCIAL_AUTOPILOT_USER_ID,
            "userEmail": SOCIAL_AUTOPILOT_USER_EMAIL,
            "topic": SOCIAL_AUTOPILOT_TOPIC,
            "objective": SOCIAL_AUTOPILOT_OBJECTIVE,
            "audience": SOCIAL_AUTOPILOT_AUDIENCE,
            "tone": SOCIAL_AUTOPILOT_TONE,
            "channels": channels,
            "postsPerChannel": SOCIAL_AUTOPILOT_POSTS_PER_CHANNEL,
            "drafts": _serialize_for_firestore(drafts),
            "modelUsed": model_used,
            "status": "queued",
            "queuedCount": enqueue["count"],
            "queuedChannels": channels,
            "scheduledFor": scheduled_at,
            "createdAt": firestore.SERVER_TIMESTAMP,
            "updatedAt": firestore.SERVER_TIMESTAMP,
            "meta": {"trigger": trigger, "actorUid": actor_uid, "dateKey": date_key},
        },
        merge=True,
    )
    return {
        "ok": True,
        "campaignId": campaign_id,
        "queuedCount": enqueue["count"],
        "queueIds": enqueue["queueIds"],
        "channels": channels,
        "modelUsed": model_used,
        "dateKey": date_key,
    }


@https_fn.on_call()
def schedule_social_autopilot_now(req: https_fn.CallableRequest) -> dict[str, Any]:
    token = _require_auth(req)
    _require_admin(token)
    summary = _run_social_autopilot_plan(trigger="manual", actor_uid=req.auth.uid)
    _audit_event(req.auth.uid, token.get("email"), "social_autopilot_scheduled", summary)
    return summary


@scheduler_fn.on_schedule(
    schedule="15 7 * * *",
    timezone=scheduler_fn.Timezone(SOCIAL_AUTOMATION_TIMEZONE),
)
def social_daily_planner_scheduler(event: scheduler_fn.ScheduledEvent) -> None:
    del event
    _run_social_autopilot_plan(trigger="scheduler")


@https_fn.on_call()
def alpaca_get_options(req: https_fn.CallableRequest) -> dict[str, Any]:
    token = _require_auth(req)
    _require_admin(token)
    data = req.data or {}
    ticker = str(data.get("ticker") or "").upper()
    expiration = data.get("expiration")
    if not ticker:
        raise https_fn.HttpsError(https_fn.FunctionsErrorCode.INVALID_ARGUMENT, "Ticker is required.")

    options = []
    try:
        headers = _alpaca_headers()
        params = {"underlying_symbols": ticker, "limit": 30}
        if expiration:
            params["expiration_date"] = expiration
        response = requests.get(
            f"{ALPACA_DATA_BASE}/v1beta1/options/contracts",
            headers=headers,
            params=params,
            timeout=12,
        )
        response.raise_for_status()
        contracts = response.json().get("option_contracts", [])
        for contract in contracts[:30]:
            delta = _safe_float(contract.get("delta"))
            implied_prob = None
            if delta is not None:
                implied_prob = max(0.0, min(abs(delta), 1.0))
            options.append(
                {
                    "symbol": contract.get("symbol"),
                    "type": contract.get("type"),
                    "strike": contract.get("strike_price"),
                    "expiration": contract.get("expiration_date"),
                    "lastPrice": contract.get("close_price"),
                    "delta": None if delta is None else round(delta, 4),
                    "impliedProbability": None if implied_prob is None else round(implied_prob * 100, 2),
                }
            )
    except Exception:
        try:
            ticker_obj = yf.Ticker(ticker)
            expirations = ticker_obj.options
            if expirations:
                exp = expiration or expirations[0]
                chain = ticker_obj.option_chain(exp)
                for _, row in chain.calls.head(10).iterrows():
                    options.append(
                        {
                            "symbol": row.get("contractSymbol"),
                            "type": "call",
                            "strike": row.get("strike"),
                            "expiration": exp,
                            "lastPrice": row.get("lastPrice"),
                            "delta": None,
                            "impliedProbability": None,
                        }
                    )
                for _, row in chain.puts.head(10).iterrows():
                    options.append(
                        {
                            "symbol": row.get("contractSymbol"),
                            "type": "put",
                            "strike": row.get("strike"),
                            "expiration": exp,
                            "lastPrice": row.get("lastPrice"),
                            "delta": None,
                            "impliedProbability": None,
                        }
                    )
        except Exception:
            options = []

    return {"options": options}


@https_fn.on_call()
def alpaca_place_order(req: https_fn.CallableRequest) -> dict[str, Any]:
    token = _require_auth(req)
    _require_admin(token)
    data = req.data or {}

    symbol = str(data.get("symbol") or "").upper().strip()
    if not symbol:
        raise https_fn.HttpsError(https_fn.FunctionsErrorCode.INVALID_ARGUMENT, "Symbol is required.")

    order_payload: dict[str, Any] = {
        "symbol": symbol,
        "qty": int(data.get("qty") or 1),
        "side": str(data.get("side") or "buy"),
        "type": str(data.get("type") or "market"),
        "time_in_force": str(data.get("timeInForce") or "day"),
    }

    limit_price = _safe_float(data.get("limitPrice"))
    if order_payload["type"] == "limit" and limit_price:
        order_payload["limit_price"] = limit_price

    response = requests.post(
        f"{ALPACA_API_BASE}/v2/orders",
        headers=_alpaca_headers(),
        json=order_payload,
        timeout=12,
    )

    if response.status_code >= 400:
        raise https_fn.HttpsError(
            https_fn.FunctionsErrorCode.FAILED_PRECONDITION,
            f"Alpaca error: {response.text}",
        )

    order_data = response.json()
    doc_ref = db.collection("alpaca_orders").document()
    doc_ref.set(
        {
            "userId": req.auth.uid,
            "userEmail": token.get("email"),
            "symbol": symbol,
            "side": order_payload["side"],
            "type": order_payload["type"],
            "qty": order_payload["qty"],
            "status": order_data.get("status"),
            "alpacaId": order_data.get("id"),
            "response": _serialize_for_firestore(order_data),
            "createdAt": firestore.SERVER_TIMESTAMP,
            "meta": data.get("meta") or {},
        }
    )

    _audit_event(
        req.auth.uid,
        token.get("email"),
        "alpaca_order_placed",
        {"orderId": doc_ref.id, "alpacaId": order_data.get("id"), "symbol": symbol},
    )

    return {"orderId": doc_ref.id, "alpacaId": order_data.get("id")}


@https_fn.on_call()
def alpaca_get_account(req: https_fn.CallableRequest) -> dict[str, Any]:
    token = _require_auth(req)
    _require_admin(token)
    response = requests.get(f"{ALPACA_API_BASE}/v2/account", headers=_alpaca_headers(), timeout=12)
    if response.status_code >= 400:
        raise https_fn.HttpsError(https_fn.FunctionsErrorCode.FAILED_PRECONDITION, response.text)
    account = response.json()
    return {
        "account": {
            "id": account.get("id"),
            "status": account.get("status"),
            "currency": account.get("currency"),
            "buying_power": account.get("buying_power"),
            "equity": account.get("equity"),
            "cash": account.get("cash"),
            "portfolio_value": account.get("portfolio_value"),
            "pattern_day_trader": account.get("pattern_day_trader"),
        }
    }


@https_fn.on_call()
def alpaca_get_positions(req: https_fn.CallableRequest) -> dict[str, Any]:
    token = _require_auth(req)
    _require_admin(token)
    response = requests.get(f"{ALPACA_API_BASE}/v2/positions", headers=_alpaca_headers(), timeout=12)
    if response.status_code >= 400:
        raise https_fn.HttpsError(https_fn.FunctionsErrorCode.FAILED_PRECONDITION, response.text)
    positions = response.json()
    clean_positions = [
        {
            "symbol": item.get("symbol"),
            "qty": item.get("qty"),
            "market_value": item.get("market_value"),
            "cost_basis": item.get("cost_basis"),
            "unrealized_pl": item.get("unrealized_pl"),
            "unrealized_plpc": item.get("unrealized_plpc"),
            "side": item.get("side"),
        }
        for item in positions
    ]
    return {"positions": clean_positions}


@https_fn.on_call()
def alpaca_list_orders(req: https_fn.CallableRequest) -> dict[str, Any]:
    token = _require_auth(req)
    _require_admin(token)
    data = req.data or {}
    status = str(data.get("status") or "all")
    limit = int(data.get("limit") or 50)
    params = {"status": status, "limit": max(1, min(limit, 100)), "direction": "desc"}

    response = requests.get(f"{ALPACA_API_BASE}/v2/orders", headers=_alpaca_headers(), params=params, timeout=12)
    if response.status_code >= 400:
        raise https_fn.HttpsError(https_fn.FunctionsErrorCode.FAILED_PRECONDITION, response.text)

    orders = response.json()
    clean_orders = [
        {
            "id": item.get("id"),
            "symbol": item.get("symbol"),
            "status": item.get("status"),
            "side": item.get("side"),
            "type": item.get("type"),
            "qty": item.get("qty"),
            "filled_qty": item.get("filled_qty"),
            "filled_avg_price": item.get("filled_avg_price"),
            "created_at": item.get("created_at"),
        }
        for item in orders
    ]
    return {"orders": clean_orders}


@https_fn.on_call()
def alpaca_cancel_order(req: https_fn.CallableRequest) -> dict[str, Any]:
    token = _require_auth(req)
    _require_admin(token)
    data = req.data or {}
    order_id = str(data.get("orderId") or "").strip()
    if not order_id:
        raise https_fn.HttpsError(https_fn.FunctionsErrorCode.INVALID_ARGUMENT, "orderId is required.")

    response = requests.delete(f"{ALPACA_API_BASE}/v2/orders/{order_id}", headers=_alpaca_headers(), timeout=12)
    if response.status_code >= 400:
        raise https_fn.HttpsError(https_fn.FunctionsErrorCode.FAILED_PRECONDITION, response.text)

    _audit_event(req.auth.uid, token.get("email"), "alpaca_order_cancelled", {"alpacaId": order_id})
    return {"cancelled": True, "orderId": order_id}


@https_fn.on_call()
def send_slack_test_message(req: https_fn.CallableRequest) -> dict[str, Any]:
    token = _require_auth(req)
    _require_admin(token)

    if not SLACK_WEBHOOK_URL:
        raise https_fn.HttpsError(
            https_fn.FunctionsErrorCode.FAILED_PRECONDITION,
            "SLACK_WEBHOOK_URL is not configured.",
        )

    data = req.data or {}
    text = str(data.get("text") or "Quantura Slack webhook test from Firebase function.")
    payload = {"text": text}
    response = requests.post(SLACK_WEBHOOK_URL, json=payload, timeout=10)
    if response.status_code >= 400:
        raise https_fn.HttpsError(
            https_fn.FunctionsErrorCode.INTERNAL,
            f"Slack webhook failed: {response.status_code} {response.text}",
        )

    _audit_event(req.auth.uid, token.get("email"), "slack_test_sent", {"message": text})
    return {"ok": True}
